#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""HUIIHONG v2 -- AI app layer investment PPT (16:9)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

C = {
    'brand_red':      RGBColor(0xC0, 0x00, 0x00),
    'dark_red':       RGBColor(0x8B, 0x00, 0x00),
    'gold':           RGBColor(0xD4, 0xAF, 0x37),
    'gold_light':     RGBColor(0xE8, 0xD5, 0x8A),
    'gold_warm':      RGBColor(0xB8, 0x9A, 0x6A),
    'gold_dark':      RGBColor(0xBF, 0x90, 0x00),
    'text_dark':      RGBColor(0x10, 0x18, 0x28),
    'text_main':      RGBColor(0x36, 0x41, 0x53),
    'text_secondary': RGBColor(0x4A, 0x55, 0x65),
    'white':          RGBColor(0xFF, 0xFF, 0xFF),
    'bg_card':        RGBColor(0xF9, 0xFA, 0xFB),
    'black':          RGBColor(0x00, 0x00, 0x00),
    'light_gray':     RGBColor(0xE0, 0xE0, 0xE0),
    'bg_blue':        RGBColor(0xE8, 0xF0, 0xFE),
    'bg_green':       RGBColor(0xE0, 0xF7, 0xE8),
    'bg_orange':      RGBColor(0xFE, 0xF0, 0xE0),
}

OUTPUT_DIR = "/Users/Admin/OpencodeWorkspace/内容输出/AI应用层投资机会"
LOGO = os.path.join(OUTPUT_DIR, "huihong_logo.png")
prs_width = Inches(13.333)
prs_height = Inches(7.5)

def _scf(run, cn="微软雅黑", en="Franklin Gothic Medium"):
    run.font.name = en
    rPr = run._r.get_or_add_rPr()
    lat = rPr.find(qn("a:latin"))
    if lat is None:
        lat = rPr.makeelement(qn("a:latin"), {}); rPr.append(lat)
    lat.set("typeface", en)
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {}); rPr.append(ea)
    ea.set("typeface", cn)

def fmt(run, size=14, bold=False, color=C["text_main"]):
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color
    _scf(run)

def ap(tf, txt, size=14, bold=False, color=C["text_main"], align=None, sb=0, sa=0):
    p = tf.add_paragraph()
    r = p.add_run(); r.text = txt; fmt(r, size=size, bold=bold, color=color)
    if align: p.alignment = align
    if sb: p.space_before = Pt(sb)
    if sa: p.space_after = Pt(sa)
    return p

def ap0(tf, txt, size=14, bold=False, color=C["text_main"], align=None):
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = txt; fmt(r, size=size, bold=bold, color=color)
    if align: p.alignment = align
    return p

def _rr(slide, left, top, w, h, fill=None, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else: sp.fill.background()
    if line: sp.line.color.rgb = line; sp.line.width = Pt(0.5)
    else: sp.line.fill.background()
    sp.adjustments[0] = 0.05
    return sp

def _rc(slide, left, top, w, h, fill=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else: sp.fill.background()
    sp.line.fill.background()
    return sp

def _ln(slide, left, top, w, color=C["brand_red"], width=Pt(1.5)):
    ln = slide.shapes.add_connector(1, left, top, left+w, top)
    ln.line.color.rgb = color; ln.line.width = width
    return ln

def _pg(slide, num):
    tx = slide.shapes.add_textbox(Inches(12.35), Inches(7.12), Inches(0.7), Inches(0.30))
    tx.text_frame.word_wrap = True
    p = tx.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(num); fmt(r, size=11, color=C["text_secondary"])

def _hb(slide):
    _rc(slide, Inches(0.642), Inches(0.483), Inches(0.465), Inches(0.464), fill=C["white"])
    _rc(slide, Inches(0.875), Inches(0.684), Inches(0.347), Inches(0.347), fill=C["brand_red"])
    _ln(slide, Inches(1.300), Inches(1.005), Inches(11.116), color=C["brand_red"], width=Pt(1.5))

def _ht(slide, t):
    tx = slide.shapes.add_textbox(Inches(1.302), Inches(0.48), Inches(11.114), Inches(0.45))
    p = tx.text_frame.paragraphs[0]; p.text = t
    p.font.size = Pt(22); p.font.bold = False; p.font.color.rgb = C["brand_red"]
    for r in p.runs: _scf(r)

def _logo(slide):
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(11.65), Inches(0.08), Inches(0.89), Inches(0.89))

def ss(slide, num):
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = C["white"]
    _hb(slide); _logo(slide); _pg(slide, num)


# ===== SLIDE DATA =====
SLIDES = [

    # 1 - Cover
    {
        "type": "cover",
        "title": 'AI应用层投资机会',
        "subtitle": '从编程、法律到生物科技，挖掘真实营收赛道',
    },

    # 2 - Core judgement
    {
        "type": "big_quote",
        "title": '核心判断',
        "quote": 'AI应用层不是“没有需求”，而是“价值分布极度不均”。编程已经被证明，法律正在被证明，生物科技将被证明。',
        "sub_quote": '投资人问“AI有没有价值”是错的问题。正确的问题是：“AI的价值在哪层？”',
    },

    # 3 - Market overview
    {
        "type": "side_by_side",
        "title": '全球AI应用层市场概览',
        "left_title": '市场规模（亿美元）',
        "left_items": [
            ('云AI市场','领先2027年$168B，CAGR 37%'),
            ('企业AI应用','2027年$99B，全球IT支出约为$4.7T'),
            ('生成式AI','2025年$67B，2030年$1.3T（Bloomberg）'),
            ('中国AI','2024年$34B，2027年$66B CAGR 25%'),
        ],
        "right_title": '渗透率分析',
        "right_items": [
            ('全球AI渗透率','今年~6%，预计2030年~25%'),
            ('普及者奖励','早期采用者可带来显著竞争优势'),
            ('中国市场特点','资金靠政府担保贷款 + 科创板企业'),
        ],
    },

    # 4 - Section: Programming
    {
        "type": "section_header",
        "title": '赛道一：编程——已被证明的黄金赛道',
    },

    # 5 - Programming data cards
    {
        "type": "data_cards",
        "title": '编程赛道核心数据',
        "cards": [
            ('Cursor', '$2B ARR（2025）', '50万付费用户，年增速>200%', C['bg_blue']),
            ('GitHub Copilot', '$1.5B营收', '累计180万付费用户，展望$3B+', C['bg_green']),
            ('Devin定价跳水', '$500→$20/月', 'AI程序员“大跳水”，Claude Code$17/月', C['bg_orange']),
            ('开发者采用', '85%已使用AI', 'JetBrains 2026：从“Vibe Coding”→“Agentic Engineering”。编程门槛消失，定价权下移', C['bg_blue']),
        ],
    },

    # 6 - Programming insight
    {
        "type": "insight",
        "title": '编程赛道投资见解',
        "body": '① 工具层已经激烈：Cursor、Copilot、Windsurf等红海，但业务场景层仍有空白。\n② AI编程从“下方AI”——你不需要购买AI，你只需要85%的AI能力。\n③ 中国编程AI市场规模2030年有望超$10B，但需要找到类似“解决以前解决不了的问题”的场景。',
    },

    # 7 - Section: Legal
    {
        "type": "section_header",
        "title": '赛道二：法律——正在被证明的下一个黄金赛道',
    },

    # 8 - Legal data cards
    {
        "type": "data_cards",
        "title": '法律赛道核心数据',
        "cards": [
            ('全球法律AI市场', '2030年$50B+', '全球法律服务市场$900B，AI渗透率~5.5%', C['bg_blue']),
            ('Harvey AI', '$300M资金+$100M营收', '全球顶级律所客户，多语言注意力管理、庭诊准备', C['bg_green']),
            ('Ironclad $150M', 'CLM领域AI化', '合同生命周期管理，从审查到签署全流程', C['bg_orange']),
            ('中国法律AI', '海国数字$10B+市值', '作为标杆，法律AI主要渗透在合规、合同管理、分析', C['bg_green']),
        ],
    },

    # 9 - Legal insight
    {
        "type": "insight",
        "title": '法律赛道投资见解',
        "body": '① 法律AI是典型的“锦上添花”（让好的变得更好）：大律所人工成本高，AI可大幅降低。\n② 法律AI的终极形态是摇号：商业律师将被重新定义，但不会被取代。\n③ 物流、金融、医疗等大宗合同场景是最佳落地场景。',
    },

    # 10 - Section: Bio
    {
        "type": "section_header",
        "title": '赛道三：生物科技——将被证明的未来巨头',
    },

    # 11 - Bio data cards
    {
        "type": "data_cards",
        "title": '生物科技AI核心数据',
        "cards": [
            ('全球AI医药市场', '2030年$67B', 'CAGR 42%，线上从药物发现到临床试验优化', C['bg_blue']),
            ('Recursion Pharma', '$10B+市值', '蓬莱屏障技术，线上AI筛选匹配药物提前十年', C['bg_green']),
            ('Isomorphic Labs', 'DeepMind分支', 'AlphaFold3升级，与药品企业合作增强AI发现药物', C['bg_orange']),
            ('中国生物科技AI', '先发药物$2B市值', '成为中国AI药物发现标杆', C['bg_green']),
        ],
    },

    # 12 - Bio insight
    {
        "type": "insight",
        "title": '生物科技赛道投资见解',
        "body": '① 生物科技AI从“讲故事”阶段转向“出药阶段”。上市药物AI发现的成分将宣告新时代。\n② 药物发现周期可从10年缩短到3-4年，成本从$2.6B降低到$500M以下。\n③ 早期药物发现是产业链中AI价值最大的环节。',
    },

    # 13 - Summary
    {
        "type": "summary",
        "title": '投资策略总结',
        "items": [
            ('编程：已被证明','工具层红海，业务场景层有机会。重点关注企业尺度AI编程、安全性与Theory of Mind系统。',C['brand_red']),
            ('法律：正在被证明','大律所与中小律所的AI应用差异将成为洗牌关键。注意力管理、合同审查是最佳落地场景。',C['dark_red']),
            ('生物科技：将被证明','从“讲故事”到“出药”的转变将是最大的投资机会。药物发现是价值最大的环节。',C['gold_warm']),
        ],
    },

    # 14 - Closing
    {
        "type": "closing",
        "title": 'THANK YOU',
        "subtitle": '汇竑股权投资，助力创新企业成长',
    },
]

# ===== BUILD FUNCTIONS =====

def build_cover(slide, data):
    ss(slide, 1)
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.0))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["title"], size=44, bold=True, color=C["brand_red"], align=PP_ALIGN.CENTER)
    ap(tf, data["subtitle"], size=18, bold=False, color=C["text_secondary"], align=PP_ALIGN.CENTER, sb=12)
    _ln(slide, Inches(5.5), Inches(4.8), Inches(2.3), color=C["gold"], width=Pt(2.5))

def build_big_quote(slide, data):
    ss(slide, 2)
    _ht(slide, data["title"])
    _rr(slide, Inches(1.3), Inches(1.8), Inches(10.7), Inches(2.5), fill=RGBColor(0xFD,0xF0,0xF0), line=C["brand_red"])
    tx = slide.shapes.add_textbox(Inches(1.6), Inches(2.0), Inches(10.1), Inches(2.0))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["quote"], size=18, bold=True, color=C["brand_red"], align=PP_ALIGN.LEFT)
    if "sub_quote" in data:
        ap(tf, "", size=10)
        ap(tf, data["sub_quote"], size=14, bold=False, color=C["text_secondary"], align=PP_ALIGN.LEFT)

def build_side_by_side(slide, data):
    ss(slide, 3)
    _ht(slide, data["title"])
    mid = Inches(6.6)
    _rr(slide, Inches(1.3), Inches(1.6), Inches(5.1), Inches(5.2), fill=C["bg_card"], line=C["light_gray"])
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(1.7), Inches(4.7), Inches(5.0))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["left_title"], size=16, bold=True, color=C["brand_red"])
    for item in data["left_items"]:
        ap(tf, "\u25b8 " + item[0], size=13, bold=True, color=C["text_dark"], sb=8)
        ap(tf, "   " + item[1], size=11, bold=False, color=C["text_secondary"], sb=2)
    _rr(slide, mid, Inches(1.6), Inches(5.1), Inches(5.2), fill=C["bg_card"], line=C["light_gray"])
    tx = slide.shapes.add_textbox(Inches(6.8), Inches(1.7), Inches(4.7), Inches(5.0))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["right_title"], size=16, bold=True, color=C["brand_red"])
    for item in data["right_items"]:
        ap(tf, "\u25b8 " + item[0], size=13, bold=True, color=C["text_dark"], sb=8)
        ap(tf, "   " + item[1], size=11, bold=False, color=C["text_secondary"], sb=2)

def build_section_header(slide, data):
    ss(slide, 0)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x28)
    _ln(slide, Inches(5.0), Inches(3.2), Inches(3.3), color=C["gold"], width=Pt(2.5))
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.5))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["title"], size=36, bold=True, color=C["gold"], align=PP_ALIGN.CENTER)

def build_data_cards(slide, data, page_num):
    ss(slide, page_num)
    _ht(slide, data["title"])
    cards = data["cards"]
    n = len(cards)
    total_w = Inches(12.0)
    gap = Inches(0.15)
    card_w = (total_w - gap * (n - 1)) / n
    for i, cd in enumerate(cards):
        left = Inches(0.65) + i * (card_w + gap)
        top = Inches(1.7)
        title, metric, desc, bg = cd
        _rr(slide, left, top, card_w, Inches(5.0), fill=bg, line=RGBColor(0xD0,0xD0,0xD0))
        tx = slide.shapes.add_textbox(left+Inches(0.12), top+Inches(0.2), card_w-Inches(0.24), Inches(0.7))
        tf = tx.text_frame; tf.word_wrap = True
        ap0(tf, metric, size=22, bold=True, color=C["brand_red"])
        tx = slide.shapes.add_textbox(left+Inches(0.12), top+Inches(0.95), card_w-Inches(0.24), Inches(0.45))
        tf = tx.text_frame; tf.word_wrap = True
        ap0(tf, title, size=13, bold=True, color=C["text_dark"])
        tx = slide.shapes.add_textbox(left+Inches(0.12), top+Inches(1.5), card_w-Inches(0.24), Inches(3.2))
        tf = tx.text_frame; tf.word_wrap = True
        ap0(tf, desc, size=11, bold=False, color=C["text_main"])

def build_insight(slide, data, page_num):
    ss(slide, page_num)
    _ht(slide, data["title"])
    _rr(slide, Inches(1.3), Inches(1.8), Inches(10.7), Inches(5.0), fill=C["bg_card"], line=C["light_gray"])
    tx = slide.shapes.add_textbox(Inches(1.6), Inches(2.0), Inches(10.1), Inches(4.5))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["body"], size=14, bold=False, color=C["text_main"])

def build_summary(slide, data, page_num):
    ss(slide, page_num)
    _ht(slide, data["title"])
    y_start = Inches(1.8)
    for i, item in enumerate(data["items"]):
        title, desc, color = item
        left = Inches(1.3); w = Inches(10.7); h = Inches(1.5)
        top = y_start + i * (h + Inches(0.15))
        _rr(slide, left, top, w, h, fill=C["bg_card"], line=C["light_gray"])
        tx = slide.shapes.add_textbox(left+Inches(0.2), top+Inches(0.1), w-Inches(0.4), h-Inches(0.2))
        tf = tx.text_frame; tf.word_wrap = True
        ap0(tf, title, size=14, bold=True, color=color)
        ap(tf, desc, size=11, bold=False, color=C["text_main"], sb=4)

def build_closing(slide, data):
    ss(slide, 0)
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = RGBColor(0x1C, 0x1C, 0x28)
    _ln(slide, Inches(5.0), Inches(3.2), Inches(3.3), color=C["gold"], width=Pt(2.5))
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.3), Inches(1.0))
    tf = tx.text_frame; tf.word_wrap = True
    ap0(tf, data["title"], size=44, bold=True, color=C["gold"], align=PP_ALIGN.CENTER)
    ap(tf, data["subtitle"], size=16, bold=False, color=C["text_secondary"], align=PP_ALIGN.CENTER, sb=8)

def build_slide(slide, data, page_num):
    t = data["type"]
    if t == "cover": build_cover(slide, data)
    elif t == "big_quote": build_big_quote(slide, data)
    elif t == "side_by_side": build_side_by_side(slide, data)
    elif t == "section_header": build_section_header(slide, data)
    elif t == "data_cards": build_data_cards(slide, data, page_num)
    elif t == "insight": build_insight(slide, data, page_num)
    elif t == "summary": build_summary(slide, data, page_num)
    elif t == "closing": build_closing(slide, data)

def main():
    prs = Presentation()
    prs.slide_width = prs_width
    prs.slide_height = prs_height
    blank = prs.slide_layouts[6]
    for i, data in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        build_slide(slide, data, i + 1)
    out_path = os.path.join(OUTPUT_DIR, "AI应用层投资机会.pptx")
    prs.save(out_path)
    print(f"PPT generated: {out_path}")
    print(f"Total {len(SLIDES)} slides")

if __name__ == "__main__":
    main()