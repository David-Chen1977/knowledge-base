#!/usr/bin/env python3
"""
智算中心节能用能体系 · 专业PPT报告生成器
使用 python-pptx 创建符合品牌视觉规范的高质量演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 品牌配色 ──
BRAND = RGBColor(0xC0, 0x00, 0x00)       # 深红 #c00000
BRAND_LIGHT = RGBColor(0xFE, 0xF2, 0xF2)  # 浅红背景
BRAND_DARK = RGBColor(0x80, 0x00, 0x00)    # 深红文字
DARK = RGBColor(0x1E, 0x1E, 0x1E)          # 主文字色
GRAY = RGBColor(0x66, 0x66, 0x66)          # 辅助文字
LIGHT_GRAY = RGBColor(0xBF, 0xBF, 0xBF)    # 边框
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF7, 0xF8, 0xFA)      # 卡片背景
POSITIVE = RGBColor(0x2E, 0x7D, 0x32)       # 正/绿色

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── 辅助函数 ──

def add_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def set_text(shape, text, size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """设置形状文字"""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p

def add_text_box(slide, left, top, width, height, text, size=14, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """添加文字框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, default_size=14, default_color=DARK):
    """添加多行文字"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(4)
    return txBox

def add_card_table(slide, left, top, width, headers, rows, col_widths=None):
    """添加卡片风格的表格"""
    row_count = len(rows) + 1
    col_count = len(headers)
    table_shape = slide.shapes.add_table(row_count, col_count, left, top, width, Inches(0.4 * row_count))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND

    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else BG_LIGHT

    return table_shape


# ════════════════════════════════════════════════════════════════
#  SLIDE 1: 封面
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BRAND)

# 主标题
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "智算中心节能用能体系", size=44, color=WHITE, bold=True)

# 副标题
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "从成本中心到竞争力核心", size=28, color=WHITE, bold=False)

# 装饰线
add_shape(slide, Inches(1), Inches(4.2), Inches(3), Inches(0.06), WHITE)

# 署名
add_text_box(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.6),
             "陈道雷 · 汇竑资本 | 2026年6月", size=16, color=WHITE)

# 底部说明
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "PE/VC 一级股权投资 · 算电协同赛道深度研究 · 机密", size=12, color=RGBColor(0xFF, 0xAA, 0xAA))


# ════════════════════════════════════════════════════════════════
#  SLIDE 2: 目录
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.6),
             "目 录", size=32, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.04), BRAND)

toc_items = [
    ("01", "三个必须面对的现实"),
    ("02", "三层架构：系统设计"),
    ("03", "产业链图谱与价值分布"),
    ("04", "五种商业模式评估"),
    ("05", "投资时机与组合策略"),
    ("06", "核心判断与行动建议"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.8) + Inches(0.75) * i
    # 数字
    add_text_box(slide, Inches(1.5), y, Inches(1), Inches(0.5),
                 num, size=28, color=BRAND, bold=True)
    # 文字
    add_text_box(slide, Inches(2.8), y + Inches(0.05), Inches(6), Inches(0.5),
                 title, size=20, color=DARK)
    # 分割线
    add_shape(slide, Inches(1.5), y + Inches(0.55), Inches(8), Inches(0.01), LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
#  SLIDE 3: 三个现实
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "三个必须面对的现实", size=30, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

# 三个卡片
cards = [
    ("⚡ 算力爆发", "2026年数据中心用电量 ~3000亿kWh\n2030年或超7000亿kWh\nCAGR 15-20%\n→ 增长最快的单一用电大户"),
    ("📋 政策硬约束", "PUE < 1.25 + 绿电 ≥ 80%\n国家枢纽节点新建数据中心红线\n→ 节能从\"加分项\"变\"准入门槛\""),
    ("💰 电费吃利润", "千P智算中心年电费 1.5-3亿\n省20% → 3000-6000万利润释放\nIRR 8% → 12%\n→ 直接决定项目回报"),
]

card_w = Inches(3.6)
card_h = Inches(4.5)
gap = Inches(0.5)
start_x = Inches(0.8)

for i, (title, body) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    y = Inches(1.6)

    # 卡片背景
    card = add_rounded_rect(slide, x, y, card_w, card_h, BG_LIGHT)

    # 标题
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), card_w - Inches(0.6), Inches(0.6),
                 title, size=18, color=BRAND, bold=True)

    # 分隔线
    add_shape(slide, x + Inches(0.3), y + Inches(1.0), Inches(0.8), Inches(0.03), BRAND)

    # 内容
    lines = body.split('\n')
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.3), card_w - Inches(0.6), card_h - Inches(1.6),
                   [(l, 13, DARK, False) for l in lines])

# 底部结论
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "▶ 三个现实同步发生 → 节能用能体系不再是\"社会责任\"，是竞争力核心", size=16, color=BRAND, bold=True)


# ════════════════════════════════════════════════════════════════
#  SLIDE 4: 三层架构总览
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "三层架构：节能是个系统，不是单点", size=30, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

# 三层金字塔（从下到上）
layers = [
    ("IT 设备级能效", "低功耗芯片 + 液冷散热\n每瓦特产出最多 Token", "6-12月见效", Inches(3.6), BRAND),
    ("设施级能效", "AI 制冷优化 + 储能系统\nPUE 从 1.5 到 1.15", "1-2年见效", Inches(2.8), RGBColor(0xE0, 0x60, 0x30)),
    ("系统级能效", "绿电直供 + VPP现货 + 算力调度\n从\"被动用电\"到\"主动管电\"", "2-3年见效", Inches(2.0), RGBColor(0x40, 0x80, 0xC0)),
]

base_y = Inches(5.0)
for i, (title, desc, timeline, height, color) in enumerate(layers):
    y = base_y - height - Inches(0.3) * i
    w = Inches(8) - Inches(1.5) * i
    x = (SLIDE_W - w) / 2

    layer_shape = add_rounded_rect(slide, x, y, w, height, color)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.3), w - Inches(1), Inches(0.5),
                 title, size=20, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.9), w - Inches(1), Inches(0.7),
                 desc, size=13, color=WHITE)
    add_text_box(slide, x + Inches(0.5), y + height - Inches(0.5), w - Inches(1), Inches(0.4),
                 f"⏱ {timeline}", size=11, color=RGBColor(0xFF, 0xDD, 0xDD))

# 右侧效果标注
add_rounded_rect(slide, Inches(10), Inches(2.0), Inches(2.8), Inches(2.0), RGBColor(0x2E, 0x7D, 0x32))
add_text_box(slide, Inches(10.3), Inches(2.3), Inches(2.2), Inches(0.5),
             "综合效果", size=20, color=WHITE, bold=True)
add_text_box(slide, Inches(10.3), Inches(2.8), Inches(2.2), Inches(1.0),
             "综合用电成本\n可降 20-40%", size=18, color=WHITE, bold=True)


# ════════════════════════════════════════════════════════════════
#  SLIDE 5: 第一层 - IT设备级
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
             "第一层 · IT设备级", size=14, color=GRAY)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
             "每瓦特产出最多 Token", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.04), BRAND)

# 左: 低功耗芯片
card1 = add_rounded_rect(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), BG_LIGHT)
add_text_box(slide, Inches(1.2), Inches(2.3), Inches(4.5), Inches(0.5),
             "🔲 低功耗推理芯片", size=20, color=BRAND, bold=True)
add_shape(slide, Inches(1.2), Inches(2.9), Inches(1.5), Inches(0.03), BRAND)

chip_lines = [
    ("推理专用 ASIC 比通用 GPU 节能 5-10 倍", 14, DARK, True),
    ("", 6, DARK, False),
    ("2026年最值得关注的不是训练芯片，是推理芯片", 14, DARK, False),
    ("推理是长期、持续、海量的——这才是真正的\"用量\"", 14, GRAY, False),
    ("", 6, DARK, False),
    ("▶ 推理芯片将是2026-2028年增长速度最快的细分", 13, BRAND, True),
]
add_multi_text(slide, Inches(1.2), Inches(3.2), Inches(4.8), Inches(3.0), chip_lines)

# 右: 液冷散热
card2 = add_rounded_rect(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(4.5), BG_LIGHT)
add_text_box(slide, Inches(7.2), Inches(2.3), Inches(5), Inches(0.5),
             "🌊 液冷散热", size=20, color=BRAND, bold=True)
add_shape(slide, Inches(7.2), Inches(2.9), Inches(1.5), Inches(0.03), BRAND)

liquid_lines = [
    ("冷板式液冷  PUE < 1.20  ✅ 已大规模商用", 14, DARK, True),
    ("浸没式液冷  PUE < 1.05  🚀 快速起量", 14, DARK, True),
    ("", 6, DARK, False),
    ("NVIDIA Rubin 系列 (2026+) → 100% 全液冷", 14, BRAND, True),
    ("整个产业链被推向液冷路线", 13, GRAY, False),
    ("", 6, DARK, False),
    ("📊 2024年液冷市场 ¥172.72亿 (+67%)", 14, DARK, True),
    ("📊 2025年新建DC液冷渗透率 > 50%", 14, DARK, True),
    ("", 6, DARK, False),
    ("▶ 液冷从\"可选\"变成\"标配\"——时间点就是现在", 13, BRAND, True),
]
add_multi_text(slide, Inches(7.2), Inches(3.2), Inches(5), Inches(3.0), liquid_lines)

# 底部
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4),
             "📌 关键转折：液冷不再是差异化优势，已经是行业准入门槛", size=14, color=POSITIVE, bold=True)


# ════════════════════════════════════════════════════════════════
#  SLIDE 6: 第二层 - 设施级 AI制冷
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
             "第二层 · 设施级能效", size=14, color=GRAY)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
             "AI 制冷优化：最大被低估的投资机会", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.04), BRAND)

# 核心数据 - 大字
add_text_box(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(0.8),
             "制冷占数据中心能耗 30-40% → AI 优化可降低 10-30%", size=22, color=DARK, bold=True)

# 左侧: 海外对标
left_card = add_rounded_rect(slide, Inches(0.8), Inches(2.8), Inches(5.8), Inches(3.8), BG_LIGHT)
add_text_box(slide, Inches(1.2), Inches(3.0), Inches(5), Inches(0.5),
             "🌍 海外对标：Phaidra", size=20, color=BRAND, bold=True)
add_shape(slide, Inches(1.2), Inches(3.5), Inches(1.5), Inches(0.03), BRAND)

phaidra_lines = [
    ("AI Agent 优化数据中心制冷系统", 14, DARK, True),
    ("", 6, DARK, False),
    ("Google 冷冻水系统 → 制冷能耗降低 30%", 14, BRAND, True),
    ("NVIDIA 液冷系统 → 热尖峰减少 ~80%", 14, BRAND, True),
    ("", 6, DARK, False),
    ("毛利率 60-80%，纯 SaaS 轻资产模式", 14, POSITIVE, True),
    ("模式已验证，独立公司运营中", 13, GRAY, False),
]
add_multi_text(slide, Inches(1.2), Inches(3.7), Inches(5), Inches(2.5), phaidra_lines)

# 右侧: 国内空白 = 机会
right_card = add_rounded_rect(slide, Inches(7.0), Inches(2.8), Inches(5.5), Inches(3.8), RGBColor(0xFF, 0xF0, 0xF0))
add_text_box(slide, Inches(7.4), Inches(3.0), Inches(4.5), Inches(0.5),
             "🇨🇳 国内：赛道空白", size=20, color=BRAND, bold=True)
add_shape(slide, Inches(7.4), Inches(3.5), Inches(1.5), Inches(0.03), BRAND)

cn_lines = [
    ("❌ 没有一家专注 AI 制冷优化的独立公司", 16, BRAND, True),
    ("", 8, DARK, False),
    ("这意味着：", 14, DARK, True),
    ("", 6, DARK, False),
    ("选择 A：投资海外 Phaidra", 14, DARK, False),
    ("选择 B：国内孵化一个", 16, BRAND, True),
    ("", 6, DARK, False),
    ("💰 国内 TAM 预估：≥500座大型DC", 14, DARK, True),
    ("    每年 ¥50-100万/座 = ¥2.5-5亿可寻址市场", 13, GRAY, False),
]
add_multi_text(slide, Inches(7.4), Inches(3.7), Inches(4.8), Inches(2.5), cn_lines)

# 底部
add_text_box(slide, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
             "▶ 行动建议：本月调研国内是否可孵化，这是最适合 PE/VC 的轻资产模式", size=14, color=BRAND, bold=True)


# ════════════════════════════════════════════════════════════════
#  SLIDE 7: 第三层 - 系统级
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
             "第三层 · 系统级能效", size=14, color=GRAY)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
             "从\"被动用电\"到\"主动管电\"", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.04), BRAND)

# 三个竖列
cols = [
    ("☀️ 绿电直供", [
        "通过 PPA 或专线直购新能源",
        "度电成本 0.2-0.3元（西部）",
        "vs 0.7-0.8元（东部工商业）",
        "",
        "案例：乌兰察布中金数据",
        "45MW/180MWh 源网荷储一体化",
        "绿电替代率 38.74%",
        "年自发自用 8.48亿kWh",
    ]),
    ("⚡ VPP + 现货交易", [
        "韶关联通 (2026.5)：国内首单",
        "数据中心 = 电网\"柔性调节器\"",
        "电价低→多算，电价高→少算",
        "",
        "📌 按钮时刻",
        "从\"理论上可行\" → \"已跑通\"",
        "千亿级市场规模正在打开",
        "适合平台型投资布局",
    ]),
    ("🔄 算力随电价调度", [
        "训练任务可中断/迁移/等待",
        "非紧急 → 移至西部低价时段",
        "推理 → 留在用户侧低延迟",
        "",
        "时空调度的经济价值巨大",
        "跨区域算力调度平台兴起",
        "效率提升 20-40%",
        "适合 PE/VC 轻资产布局",
    ]),
]

col_w = Inches(3.8)
for i, (title, items) in enumerate(cols):
    x = Inches(0.8) + (col_w + Inches(0.3)) * i
    y = Inches(1.9)

    card = add_rounded_rect(slide, x, y, col_w, Inches(5.0), BG_LIGHT)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), col_w - Inches(0.6), Inches(0.5),
                 title, size=18, color=BRAND, bold=True)
    add_shape(slide, x + Inches(0.3), y + Inches(0.75), Inches(1), Inches(0.03), BRAND)

    item_lines = [(item, 12, DARK, False) for item in items]
    add_multi_text(slide, x + Inches(0.3), y + Inches(1.0), col_w - Inches(0.6), Inches(3.8), item_lines)


# ════════════════════════════════════════════════════════════════
#  SLIDE 8: 产业链图谱
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "产业链：七个环节，一个被低估的价值高地", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

headers = ["环节", "毛利率", "资本密度", "PE/VC适合度", "代表公司"]
rows = [
    ["AI 低功耗芯片", "50-60%", "极高", "⭐⭐⭐", "寒武纪/算能"],
    ["液冷散热系统", "25-40%", "中", "⭐⭐⭐", "英维克/高澜/曙光数创"],
    ["高效供电配电", "20-30%", "中", "⭐⭐", "科华数据"],
    ["储能系统", "15-25%", "高", "⭐⭐", "宁德/阳光电源"],
    ["智算中心运营", "30-50%", "极高", "⭐⭐⭐", "万国数据/世纪互联"],
    ["能效优化服务", "60-80%", "极低", "⭐⭐⭐⭐⭐", "国内空白"],
    ["下游算力消费", "40-70%", "中", "⭐⭐⭐⭐", "阿里云/字节跳动"],
]

tbl = add_card_table(slide, Inches(0.8), Inches(1.4), Inches(11.5), headers, rows,
                     col_widths=[Inches(2.5), Inches(1.8), Inches(1.8), Inches(2.5), Inches(2.9)])

# 高亮能效优化行
# (table rows: row 0 = header, rows 1-7 = data, row 6 = 能效优化服务)
for cell in tbl.table.rows[6].cells:
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = BRAND
        p.font.bold = True

# 底部判断
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.8),
             "核心判断：价值最高的环节不在设备和运营，在能效优化服务（毛利率60-80%，轻资产可复制）",
             size=16, color=BRAND, bold=True)

# 竞争格局
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.6),
             "液冷散热三龙头占据大部分市场（英维克/高澜/曙光数创）| 能效SaaS赛道国内空白 → 最佳布局窗口",
             size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════
#  SLIDE 9: 五种商业模式（横向对比）
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "五种商业模式：你的钱该放在哪里？", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

headers2 = ["模式", "描述", "资产", "IRR", "PE/VC适合度"]
rows2 = [
    ["A 源网荷储一体化", "自建风光+储能+数据中心", "重资产", "6-10%", "❌ 适合产业资本"],
    ["B 分布式算力仓", "绿电富集区模块化DC部署", "中资产", "12-18%", "✅ 有资源玩家"],
    ["C 能效优化SaaS", "AI优化制冷/用电/交易", "轻资产", "30%+", "✅✅ 最适PE/VC"],
    ["D VPP聚合+交易", "聚合DC负荷参与电力市场", "轻资产", "20-30%", "✅✅ 窗口期"],
    ["E 存量改造", "老旧DC液冷+能效升级", "中资产", "15-25%", "✅ 政策驱动"],
]

tbl2 = add_card_table(slide, Inches(0.8), Inches(1.4), Inches(11.5), headers2, rows2,
                      col_widths=[Inches(2.8), Inches(3.5), Inches(1.5), Inches(1.5), Inches(2.2)])

# 高亮模式C和D
for r_idx in [2, 3]:  # row 3 = 模式C, row 4 = 模式D
    for cell in tbl2.table.rows[r_idx].cells:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True

# 底部 - 关键洞察
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.8),
             "▶ 模式 C（能效SaaS）和模式 D（VPP聚合）最适合 PE/VC：轻资产、高壁垒、毛利率 60-80%",
             size=16, color=BRAND, bold=True)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.6),
             "重资产模式（A/E）适合产业资本和央企，财务投资人应优先考虑平台型和 SaaS 模式",
             size=13, color=GRAY)


# ════════════════════════════════════════════════════════════════
#  SLIDE 10: 模式C深度
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.5),
             "深度拆解 · 模式 C", size=14, color=GRAY)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
             "能效优化 SaaS —— 最适合 PE/VC 的标的", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.04), BRAND)

# 左侧：商业模式画布
left = add_rounded_rect(slide, Inches(0.8), Inches(1.9), Inches(6.0), Inches(4.8), BG_LIGHT)
add_text_box(slide, Inches(1.2), Inches(2.1), Inches(5), Inches(0.5),
             "📋 商业模式画布", size=18, color=BRAND, bold=True)

canvas_items = [
    ("客户", "智算中心运营方（万国数据/世纪互联/秦淮数据）"),
    ("价值主张", "省电 10-30% = 每座 DC 每年节省 ¥1000-3000万"),
    ("收费模式", "按节省电费分成 20-30% / SaaS 年订阅 ¥50-100万"),
    ("核心壁垒", "算法精度 + 跨平台兼容 + 数据飞轮效应"),
    ("增长引擎", "每多一个客户 → 多一组训练数据 → 模型更准"),
]
for i, (label, value) in enumerate(canvas_items):
    y = Inches(2.8) + Inches(0.65) * i
    add_text_box(slide, Inches(1.2), y, Inches(1.5), Inches(0.5),
                 label, size=13, color=BRAND, bold=True)
    add_text_box(slide, Inches(2.8), y, Inches(3.8), Inches(0.5),
                 value, size=12, color=DARK)

# 右侧：市场预估
right = add_rounded_rect(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(4.8), RGBColor(0xFF, 0xF0, 0xF0))
add_text_box(slide, Inches(7.6), Inches(2.1), Inches(4.5), Inches(0.5),
             "💰 市场预估", size=18, color=BRAND, bold=True)

market_lines = [
    ("TAM（总可寻址市场）", 14, BRAND, True),
    ("国内 ≥500 座大型数据中心", 13, DARK, False),
    ("每座年费 ¥50-100万 → ¥2.5-5亿", 13, DARK, False),
    ("", 8, DARK, False),
    ("SAM（可服务市场）", 14, BRAND, True),
    ("目标前 30% 高端客户 → ¥0.75-1.5亿", 13, DARK, False),
    ("", 8, DARK, False),
    ("SOM（可获得市场）", 14, BRAND, True),
    ("3年目标 50座 DC → ¥2500-5000万 ARR", 13, DARK, False),
    ("", 8, DARK, False),
    ("毛利率 60-80%", 16, POSITIVE, True),
    ("▶ 海外已验证，国内空白，3年窗口期", 14, BRAND, True),
]
add_multi_text(slide, Inches(7.6), Inches(2.7), Inches(4.5), Inches(3.8), market_lines)

# 底部
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.4),
             "▶ 行动：本月启动国内对标公司调研，如空白则评估内部孵化可行性", size=14, color=BRAND, bold=True)


# ════════════════════════════════════════════════════════════════
#  SLIDE 11: 投资时机判断
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "投资时机：窗口期正在打开", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

# 时间轴
timeline_y = Inches(2.0)
add_text_box(slide, Inches(0.8), timeline_y - Inches(0.4), Inches(11), Inches(0.4),
             "赛道成熟度时间轴", size=16, color=GRAY)

year_positions = [(Inches(1.5), "2024"), (Inches(4.0), "2025"), (Inches(6.5), "2026 🔥"), (Inches(9.0), "2027"), (Inches(11.5), "2028")]
for x, year in year_positions:
    add_text_box(slide, x, timeline_y, Inches(1.2), Inches(0.4),
                 year, size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

# 时间轴线条
add_shape(slide, Inches(1.5), timeline_y + Inches(0.5), Inches(10.8), Inches(0.03), BRAND)

# 节点
node_labels = [
    (Inches(1.5), "概念期\n液冷方案POC"),
    (Inches(4.0), "早期\n液冷渗透50%"),
    (Inches(6.5), "成长期 🔥\nVPP首单落地\n政策硬约束"),
    (Inches(9.0), "快速扩散\n并购开始出现"),
    (Inches(11.5), "成熟\n竞争白热化\n退出窗口"),
]
for x, label in node_labels:
    # 圆点
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.35), timeline_y + Inches(0.42), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BRAND
    dot.line.fill.background()
    # 标签
    add_text_box(slide, x - Inches(0.4), timeline_y + Inches(0.8), Inches(1.8), Inches(1.0),
                 label, size=10, color=DARK, alignment=PP_ALIGN.CENTER)

# 判断卡片
judge_y = Inches(4.5)
judges = [
    ("🟢 液冷渗透 > 50%", "确定性最强，竞争最激烈"),
    ("🟢 VPP 首单落地", "模式已跑通，赶紧布局"),
    ("🔵 AI制冷优化国内空白", "要么第一个进，要么等 Phaidra"),
]

for i, (title, desc) in enumerate(judges):
    x = Inches(0.8) + Inches(4.0) * i
    card = add_rounded_rect(slide, x, judge_y, Inches(3.6), Inches(1.5), BG_LIGHT)
    add_text_box(slide, x + Inches(0.3), judge_y + Inches(0.2), Inches(3), Inches(0.4),
                 title, size=16, color=BRAND, bold=True)
    add_text_box(slide, x + Inches(0.3), judge_y + Inches(0.7), Inches(3), Inches(0.6),
                 desc, size=12, color=DARK)

# 结论
add_rounded_rect(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6), BRAND)
add_text_box(slide, Inches(1.2), Inches(6.55), Inches(10.5), Inches(0.5),
             "当前时间点：未来 12-18 个月是 PE/VC 的黄金窗口期", size=16, color=WHITE, bold=True)


# ════════════════════════════════════════════════════════════════
#  SLIDE 12: 组合策略
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "组合策略：1 个亿怎么投？", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

# 投资组合卡片
allocations = [
    ("40%", "能效 SaaS + VPP 聚合", "轻资产高壁垒\n平台型布局", BRAND),
    ("30%", "电力交易 AI + 算力调度", "窗口期\n先发优势", RGBColor(0xE0, 0x60, 0x30)),
    ("20%", "浸没式液冷 + 分布式算力仓", "前瞻布局\n技术路线押注", RGBColor(0x40, 0x80, 0xC0)),
    ("10%", "极早期技术 (光子计算等)", "期权价值\n生态卡位", GRAY),
]

for i, (pct, title, desc, color) in enumerate(allocations):
    x = Inches(0.8) + Inches(3.1) * i

    # 百分比大数字
    pct_card = add_rounded_rect(slide, x, Inches(1.5), Inches(2.7), Inches(1.5), color)
    add_text_box(slide, x, Inches(1.6), Inches(2.7), Inches(1.0),
                 pct, size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.5), Inches(2.7), Inches(0.4),
                 title, size=11, color=WHITE, alignment=PP_ALIGN.CENTER)

    # 说明文字
    add_text_box(slide, x, Inches(3.3), Inches(2.7), Inches(0.8),
                 desc, size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# 核心原则
add_rounded_rect(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5), BG_LIGHT)
add_text_box(slide, Inches(1.2), Inches(4.7), Inches(10.5), Inches(0.5),
             "🎯 核心原则", size=20, color=BRAND, bold=True)

principles = [
    "一级市场赌平台型机会，设备留给二级市场",
    "轻资产优先（SaaS + VPP），重资产做配套",
    "优先布局国内空白赛道（AI 制冷优化）",
]
for i, p_text in enumerate(principles):
    add_text_box(slide, Inches(1.2), Inches(5.3) + Inches(0.4) * i, Inches(10), Inches(0.4),
                 f"▸ {p_text}", size=14, color=DARK)


# ════════════════════════════════════════════════════════════════
#  SLIDE 13: DD 五问
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "DD 必问的五个问题", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

# 五个问题卡片
questions = [
    ("1", "收入质量", "客户是试点项目还是持续付费？", "续约率多少？能否 predict revenue？"),
    ("2", "竞争壁垒", "为什么电网和数据中心不自己做？", "你的差异化到底在哪里？"),
    ("3", "数据飞轮", "用的人越多，效果越好吗？", "有没有网络效应和数据积累？"),
    ("4", "政策依赖", "如果电改放缓，商业模式还成立吗？", "有没有政策以外的真实价值？"),
    ("5", "团队背景", "有没有既懂电力又懂 AI 的复合人才？", "创始团队的行业积累有多深？"),
]

for i, (num, title, q1, q2) in enumerate(questions):
    y = Inches(1.5) + Inches(1.05) * i

    # 编号圈
    num_shape = add_rounded_rect(slide, Inches(0.8), y, Inches(0.5), Inches(0.5), BRAND)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4),
                 num, size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 标题
    add_text_box(slide, Inches(1.5), y, Inches(2), Inches(0.5),
                 title, size=16, color=BRAND, bold=True)

    # 问题
    add_text_box(slide, Inches(3.5), y, Inches(4), Inches(0.5),
                 q1, size=13, color=DARK)
    add_text_box(slide, Inches(7.5), y, Inches(4), Inches(0.5),
                 q2, size=12, color=GRAY)

    # 分割线
    add_shape(slide, Inches(1.5), y + Inches(0.55), Inches(10), Inches(0.01), LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
#  SLIDE 14: 核心判断
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "三个核心判断", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

# 三个结论卡片
conclusions = [
    ("⚡", "双 80 是硬约束不是口号", "绿电 80% + PUE 1.25\n节能用能体系从\"加分项\"变\"准入门槛\""),
    ("🔄", "广东首单 = 按钮时刻", "数据中心 VPP 从\"理论上可行\"\n到\"商业上已跑通\"，千亿级市场"),
    ("💰", "最性感的投资不是盖 DC", "液冷 ¥ 重 vs SaaS ¥¥ 轻\n一级市场赌后者，毛利率 60-80%"),
]

for i, (icon, title, desc) in enumerate(conclusions):
    x = Inches(0.8) + Inches(4.0) * i

    card = add_rounded_rect(slide, x, Inches(1.5), Inches(3.6), Inches(5.0), BG_LIGHT)

    # Icon
    add_text_box(slide, x, Inches(1.7), Inches(3.6), Inches(0.8),
                 icon, size=48, color=BRAND, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(3.0), Inches(0.6),
                 title, size=16, color=BRAND, bold=True, alignment=PP_ALIGN.CENTER)

    add_shape(slide, x + Inches(1.0), Inches(3.3), Inches(1.6), Inches(0.03), BRAND)

    # Description
    add_text_box(slide, x + Inches(0.3), Inches(3.6), Inches(3.0), Inches(1.5),
                 desc, size=13, color=DARK, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
#  SLIDE 15: 行动建议（通用版）
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6),
             "后续跟踪方向", size=28, color=BRAND, bold=True)
add_shape(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.04), BRAND)

actions = [
    ("P0", "AI 制冷优化 SaaS 标的扫描", "国内空白赛道，评估自研/投资/孵化路径"),
    ("P0", "数据中心 VPP 现货交易案例验证", "2026年5月首单落地，验证模式 D 财务模型"),
    ("P1", "绿电直供/源网荷储项目评估", "PPA vs 自建新能源，两条路径经济性对比"),
    ("P1", "存量数据中心液冷改造机会", "政策确定性强，关注改造服务商和材料供应商"),
    ("P2", "NVIDIA Rubin 全液冷供应链跟踪", "确认新技术路线下的投资机会窗口"),
    ("P2", "AI 电力交易平台早期布局", "窗口期 12-18 个月，先建立行业认知"),
]

for i, (time, action, detail) in enumerate(actions):
    y = Inches(1.5) + Inches(0.8) * i

    # 时间标签
    time_badge = add_rounded_rect(slide, Inches(0.8), y, Inches(1.0), Inches(0.5), BRAND)
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(1.0), Inches(0.4),
                 time, size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 行动
    add_text_box(slide, Inches(2.0), y, Inches(4.5), Inches(0.5),
                 action, size=14, color=DARK, bold=True)

    # 细节
    add_text_box(slide, Inches(6.5), y, Inches(5.5), Inches(0.5),
                 detail, size=12, color=GRAY)

    # 分割线
    if i < len(actions) - 1:
        add_shape(slide, Inches(2.0), y + Inches(0.6), Inches(10), Inches(0.01), LIGHT_GRAY)


# ════════════════════════════════════════════════════════════════
#  SLIDE 16: 结尾
# ════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BRAND)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
             "谢谢", size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(5.5), Inches(3.6), Inches(2.5), Inches(0.04), WHITE)

add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "陈道雷 · 汇竑资本", size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
             "以上为个人基于公开信息和行业研究的分析推演，不构成投资建议。", size=12, color=RGBColor(0xFF, 0xAA, 0xAA), alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
#  保存
# ════════════════════════════════════════════════════════════════

output_path = "/Users/Admin/OpencodeWorkspace/内容输出/智算中心节能用能/智算中心节能用能体系_完整版.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成: {output_path}")
print(f"   共 {len(prs.slides)} 页幻灯片")
