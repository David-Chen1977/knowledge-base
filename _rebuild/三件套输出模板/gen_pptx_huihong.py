#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""汇竑资本模板 PPT 生成器 v2"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

H_RED=RGBColor(0xC0,0x00,0x00); H_DARK=RGBColor(0x10,0x18,0x28)
H_BODY=RGBColor(0x2C,0x2C,0x2C); H_GOLD=RGBColor(0xB8,0x9A,0x6A)
H_MAROON=RGBColor(0x8B,0x00,0x00); H_WHITE=RGBColor(0xFF,0xFF,0xFF)
H_EMPH=RGBColor(0xC4,0x1E,0x3A); H_CARD=RGBColor(0x2C,0x3E,0x50)
H_BGRAY=RGBColor(0xF5,0xF7,0xFA); H_GRAY=RGBColor(0x36,0x41,0x53)
H_LGRAY=RGBColor(0xE8,0xEC,0xF0); H_LSLATE=RGBColor(0x8D,0x99,0xAE)
H_SLATE=RGBColor(0x34,0x3A,0x40)
FONT='微软雅黑'
OUT='/Users/Admin/三件套输出'
LOGO=os.path.join(OUT,'huihong_logo_red.png')

def nprs():
    p=Presentation(); p.slide_width=Inches(13.333); p.slide_height=Inches(7.5); return p
def bgs(p,c=H_WHITE):
    s=p.slides.add_slide(p.slide_layouts[6])
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),Inches(13.333),Inches(7.5))
    r.fill.solid();r.fill.fore_color.rgb=c;r.line.fill.background()
    sp=r._element;sp.getparent().remove(sp);s.shapes._spTree.insert(2,sp); return s
def tx(s,l,t,w,h,txt,sz=12,c=H_BODY,b=False,f=FONT,a=PP_ALIGN.LEFT):
    bx=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf=bx.text_frame;tf.word_wrap=True;p=tf.paragraphs[0];p.text=txt;p.alignment=a
    p.font.name=f;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b
def rx(s,l,t,w,h,c):
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    r.fill.solid();r.fill.fore_color.rgb=c;r.line.fill.background()
def pi(s,pth,l,t,w=None,h=None):
    if os.path.exists(pth):
        if w and h: s.shapes.add_picture(pth,Inches(l),Inches(t),Inches(w),Inches(h))
def ph(s,ttl,pn):
    sq=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0.63),Inches(0.33),Inches(0.65),Inches(0.65))
    sq.fill.solid();sq.fill.fore_color.rgb=H_WHITE
    sq.line.color.rgb=RGBColor(0x2C,0x2C,0x2C);sq.line.width=Pt(1.5)
    rx(s,1.05,0.77,0.25,0.25,H_RED)
    tx(s,1.30,0.37,9.5,0.65,ttl,26,H_RED,True)
    rx(s,1.30,0.99,11.12,0.02,H_RED)
    if os.path.exists(LOGO): pi(s,LOGO,11.67,0.12,0.75,0.75)
    tx(s,12.0,7.05,0.8,0.3,'第%d页'%pn,10,H_LSLATE,a=PP_ALIGN.RIGHT)
def tb(s,x,y,w,h,d,cw=None):
    rows=len(d);cols=len(d[0])
    ts=s.shapes.add_table(rows,cols,Inches(x),Inches(y),Inches(w),Inches(h));t=ts.table
    if cw:
        for ci,cv in enumerate(cw): t.columns[ci].width=Inches(cv)
    for r in range(rows):
        for c in range(cols):
            cl=t.cell(r,c);cl.text=str(d[r][c])
            for p in cl.text_frame.paragraphs:
                p.font.name=FONT;p.font.size=Pt(10);p.alignment=PP_ALIGN.CENTER
                p.font.color.rgb=H_WHITE if r==0 else H_BODY;p.font.bold=(r==0)
            cl.text_frame.word_wrap=True;cl.fill.solid()
            if r==0: cl.fill.fore_color.rgb=H_DARK
            elif r%2==0: cl.fill.fore_color.rgb=H_BGRAY
            else: cl.fill.fore_color.rgb=H_WHITE

p=nprs();pg=0
# S1
pg+=1;s=bgs(p);rx(s,0,0,0.06,3.5,H_RED)
pi(s,LOGO,11.5,0.3,1.2,1.2)
tx(s,1.5,1.8,10,0.8,'液冷赛道 PE 投资视角',38,H_RED,True)
tx(s,1.5,2.7,10,0.5,'千亿市场的放量逻辑与投资路径',22,H_DARK)
rx(s,1.5,3.5,3.0,0.03,H_GOLD)
tx(s,1.5,3.9,10,0.4,'GPU液冷市场 ¥690-970亿 (2026E) · CAGR 31.5%',18,H_GOLD,True)
tx(s,1.5,4.4,10,0.3,'"当单芯片功耗突破2300W，液冷不再是备选项而是必选项"',13,H_BODY)
rx(s,0,7.0,13.333,0.02,H_RED)
tx(s,1.5,7.1,10,0.3,'本材料仅供内部讨论，不构成投资建议  |  2026.06',10,H_LSLATE,a=PP_ALIGN.CENTER)

# S2
pg+=1;s=bgs(p);ph(s,'为什么液冷？芯片功耗突破风冷极限',pg)
tx(s,1.3,1.2,10,0.3,'NVIDIA GPU功耗3年增长229%，风冷天花板已到',13,H_LSLATE)
tb(s,1.3,1.7,6.0,2.5,
    [['GPU型号','年份','功耗(W)','vs 风冷极限'],['H100','2024','700','✅ 可行'],
     ['B200','2025','1,000','⚠️ 超25%'],['B300','2026','1,400','❌ 超75%'],
     ['R200','2027','2,300','❌ 超188%']],[1.5,1.0,1.5,2.0])
cl=[('核心结论',H_DARK,16,True),('',H_WHITE,6,False),
    ('• 液冷需求 = GPU出货量×功耗增速',H_BODY,13,False),
    ('• 2026起所有新建DC必须液冷',H_BODY,13,False),('',H_WHITE,4,False),
    ('B300 (1400W)',H_EMPH,14,True),('  超风冷极限75%',H_GRAY,12,False),
    ('',H_WHITE,4,False),('R200 (2300W)',H_RED,14,True),
    ('  超风冷极限188%，无替代方案',H_GRAY,12,False)]
for i,(t,c,sz,b) in enumerate(cl):
    tx(s,8.0,1.7+i*0.35,4.5,0.35,t,sz,c,b)

# S3
pg+=1;s=bgs(p);ph(s,'液冷市场规模：全球与中国共振',pg)
tx(s,1.3,1.2,10,0.3,'三重视角分别验证千亿级市场空间',13,H_LSLATE)
tb(s,1.3,1.6,4.8,2.0,
    [['市场','2025','2026E','2027E','CAGR'],
     ['全球DC液冷','$4.9B','$6.8B','$9.2B','19.5%'],
     ['全球浸没式','$1.7B','$2.3B','$3.1B','19.8%'],
     ['中国智算液冷','¥765亿','¥920亿','¥1,100亿','21%+'],
     ['中国GPU液冷','¥500亿','¥690亿','¥850亿','~30%']],[1.6,1.0,1.0,1.0,0.8])
tb(s,6.8,1.6,4.8,2.0,
    [['指标','2024','2025','2026E','2027E'],
     ['液冷DC渗透率','14%','20%','35%','50%+'],
     ['全球液冷市场($B)','3.8','4.9','6.8','9.2'],
     ['PUE达标DC占比','32%','45%','65%','85%']],[1.6,1.0,1.0,1.0,0.8])
tx(s,1.3,3.8,11,0.3,'数据来源: MarketsandMarkets, GM Insights, 券商研报',9,H_LSLATE)
for i,(n,te,ev) in enumerate([('1','2026 Q2-Q4','NVIDIA GB300/R200批量出货'),
    ('2','2026全年','PUE≤1.15硬约束生效')]):
    cx=1.3+i*3.6
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx),Inches(4.3),Inches(0.30),Inches(0.30))
    c.fill.solid();c.fill.fore_color.rgb=H_MAROON;c.line.fill.background()
    tx(s,cx,4.3,0.30,0.30,n,12,H_WHITE,True,a=PP_ALIGN.CENTER)
    tx(s,cx+0.40,4.3,1.2,0.3,te,11,H_CARD,True)
    tx(s,cx+0.40,4.7,2.5,0.5,ev,10,H_BODY)

# S4
pg+=1;s=bgs(p);ph(s,'技术路线：冷板为主，浸没为势',pg)
tx(s,1.3,1.2,10,0.3,'冷板式~70%份额，浸没式增速30%+',13,H_LSLATE)
tb(s,1.3,1.6,5.5,2.0,
    [['维度','冷板式','浸没式'],['PUE','1.10-1.15','1.04-1.10'],
     ['单机柜散热','40-100 kW','100 kW+'],['改造成本','低(可存量改造)','高(需新建)'],
     ['代表企业','英维克/浪潮/超聚变','曙光数创/高澜/阿里']],[1.5,2.0,2.0])
tx(s,7.0,1.6,5.5,0.3,'趋势判断',16,H_DARK,True)
tx(s,7.0,2.0,5.0,0.3,'• 2026-2027: 冷板仍为主流',12,H_BODY)
tx(s,7.0,2.35,5.0,0.3,'• 2027-2028: 浸没式成本打平冷板',12,H_EMPH)
tx(s,7.0,2.7,5.0,0.3,'• 2028+: 浸没式份额升至40%+',12,H_EMPH)
tx(s,7.0,3.2,5.0,0.3,'投资启示：双线布局，当前优先冷板链',13,H_RED,True)

# S5
pg+=1;s=bgs(p);ph(s,'产业链价值分布与投资热度',pg)
tx(s,1.3,1.2,10,0.3,'冷却液(耗材属性) — 最值得关注的长周期环节',13,H_LSLATE)
tb(s,1.3,1.6,5.0,3.0,
    [['环节','价值占比','技术壁垒','PE热度','推荐度'],
     ['系统集成','30%','中','★★★','★★★'],
     ['CDU','25%','中高','★★★★','★★★★'],
     ['管路/机柜','20%','低','★★','★★'],
     ['冷板','15%','高','★★★','★★★'],
     ['冷却液','10%','高','★★★★★','★★★★★']],[1.2,1.0,0.8,1.0,1.0])
tx(s,7.0,1.6,5.5,0.3,'投资逻辑',16,H_DARK,True)
tx(s,7.0,2.0,5.0,0.3,'• 冷却液具耗材属性，复购率最高',12,H_BODY)
tx(s,7.0,2.35,5.0,0.3,'• CDU占价值量最大，技术壁垒高',12,H_BODY)
tx(s,7.0,2.7,5.0,0.3,'• 系统集成适合产业资本布局',12,H_BODY)
tx(s,7.0,3.2,5.0,0.3,'聚焦：冷却液 > CDU > 冷板',13,H_RED,True)

# S6
pg+=1;s=bgs(p);ph(s,'冷却液：耗材属性驱动的长周期机会',pg)
tx(s,1.3,1.2,10,0.3,'3M退出PFAS → 催生约¥50亿国产替代空间',13,H_LSLATE)
tb(s,1.3,1.6,11.0,2.0,
    [['类型','应用路线','格局','国产替代','关注度'],
     ['电子氟化液(单相)','浸没式','3M退出/巨化入局','试产验证','★★★★★'],
     ['Novec类(两相)','浸没式','3M独家退出','实验室','★★★'],
     ['水基冷却液','冷板式','多供应商成熟','已完成','★★'],
     ['AEC密封材料','冷板辅材','进口+国产并行','部分国产','★★★']],[2.5,1.5,2.5,1.5,1.5])
tx(s,1.3,3.9,5.5,0.3,'3M退出PFAS时间线',16,H_DARK,True)
tx(s,1.3,4.3,5.5,0.3,'• 2025.12: 宣布2026年底全面停产',12,H_BODY)
tx(s,1.3,4.65,5.5,0.3,'• 2026E: 全球冷却液供应缺口显现',12,H_BODY)
tx(s,1.3,5.0,5.5,0.3,'• 2027E: 国产替代窗口打开(~¥50亿)',12,H_EMPH,True)
tx(s,7.0,3.9,5.5,0.3,'PE启示',16,H_DARK,True)
tx(s,7.0,4.3,5.5,0.3,'• 耗材属性→复购率极高→现金流稳',12,H_BODY)
tx(s,7.0,4.65,5.5,0.3,'• 具备OEM认证的厂商是最优标的',12,H_RED,True)
tx(s,7.0,5.0,5.5,0.3,'• 关注已进英维克/曙光供应链的企业',12,H_BODY)

# S7
pg+=1;s=bgs(p);ph(s,'竞争格局：国内龙头各占山头',pg)
tx(s,1.3,1.2,10,0.3,'英维克(冷板) vs 曙光数创(浸没) vs 高澜(双线)',13,H_LSLATE)
tb(s,1.3,1.6,11.0,3.5,
    [['企业','技术路线','核心客户','收入(2025E)','核心优势','热度'],
     ['英维克','冷板式','运营商/华为','~¥40亿','冷板龙头，全栈能力','★★★★'],
     ['曙光数创','浸没式','中科系/运营商','~¥25亿','浸没市占率>40%','★★★★'],
     ['高澜股份','冷板+浸没','电网/IDC','~¥15亿','双线布局，性价比','★★★'],
     ['浪潮信息','冷板(整机)','互联网/运营商','整机~¥800亿','整机厂商液冷增量','★★'],
     ['超聚变','冷板(整机)','互联网/运营商','整机~¥200亿','华为系技术实力','★★★'],
     ['维谛Vertiv','系统层(全球)','全球DC','~$80亿','全球系统集成龙头','★★★']],[2.0,1.5,2.0,2.0,2.5,1.5])
tx(s,1.3,5.5,11,0.3,'投资聚焦：核心零部件(CDU/冷板/冷却液) > 整机组装 > 系统集成',12,H_DARK,True)

# S8
pg+=1;s=bgs(p);ph(s,'PE投资策略：三阶段布局路线图',pg)
tx(s,1.3,1.2,10,0.3,'从行业深耕到产业整合的三年路线',13,H_LSLATE)
H_CARD_BG=RGBColor(0xF0,0xF4,0xF8)
for i,(n,tit,per,items) in enumerate([('1','行业深耕','0-6月',['完成产业链mapping','拜访20+目标公司','锁定3-5核心标的','评估技术壁垒+客户粘性']),('2','投资布局','6-18月',['领投/跟投冷却液+CDU','布局国产替代龙头','关注pre-IPO轮次','估值区间15-25x PE']),('3','产业整合','18-36月',['打造冷却液+模块平台','推动被投企业协同','对接下游客户资源','退出规划: IPO/并购'])]):
    cx=0.5+i*4.0
    cd=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(cx),Inches(1.8),Inches(3.6),Inches(4.0))
    cd.fill.solid();cd.fill.fore_color.rgb=H_CARD_BG;cd.line.fill.background()
    rx(s,cx,1.8,3.6,0.035,H_RED)
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx+0.2),Inches(2.1),Inches(0.40),Inches(0.40))
    c.fill.solid();c.fill.fore_color.rgb=H_MAROON;c.line.fill.background()
    tx(s,cx+0.2,2.1,0.40,0.40,n,15,H_WHITE,True,a=PP_ALIGN.CENTER)
    tx(s,cx+0.7,2.15,2.5,0.3,tit,20,H_DARK,True)
    tx(s,cx+0.7,2.5,2.5,0.25,per,11,H_LSLATE)
    for j,item in enumerate(items):
        tx(s,cx+0.2,2.9+j*0.45,3.2,0.4,'· '+item,11,H_BODY)

# S9
pg+=1;s=bgs(p);ph(s,'2026-2027 关键催化剂',pg)
tx(s,1.3,1.2,10,0.3,'密集催化剂推动液冷渗透率从20%到50%+',13,H_LSLATE)
tb(s,1.3,1.6,11.0,3.0,
    [['时间','事件','影响','确定性'],
     ['2026 Q2-Q4','NVIDIA GB300/R200批量出货','全球GPU液冷$6.9B启动','高'],
     ['2026全年','PUE≤1.15硬约束生效','新建DC必须液冷','高'],
     ['2026 H2','华为CloudMatrix超节点','国产算力液冷增量','中高'],
     ['2026-2027','3M PFAS全面停产','~¥50亿冷却液替代','高'],
     ['2027','浸没式成本打平冷板','技术迭代加速','中'],
     ['2027','液冷DC渗透率20%→50%+','¥1,300亿市场打开','中高']],[1.5,3.5,3.5,1.5])
tx(s,1.3,5.0,11,0.3,'催化剂密集度评估',16,H_DARK,True)
tx(s,1.3,5.4,11,0.3,'• 2026年是催化剂最密集的一年',12,H_BODY)
tx(s,1.3,5.75,11,0.3,'• PE应在上半年完成行业mapping，下半年进入重点谈判',12,H_RED,True)

# S10
pg+=1;s=bgs(p);ph(s,'风险因素与缓释措施',pg)
for i,(tit,desc,miti) in enumerate([('技术路线风险','冷板→浸没切换可能冲击现有冷板龙头','双线布局，关注全栈能力'),('估值过热风险','液冷赛道融资火热导致一二级倒挂','聚焦技术壁垒高的环节'),('3M退出不及预期','若延迟退出则国产替代节奏放缓','关注已获认证替代厂商'),('下游需求波动','AI capex放缓或影响液冷出货','锁定运营商+互联网长期合约')]):
    cy=1.6+i*1.15
    rx(s,1.3,cy,0.05,0.7,H_RED)
    tx(s,1.6,cy,6,0.3,tit,16,H_DARK,True)
    tx(s,1.6,cy+0.35,6,0.3,desc,12,H_GRAY)
    tx(s,8.0,cy+0.35,4,0.3,'缓释：'+miti,10,H_LSLATE)

# S11
pg+=1;s=bgs(p);ph(s,'液冷赛道投资总结',pg)
for i,(tit,sub) in enumerate([
    ('液冷是AI基础设施中确定性最高的方向','CAGR 31.5% | 催化剂密集 | 国产替代空间明确'),
    ('冷却液(耗材复购)+CDU(核心系统)产业链价值最高','复购率高 + 国产替代¥50亿市场'),
    ('PE策略：三阶段布局 — 深耕→投资→整合','当前Phase1: 完成产业链mapping，锁定核心标的'),
    ('建议关注: 英维克/曙光数创/高澜 + 冷却液新锐','')]):
    cy=1.4+i*1.3
    rx(s,1.3,cy,0.05,0.8,H_RED)
    tx(s,1.6,cy,10,0.45,tit,15,H_DARK,True)
    if sub:
        tx(s,1.6,cy+0.5,10,0.35,sub,12,H_GRAY)

path=os.path.join(OUT,'02_液冷赛道PE投资视角_PPT_汇竑模板.pptx')
p.save(path)
print('✅ 已生成: %s (%d页)' % (path,pg))
