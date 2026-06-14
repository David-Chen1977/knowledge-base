#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = '/Users/Admin/三件套输出/03_算电协同公司深度对比_PPT.pptx'
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

DARK_BG=RGBColor(0x1A,0x1A,0x2E); BLU=RGBColor(0x00,0x78,0xD4); GRN=RGBColor(0x00,0xB3,0x6F)
ORG=RGBColor(0xFF,0x8C,0x00); W=RGBColor(0xFF,0xFF,0xFF); LG=RGBColor(0xD0,0xD0,0xE0)
MG=RGBColor(0x60,0x60,0x80); DT=RGBColor(0x2D,0x2D,0x3F); TH=RGBColor(0x2D,0x2D,0x3F); TA=RGBColor(0xF0,0xF4,0xF8)

def bg(sl,c=DARK_BG):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),prs.slide_width,prs.slide_height)
    s.fill.solid();s.fill.fore_color.rgb=c;s.line.fill.background()
    sp=s._element;sp.getparent().remove(sp);sl.shapes._spTree.insert(2,sp)

def tb(sl,l,t,w,h,txt,sz=18,c=W,b=False,al=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    bx.text_frame.word_wrap=True
    p=bx.text_frame.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name='Microsoft YaHei';p.alignment=al

def bar(sl,l,t,w=0.06,h=0.5,c=BLU):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid();s.fill.fore_color.rgb=c;s.line.fill.background()

def tbl(sl,l,t,w,h,rc,cc,data,cw=None):
    ts=sl.shapes.add_table(rc,cc,Inches(l),Inches(t),Inches(w),Inches(h));tbl=ts.table
    if cw:
        for i,x in enumerate(cw):tbl.columns[i].width=Inches(x)
    for ri,rd in enumerate(data):
        for ci,ct in enumerate(rd):
            cell=tbl.cell(ri,ci);cell.text=str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(11);p.font.name='Microsoft YaHei';p.font.color.rgb=W if ri==0 else DT;p.font.bold=ri==0
            cell.fill.solid();cell.fill.fore_color.rgb=TH if ri==0 else (TA if ri%2==0 else W)

# S1: Cover
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
bar(sl,0.8,2.0,0.08,2.5,BLU)
tb(sl,1.2,2.0,10,1.2,'AIDC\u8d5b\u9053\u4e24\u5f3a\u5bf9\u6bd4',38,W,True)
tb(sl,1.2,3.3,10,0.8,'\u6da6\u6cfd\u79d1\u6280 vs \u6570\u636e\u6e2f',22,BLU)
tb(sl,1.2,4.5,10,0.5,'\u7b97\u7535\u534f\u540c\u65f6\u4ee3\u7684\u516c\u53f8\u6df1\u5ea6\u5bf9\u6bd4  |  2026.06',16,LG)
tb(sl,1.2,5.5,10,0.5,'Runze Technology vs DataPort: A Tale of Two AIDC Strategies',14,MG)

# S2: Industry Backdrop
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u884c\u4e1a\u80cc\u666f: \u7b97\u7535\u534f\u540c\u4ece\u653f\u7b56\u8d70\u5411\u843d\u5730',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tb(sl,0.8,1.3,11.5,0.5,'2026.05.15 \u56fd\u5bb6\u80fd\u6e90\u5c40\u53d1\u6587\u300a\u4ee5\u7b97\u7535\u534f\u540c\u8d4b\u80fd\u65b0\u578b\u80fd\u6e90\u4f53\u7cfb\u300b - PUE\u22641.15\u786c\u7ea6\u675f\u751f\u6548',16,DT,True)
tbl(sl,0.8,2.0,11.5,2.5,6,4,
    [['\u6307\u6807','\u5f53\u524d','2030E','\u8d8b\u52bf'],
     ['\u4e2d\u56fdDC\u7528\u7535(\u4ebfkWh)','~1,000','4,500','+350%'],
     ['\u5168\u7403DC\u7528\u7535(\u4ebfkWh)','4,150','9,450','+128%'],
     ['\u4e2d\u56fdDC\u5360\u5168\u793e\u4f1a\u7528\u7535','1.9%','3.5%','\u6301\u7eed\u4e0a\u5347'],
     ['IDC\u5e02\u573a\u89c4\u6a21(\u5168\u7403)','~$900B(2025)','~$1.6T(2030)','+78%'],
     ['PUE\u8981\u6c42','\u65e7DC: >1.5','\u65b0\u5efaDC: <1.15','\u6db2\u51b7\u521a\u9700']],
    [3.5,3.0,2.5,2.5])
tb(sl,0.8,4.8,11,0.5,'\u6838\u5fc3\u903b\u8f91: AI\u7b97\u529b\u7206\u53d1 -> DC\u7528\u7535\u6fc0\u589e -> \u7eff\u7535\u76f4\u4f9b+\u50a8\u80fd\u8c03\u5cf0+\u6db2\u51b7\u6563\u70ed\u4e09\u6761\u819d\u540c\u65f6\u63a8\u8fdb',14,BLU,True)

# S3: Runze Technology - Financials
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u6da6\u6cfd\u79d1\u6280: AIDC\u8f6c\u578b\u6807\u6746 (300442)',28,BLU,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,2.5,6,5,
    [['\u6307\u6807','2025\u5168\u5e74','YoY','2026Q1','YoY'],
     ['\u8425\u4e1a\u6536\u5165','56.74\u4ebf','+30%','18.40\u4ebf','+54%'],
     ['\u5f52\u6bcd\u51c0\u5229\u6da6','50.50\u4ebf','+182%','5.82\u4ebf','+35%'],
     ['AIDC\u6536\u5165','25.1\u4ebf(44%)','+44%','-','-'],
     ['AIDC\u6bdb\u5229\u7387','48.50%','-','-','-'],
     ['\u7814\u53d1\u6295\u5165','-','+39.75%','-','-']],
    [3.0,2.5,1.5,2.5,1.5])
tb(sl,0.8,4.0,11,0.5,'\u6838\u5fc3\u7ade\u4e89\u4f18\u52bf',18,DT,True)
for i,(p,c) in enumerate([
    ('AIDC\u5360\u6bd444%(\u5df2\u6210\u7b2c\u4e00\u5927\u4e1a\u52a1)', '\u25cf \u6bdb\u5229\u738748.50% \u8fdc\u9ad8\u4e8e\u4f20\u7edfIDC'),
    ('\u300c\u56db\u81ea\u300d\u6a21\u5f0f(\u81ea\u6295/\u81ea\u5efa/\u81ea\u6301/\u81ea\u8fd0\u7ef4)', '\u25cf \u9ad8\u58c1\u5792, \u5ba2\u6237\u4fe1\u4efb\u6210\u672c\u6781\u9ad8'),
    ('\u5168\u74039\u4e2a\u96c6\u7fa4, \u89c4\u5212~6GW, \u5df2\u8fd0\u8425~750MW', '\u25cf \u9999\u6e2f\u6c99\u5cad\u5df2\u843d\u5730(\u56fd\u9645\u5316\u7b2c\u4e00\u6b65)'),
    ('\u6db2\u51b7\u6280\u672f\u5148\u53d1(2023\u5e74\u9996\u680b\u7eaf\u6db2\u51b7\u667a\u7b97\u4e2d\u5fc3)', '\u25cf 81\u540dUPTIME AOS\u8ba4\u8bc1\u5de5\u7a0b\u5e08'),
]):
    tb(sl,1.0,4.5+i*0.6,5.5,0.5,p,12,DT,True)
    tb(sl,6.5,4.5+i*0.6,5.5,0.5,c,11,BLU)

# S4: DataPort - Financials
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u6570\u636e\u6e2f: \u8f6c\u578b\u8ffd\u8d76\u671f (603881)',28,ORG,True);bar(sl,0.8,1.0,2.5,0.04,ORG)
tbl(sl,0.8,1.3,11.5,2.5,6,5,
    [['\u6307\u6807','2025\u5168\u5e74','YoY','2026Q1','YoY'],
     ['\u8425\u4e1a\u6536\u5165','17.21\u4ebf','\u6301\u5e73','3.80\u4ebf','-4%'],
     ['\u5f52\u6bcd\u51c0\u5229\u6da6','1.39\u4ebf','+5%','0.45\u4ebf','+2%'],
     ['\u62e3\u975e\u51c0\u5229\u6da6','\u4e0b\u964d15.17%','-','-','-'],
     ['\u7ecf\u8425\u73b0\u91d1\u6d41','12.51\u4ebf','-','-','-'],
     ['\u5f53\u524dPE','~159x','-','-','-']],
    [3.0,2.5,1.5,2.5,1.5])
tb(sl,0.8,4.0,11,0.5,'\u9690\u5fe7\u4e0e\u7f3a\u53e3\u5206\u6790',18,DT,True)
for i,(p,c) in enumerate([
    ('\u8425\u6536\u589e\u957f\u505c\u6ede: 2025\u5e74\u6301\u5e73, 2026Q1\u53cd\u800c-4%', '\u26a0 \u4e0e\u6da6\u6cfd+54%\u5f62\u6210\u9c9c\u660e\u5bf9\u6bd4'),
    ('AIDC\u8f6c\u578b\u521a\u521a\u8d77\u6b65: \u667a\u7b97\u90e8\u95e82025.09\u624d\u6210\u7acb', '\u26a0 \u5c1a\u65e0\u5b9e\u8d28\u6027\u6536\u5165\u8d21\u732e'),
    ('\u4f30\u503c\u6781\u5ea6\u6602\u8d35: PE~159x', '\u26a0 \u8fdc\u8d85\u884c\u4e1a\u5e73\u5747, \u9700\u8d85\u9ad8\u589e\u957f\u8bc1\u660e'),
    ('\u56fd\u4f01\u51b3\u7b56\u673a\u5236: \u7a33\u5065\u4f46\u6162', '\u25cf \u56fd\u8d44\u80cc\u666f, \u878d\u8d44\u6210\u672c\u4f4e'),
]):
    tb(sl,1.0,4.5+i*0.6,5.5,0.5,p,12,DT,True)
    tb(sl,6.5,4.5+i*0.6,5.5,0.5,c,11,ORG)

# S5: Head-to-Head Comparison
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.7,'\u5bf9\u6bd4\u7eb5\u89c8: \u4e24\u5bb6\u516c\u53f8\u5173\u952e\u6307\u6807\u5bf9\u7167',28,W,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,4.0,10,3,
    [['','\u6da6\u6cfd\u79d1\u6280','\u6570\u636e\u6e2f'],
     ['\u8f6c\u578b\u9636\u6bb5','\u4e1a\u7ee9\u91ca\u653e\u671f','\u63a2\u7d22\u671f'],
     ['AIDC\u5360\u6bd4','44% (\u5df2\u6210\u7b2c\u4e00\u5927\u4e1a\u52a1)','~0% (\u667a\u7b97\u90e8\u95e8\u521a\u6210\u7acb)'],
     ['\u8425\u6536\u589e\u901f(2025)','+30%','\u6301\u5e73'],
     ['2026Q1\u589e\u901f','+54%','-4%'],
     ['AIDC\u6bdb\u5229\u7387','48.50%','\u5c1a\u672a\u62ab\u9732'],
     ['\u5168\u7403\u5e03\u5c40','9\u96c6\u7fa4~6GW','\u6682\u65e0'],
     ['\u6db2\u51b7\u6280\u672f','2023\u5e74\u9996\u680b\u7eaf\u6db2\u51b7','\u4f18\u5316\u4e2d'],
     ['\u4f30\u503c(PE)','~40x','~159x'],
     ['\u786e\u5b9a\u6027','**** (\u5df2\u5151\u73b0)','** (\u672a\u9a8c\u8bc1)']],
    [3.0,5.0,5.0])
tb(sl,0.8,5.5,11,0.5,'\u4e00\u53e5\u8bdd\u7ed3\u8bba: \u6da6\u6cfd\u79d1\u6280\u662f\u7b97\u7535\u534f\u540c\u8d5b\u9053\u5f53\u524d\u6700\u786e\u5b9a\u7684A\u80a1\u6807\u7684, \u6570\u636e\u6e2f\u9700\u8981\u7b49\u5f85\u667a\u7b97\u8f6c\u578b\u53d6\u5f97\u5b9e\u8d28\u6027\u8fdb\u5c55\u624d\u80fd\u786e\u8ba4\u62d0\u70b9',15,ORG,True)

# S6: Green Power Comparison
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u7eff\u8272\u7535\u529b\u5e03\u5c40\u5bf9\u6bd4',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,GRN)
tbl(sl,0.8,1.3,11.5,2.5,6,3,
    [['\u7ef4\u5ea6','\u6da6\u6cfd\u79d1\u6280','\u6570\u636e\u6e2f'],
     ['\u7eff\u7535\u91c7\u8d2d','\u5404\u56ed\u533a\u914d\u5957\u7eff\u7535\u4ea4\u6613','\u5e38\u89c4\u7eff\u7535\u91c7\u8d2d'],
     ['\u53ef\u518d\u751f\u80fd\u6e90\u76ee\u6807','2025\u5e74\u6d88\u7eb3\u7eff\u75355.3\u4ebfkWh','\u7eff\u8272DC\u8ba4\u8bc1'],
     ['PUE\u6c34\u5e73','\u6db2\u51b7\u6280\u672f\u52a9\u529b\u4f4e\u4e8e\u884c\u4e1a','\u6301\u7eed\u4f18\u5316\u4e2d'],
     ['\u50a8\u80fd\u914d\u5957','\u90e8\u5206\u56ed\u533a\u914d\u5957\u50a8\u80fd','\u672a\u89c1\u5927\u89c4\u6a21\u50a8\u80fd'],
     ['\u7b97\u7535\u534f\u540c\u53c2\u4e0e\u5ea6','\u9ad8(\u7eff\u7535+\u50a8\u80fd+\u6db2\u51b7\u4e00\u7ad9\u5f0f)','\u521d\u671f']],
    [3.0,5.0,5.0])
tb(sl,0.8,4.0,11,0.5,'\u7b97\u7535\u534f\u540c\u7684\u6838\u5fc3\u8d4b\u80fd\u4ef7\u503c',18,DT,True)
tb(sl,0.8,4.5,11,0.5,'\u2022 \u7eff\u7535\u76f4\u4f9b: \u964d\u4f4e\u7535\u529b\u6210\u672c(\u6570\u636e\u4e2d\u5fc3\u7535\u529b\u5360OPEX~60%)',13,DT)
tb(sl,0.8,5.0,11,0.5,'\u2022 \u50a8\u80fd\u8c03\u5cf0: \u964d\u4f4e\u5bb9\u91cf\u7535\u8d39, \u53c2\u4e0e\u9891\u7387\u8c03\u5e02\u573a',13,DT)
tb(sl,0.8,5.5,11,0.5,'\u2022 PUE\u4f18\u5316: \u6db2\u51b7+\u7eff\u7535\u53ef\u8fbePUE<1.1, \u7b26\u5408\u653f\u7b56\u8981\u6c42',13,DT)

# S7: PE Perspective
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.7,'PE\u89c6\u89d2: \u4ece\u884c\u4e1a\u98ce\u5411\u6807\u5230\u4e00\u7ea7\u673a\u4f1a',28,W,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tb(sl,0.8,1.3,11.5,0.5,'\u5df2\u4e0a\u5e02\u516c\u53f8->\u8d5b\u9053\u666f\u6c14\u5ea6\u98ce\u5411\u6807  |  \u771f\u6b63PE\u673a\u4f1a->AIDC\u914d\u5957\u7684\u80fd\u6e90\u670d\u52a1\u5546+\u50a8\u80fd\u8fd0\u8425\u5546',14,BLU,True)
tbl(sl,0.8,2.0,11.5,3.5,7,4,
    [['\u6295\u8d44\u5c42\u9762','\u7c7b\u578b','\u4ee3\u8868\u673a\u4f1a','PE\u9002\u5408\u5ea6'],
     ['AIDC\u56ed\u533a\u914d\u5957\u80fd\u6e90\u670d\u52a1\u5546','\u4e00\u7ea7','IDC\u9644\u8fd1\u5efa\u7eff\u7535+\u50a8\u80fd\u8d44\u4ea7','****'],
     ['\u5206\u5e03\u5f0f\u7eff\u7535+\u50a8\u80fd\u4e00\u4f53\u5316\u8fd0\u8425\u5546','\u4e00\u7ea7','\u4e3aAIDC\u63d0\u4f9bPPA+\u8c03\u5cf0\u670d\u52a1','*****'],
     ['\u7b97\u529b+\u7535\u529b\u8054\u5408\u8c03\u5ea6\u5e73\u53f0','\u65e9\u671f','\u201c\u5148\u77e5\u7535\u91cf, \u540e\u8c03\u7b97\u529b\u201d','***'],
     ['\u6db2\u51b7\u6280\u672f\u516c\u53f8(\u914d\u5957AIDC)','\u4e00\u7ea7/\u65e9\u671f','\u4ece\u6db2\u51b7\u5230\u7eff\u7535\u7684\u5168\u6808\u65b9\u6848','****'],
     ['\u50a8\u80fd\u8fd0\u8425\u5546(\u4e13\u6ce8AIDC\u573a\u666f)','\u4e00\u7ea7','AIDC\u5907\u7535+\u8c03\u5cf0\u53cc\u6536\u76ca','****'],
     ['\u77ff\u673a\u7b97\u529b\u8d4b\u80fd\u5e73\u53f0(\u6db2\u51b7+\u7eff\u7535)','\u65e9\u671f','Crusoe\u4e2d\u56fd\u7248','***']],
    [4.0,1.5,4.0,1.5])

# S8: Catalysts Timeline
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u5173\u952e\u50ac\u5316\u5242\u65f6\u95f4\u7ebf',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,2.5,6,3,
    [['\u65f6\u95f4','\u4e8b\u4ef6','\u5f71\u54cd\u6807\u7684'],
     ['2026H2','\u6da6\u6cfd\u79d1\u6280\u6d77\u5916\u56ed\u533a(\u9999\u6e2f\u6c99\u5cad)\u6295\u4ea7','\u6da6\u6cfd\u79d1\u6280'],
     ['2026H2','PUE\u22641.15\u786c\u7ea6\u675f\u5168\u9762\u6267\u884c','\u5168\u884c\u4e1a'],
     ['2026-2027','\u4e2d\u56fd\u6db2\u51b7DC\u6e17\u900f\u738720%\u219250%+','\u6da6\u6cfd(\u6db2\u51b7\u5148\u53d1)'],
     ['2026H2','\u6570\u636e\u6e2f\u667a\u7b97\u4e1a\u52a1\u80fd\u5426\u8d21\u732e\u6536\u5165?','\u6570\u636e\u6e2f(\u5173\u952e\u9a8c\u8bc1\u70b9)'],
     ['2027','\u65e5\u5747Token\u8c03\u7528\u91cf\u7a81\u7834500\u4e07\u4ebf','\u5168\u884c\u4e1aAI\u7b97\u529b']],
    [2.0,6.0,3.0])
tb(sl,0.8,4.0,11.5,0.5,'\u6ce8\u610f: \u6570\u636e\u6e2f\u7684\u7b2c\u4e00\u4e2a\u667a\u7b97\u6536\u5165\u5c06\u662f\u5173\u952e\u9a8c\u8bc1\u70b9, 2026H2\u81f32027H1\u5c06\u51b3\u5b9a\u5176\u8f6c\u578b\u53ef\u4fe1\u5ea6',14,ORG,True)

# S9: Risk
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.7,'\u98ce\u9669\u63d0\u793a',28,W,True);bar(sl,0.8,1.0,2.0,0.04,ORG)
for i,r in enumerate([
    'AI\u7b97\u529b\u589e\u901f\u4e0d\u53ca\u9884\u671f -> Token\u7ecf\u6d4e\u6ce1\u6cab\u7834\u88c2 -> AIDC\u9700\u6c42\u56de\u8c03',
    '\u7ade\u4e89\u52a0\u5267: \u4e07\u56fd\u6570\u636e/\u79e6\u6dee\u6570\u636e/\u4e16\u7eaa\u4e92\u8054\u90fd\u5728\u52a0\u7801AIDC',
    '\u8d44\u672c\u652f\u51fa\u538b\u529b: \u6da6\u6cfd\u79d1\u6280\u5168\u74036GW\u89c4\u5212\u9700\u8981\u6301\u7eed\u878d\u8d44',
    '\u7535\u529b\u653f\u7b56\u53d8\u5316: \u7eff\u7535\u4ea4\u6613/\u7535\u4ef7\u4f18\u60e0\u8c03\u6574',
    '\u6570\u636e\u6e2f\u4f30\u503c\u98ce\u9669: 159x PE\u4e0b\u8dcc\u7a7a\u95f4\u6781\u5927, \u8f6c\u578b\u4e0d\u53ca\u9884\u671f\u53ef\u80fd\u8170\u65a9',
]):
    tb(sl,1.2,1.3+i*0.8,11,0.6,'\u26a0 '+r,13,W)

# S10: Summary
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
bar(sl,0.8,2.0,0.08,2.5,BLU)
tb(sl,1.2,2.0,10,0.8,'\u603b\u7ed3: \u7b97\u7535\u534f\u540c\u8d5b\u9053\u7684\u6295\u8d44\u542b\u4e49',36,W,True)
tb(sl,1.2,3.0,10,0.5,'\u6da6\u6cfd\u79d1\u6280: AIDC\u8f6c\u578b\u5df2\u5151\u73b0, 44%\u6536\u5165\u5360\u6bd4, 48.5%\u6bdb\u5229\u7387, \u786e\u5b9a\u6027\u6700\u9ad8',16,BLU,True)
tb(sl,1.2,3.7,10,0.5,'\u6570\u636e\u6e2f: \u8f6c\u578b\u8ffd\u8d76\u671f, \u667a\u7b97\u5c1a\u65e0\u6536\u5165, 159x PE\u4f30\u503c\u542b\u8d85\u9ad8\u9884\u671f',16,ORG,True)
tb(sl,1.2,4.4,10,0.5,'PE\u89c6\u89d2: \u771f\u6b63\u4e00\u7ea7\u673a\u4f1a\u5728\u914d\u5957\u80fd\u6e90\u670d\u52a1\u5546/\u50a8\u80fd\u8fd0\u8425\u5546, \u800c\u975e\u5df2\u4e0a\u5e02\u516c\u53f8\u672c\u8eab',16,GRN,True)
tb(sl,1.2,5.5,10,0.5,'\u6570\u636e\u6765\u6e90: \u5404\u516c\u53f82025\u5e74\u62a5/2026\u5e74\u4e00\u5b63\u62a5/\u65b0\u534e\u7f51/\u5238\u5546\u7814\u62a5/\u4e1c\u65b9\u8d22\u5bcc',11,MG)

prs.save(OUTPUT)
print('PPT 03 saved:', OUTPUT)
