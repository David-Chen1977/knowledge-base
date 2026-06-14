#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = '/Users/Admin/三件套输出/02_液冷赛道PE投资视角_PPT.pptx'
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1A,0x1A,0x2E); BLU=RGBColor(0x00,0x78,0xD4); GRN=RGBColor(0x00,0xB3,0x6F)
ORG=RGBColor(0xFF,0x8C,0x00); W=RGBColor(0xFF,0xFF,0xFF); LG=RGBColor(0xD0,0xD0,0xE0)
MG=RGBColor(0x60,0x60,0x80); DT=RGBColor(0x2D,0x2D,0x3F); TH=RGBColor(0x2D,0x2D,0x3F); TA=RGBColor(0xF0,0xF4,0xF8)

def bg(sl,c=DARK_BG):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(0),Inches(0),prs.slide_width,prs.slide_height)
    s.fill.solid();s.fill.fore_color.rgb=c;s.line.fill.background()
    sp=s._element;sp.getparent().remove(sp);sl.shapes._spTree.insert(2,sp)

def tb(sl,l,t,w,h,txt,sz=18,c=W,b=False,al=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    bx.text_frame.word_wrap=True
    p=bx.text_frame.paragraphs[0];p.text=txt;p.font.size=Pt(sz);p.font.color.rgb=c;p.font.bold=b;p.font.name='Microsoft YaHei';p.alignment=al

def bar(sl,l,t,w=0.06,h=0.5,c=BLU):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid();s.fill.fore_color.rgb=c;s.line.fill.background()

def tbl(sl,l,t,w,h,rc,cc,data,cw=None):
    ts=sl.shapes.add_table(rc,cc,Inches(l),Inches(t),Inches(w),Inches(h));tbl=ts.table
    if cw:
        for i,x in enumerate(cw):tbl.columns[i].width=Inches(x)
    for ri,rd in enumerate(data):
        for ci,ct in enumerate(rd):
            cell=tbl.cell(ri,ci);cell.text=str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(11);p.font.name='Microsoft YaHei';p.font.color.rgb=W if ri==0 else DT;p.font.bold=ri==0
            cell.fill.solid();cell.fill.fore_color.rgb=TH if ri==0 else (TA if ri%2==0 else W)

# S1: Cover
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
bar(sl,0.8,2.0,0.08,2.5,BLU)
tb(sl,1.2,2.0,10,1.2,'\u6db2\u51b7\u8d5b\u9053PE\u6295\u8d44\u5730\u56fe',38,W,True)
tb(sl,1.2,3.3,10,0.8,'\u516d\u5927\u6295\u8d44\u4e3b\u9898\u4e0e\u5168\u7403\u6807\u7684\u77e9\u9635',22,BLU)
tb(sl,1.2,4.5,10,0.5,'\u5343\u4ebf\u5e02\u573a\u7684\u653e\u91cf\u8f6c\u6298\u70b9  |  2026.06',16,LG)
tb(sl,1.2,5.5,10,0.5,'Liquid Cooling: The Next Infrastructure Frontier',14,MG)

# S2: Why Now
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u4e3a\u4ec0\u4e48\u6db2\u51b7\u662f PE \u5fc5\u770b\u8d5b\u9053',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tb(sl,0.8,1.3,11.5,0.5,'\u6db2\u51b7\u4ece\u201c\u53ef\u9009\u9879\u201d\u53d8\u4e3a\u201c\u5fc5\u9009\u9879\u201d\u2014\u2014\u6280\u672f\u521b\u65b0\u9a71\u52a8\u7684\u57fa\u7840\u8bbe\u65bd\u8303\u5f0f\u8f6c\u79fb',20,DT,True)
tbl(sl,0.8,2.0,11.5,2.5,6,5,
    [['\u6307\u6807','2024','2025','2026E','2027E'],
     ['\u5168\u7403\u6db2\u51b7DC\u5e02\u573a','~$45B','$56.7B','$67.5B','$80.3B'],
     ['\u6db2\u51b7\u6e17\u900f\u7387','~14%','~26%','~40%','~50%+'],
     ['\u4e2d\u56fd\u6db2\u51b7\u670d\u52a1\u5668','\u00a5172.7\u4ebf(+67%)','-','-','-'],
     ['\u82af\u7247\u529f\u8017', '1000W', '1400W', '2300W','-'],
     ['PUE\u653f\u7b56\u7ea6\u675f','-','-','\u22641.15\u751f\u6548','-']],
    [3.5,2.0,2.0,2.0,2.0])
tb(sl,0.8,4.8,11.5,1.0,'\u4e09\u91cd\u5200\u5203: (1) \u82af\u7247\u529f\u8017\u66b4\u529b\u589e\u957f B200(1000W)->R200(2300W) (2) \u4e2d\u56fdPUE\u22641.15\u786c\u7ea6\u675f (3) 3M\u9000\u51faPFAS\u91cd\u5851\u4f9b\u5e94\u94fe',13,DT)

# S3: Market Size
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u5e02\u573a\u89c4\u6a21\u4e0e\u300c\u6db2\u51b7+\u300d\u8d5b\u9053\u62d3\u5c55',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,3.0,6,5,
    [['\u7ec6\u5206\u5e02\u573a','2025','2026E','2030/2033E','CAGR'],
     ['\u5168\u7403DC\u6db2\u51b7','$4.9B','$5.8B','$27.65B','31.5%'],
     ['\u5168\u7403\u6d78\u6ca1\u5f0f','$1.7B','$2.1B','$10.9B','19.8%'],
     ['\u5168\u7403\u6d78\u6ca1\u51b7\u5374\u6db2','$0.18B','-','$0.83B','23.9%'],
     ['\u4e2d\u56fd\u667a\u7b97\u4e2d\u5fc3\u6db2\u51b7','\u00a5765\u4ebf','\u00a5920\u4ebf','\u00a51300\u4ebf','21%+'],
     ['GPU\u6db2\u51b7(NVIDIA\u94fe)','-','\u00a5690-970\u4ebf','-','-']],
    [3.5,2.0,2.0,2.5,1.5])
tb(sl,0.8,4.5,11,0.5,'\u300c\u6db2\u51b7+\u300d\u8d5b\u9053\u62d3\u5c55',18,DT,True)
tb(sl,0.8,5.0,5.5,0.5,'\u2022 \u6db2\u51b7+\u7b97\u529b\u4e00\u4f53\u5316: Crusoe $100\u4ebf\u4f30\u503c',13,DT)
tb(sl,0.8,5.4,5.5,0.5,'\u2022 \u6db2\u51b7+\u50a8\u80fd: \u591a\u80fd\u7f51\u8377\u534f\u540c\u8c03\u5ea6',13,DT)
tb(sl,0.8,5.8,5.5,0.5,'\u2022 \u6db2\u51b7+\u7eff\u7535: \u6e90\u7f51\u8377\u50a8\u4e00\u4f53\u5316',13,DT)

# S4: Value Chain
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.7,'\u4ea7\u4e1a\u94fe\u4ef7\u503c\u5206\u5e03\u4e0ePE\u8fdb\u653b\u70b9',28,W,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,8.0,3.0,6,4,
    [['\u73af\u8282','\u4ef7\u503c\u5360\u6bd4','\u6280\u672f\u58c1\u5792','PE\u7b56\u7565'],
     ['\u82af\u7247\u7ea7\u51b7\u677f','~15%','\u9ad8(\u5fae\u901a\u9053)', '\u6295\u8d44\u7cbe\u5bc6\u5236\u9020\u56e2\u961f'],
     ['\u51b7\u5374\u6db2(\u8017\u6750)','~10%','\u9ad8(\u4ecb\u7535+\u73af\u4fdd)', '\u2b50\u590d\u8d2d\u7387\u6700\u9ad8'],
     ['CDU','~25%','\u4e2d\u9ad8', '\u6838\u5fc3\u7cfb\u7edf \u4ece\u96f6\u4ef6\u2192\u5168\u6808'],
     ['\u7ba1\u8def/\u63a5\u5934/\u673a\u6a71','~20%','\u4e2d', '\u968f\u89c4\u6a21\u653e\u91cf\u964d\u672c'],
     ['\u7cfb\u7edf\u96c6\u6210/\u8fd0\u7ef4','~30%','\u4e2d', '\u5168\u6808\u4e00\u4f53\u5316\u8d8b\u52bf']],
    [2.5,2.0,2.5,2.0])
tb(sl,0.8,4.5,11,0.7,'\u9ad8\u4ef7\u503c/\u9ad8\u58c1\u5792 = PE\u4f18\u5148\u6295\u8d44  |  \u4f4e\u4ef7\u503c/\u4f4e\u58c1\u5792 = \u5927\u5b97\u5546\u54c1\u5316\u98ce\u9669',16,W,True)

# S5: 6 Investment Themes (simplified as text blocks)
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.6,'PE/VC \u516d\u5927\u6295\u8d44\u4e3b\u9898',28,W,True);bar(sl,0.8,0.9,2.5,0.04,BLU)
themes=[
    ('1. \u6838\u5fc3\u90e8\u4ef6\u56fd\u4ea7\u5316\u66ff\u4ee3','3M\u9000\u51fa\u540e\u56fd\u4ea7\u6c1f\u5316\u6db2\u5c06\u4eab\u53d7\u7ea2\u5229, \u8fdb\u53e3\u66ff\u4ee3\u738741%',BLU),
    ('2. \u6280\u672f\u8def\u7ebf\u5dee\u5f02\u5316\u516c\u53f8','\u591a\u79cd\u8def\u7ebf\u5e76\u5b58=PE\u6700\u4f73\u6295\u5165\u65f6\u70b9, \u51b7\u677f\u534670%\u4f46\u6d78\u6ca1\u589e\u901f30%+',GRN),
    ('3. \u6db2\u51b7+\u7b97\u529b\u4e00\u4f53\u5316\u8fd0\u8425','Crusoe($100\u4ebf)\u6a21\u5f0f: \u5929\u7136\u6c14->\u53d1\u7535->\u6db2\u51b7DC->\u7b97\u529b',ORG),
    ('4. \u5b58\u91cfDC\u6db2\u51b7\u6539\u9020','\u5168\u7403~60%\u5b58\u91cfDC\u9700\u6539\u9020, Iceotope\u7684\u4e0d\u6539\u9020\u670d\u52a1\u5668\u65b9\u6848\u521b\u65b0',BLU),
    ('5. \u51b7\u5374\u6db2\u6750\u6599(\u8017\u6750\u903b\u8f91)','\u6253\u5370\u673a\u58a8\u76d2\u6a21\u5f0f: \u5934\u90e8\u5ba2\u6237\u4e00\u65e6\u91c7\u7528\u5207\u6362\u6210\u672c\u6781\u9ad8',GRN),
    ('6. \u5e76\u8d2d\u6574\u5408(\u4ea7\u4e1a\u6574\u5408\u8005)','CoolIT($4.75B)\u88abEcolab\u6536\u8d2d\u9a8c\u8bc1\u6db2\u51b7\u6280\u672f\u516c\u53f8\u5546\u4e1a\u4ef7\u503c',ORG),
]
for i,(t,d,c) in enumerate(themes):
    row=i//2;col=i%2
    x=0.8+col*6.2;y=1.3+row*1.8
    bar(sl,x,y,3.7,0.04,c)
    tb(sl,x,y+0.1,3.7,0.4,t,16,W,True)
    tb(sl,x,y+0.5,3.7,0.8,d,12,LG)

# S6: Global Targets
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u5168\u7403\u4e00\u7ea7\u5e02\u573a\u6807\u7684\u77e9\u9635',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,4.0,9,5,
    [['\u516c\u53f8','\u56fd\u5bb6','\u6280\u672f\u8def\u7ebf','\u878d\u8d44/\u4f30\u503c','PE\u9636\u6bb5'],
     ['Crusoe','US','\u6db2\u51b7+\u7b97\u529b\u4e00\u4f53\u5316','$1.4B+ / $10B','\u540e\u671f/Growth'],
     ['Corintis','CH','\u5fae\u6d41\u4f53\u82af\u7247\u7ea7','~$50M','\u6210\u957f\u671f'],
     ['Iceotope','UK','\u7cbe\u5bc6\u5168\u7ec4\u4ef6','$26M','\u6210\u957f\u671f'],
     ['Accelsius','US','\u4e24\u76f8DLC','$65M','\u6210\u957f\u671f'],
     ['ZutaCore','US/IL','\u65e0\u6c34\u4e24\u76f8DLC','\u672a\u516c\u5f00','\u6210\u957f\u671f'],
     ['Submer','ES','\u6d78\u6ca1\u5f0f','$131M','\u6210\u957f\u671f'],
     ['Colovore','US','\u6db2\u51b7DC\u8fd0\u8425\u5546','$925M\u503a','\u57fa\u7840\u8bbe\u65bd'],
     ['EcoDataCenter','SE','\u6c34\u7535+\u6db2\u51b7DC','\u20ac450M\u503a','\u57fa\u7840\u8bbe\u65bd']],
    [3.0,1.0,3.0,2.5,2.0])
tb(sl,0.8,5.5,11,0.5,'\u5173\u952e\u5224\u65ad: \u4f20\u7edf\u5de5\u4e1a\u5de8\u5934(Ecolab/JCI/Carrier)\u6b63\u5728\u79ef\u6781\u5bfb\u627e\u6db2\u51b7\u6280\u672f\u5165\u53e3',14,BLU,True)

# S7: Exit Paths
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.7,'\u9000\u51fa\u8def\u5f84\u5206\u6790',28,W,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,2.0,5,5,
    [['\u9000\u51fa\u8def\u5f84','\u9002\u5408\u6807\u7684','\u5178\u578b\u6848\u4f8b','\u9884\u671f\u56de\u62a5','\u65f6\u95f4\u5468\u671f'],
     ['\u6218\u7565\u6536\u8d2d(\u6700\u53ef\u80fd)','\u6280\u672f/\u6750\u6599\u516c\u53f8','CoolIT->Ecolab $4.75B','3-5x','3-5\u5e74'],
     ['SPAC/IPO','\u89c4\u6a21\u6536\u5165\u5e73\u53f0','Crusoe $10B','-','2-4\u5e74'],
     ['PE Secondary','\u6210\u957f\u671f\u516c\u53f8','-','2-3x','2-3\u5e74'],
     ['\u4e2d\u56fd\u4ea7\u4e1a\u8d44\u672c','\u56fd\u4ea7\u66ff\u4ee3\u516c\u53f8','-','3-5x','2-4\u5e74']],
    [3.0,2.5,3.0,2.0,1.5])
# Case study highlight
tb(sl,0.8,3.6,11,0.5,'\u6807\u6746\u6848\u4f8b: CoolIT -> Ecolab (\u00a34.75B)',20,W,True)
tb(sl,0.8,4.2,11,0.5,'2026\u5e743\u6708, \u6db2\u51b7\u8d5b\u9053\u6700\u5927\u9000\u51fa\u6848\u4f8b, \u9a8c\u8bc1\u4e86\u6280\u672f\u516c\u53f8\u7684\u5546\u4e1a\u5316\u4ef7\u503c',14,LG)
tb(sl,0.8,4.8,11,0.5,'CoolIT\u7684$4.75B\u662f\u4ef7\u683c\u951a\u70b9 - \u6db2\u51b7\u8d5b\u9053\u7684\u5e76\u8d2d\u4f30\u503c\u4e2d\u67a2\u5df2\u7ecf\u786e\u7acb',14,ORG,True)

# S8: China Opportunities
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u4e2d\u56fd\u6db2\u51b7\u751f\u6001\u4e0ePE\u673a\u4f1a',28,DT,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,3.0,6,4,
    [['\u516c\u53f8','\u65b9\u5411','\u9636\u6bb5','PE\u5173\u6ce8\u70b9'],
     ['\u6db2\u51b7\u65f6\u4ee3(\u5e7f\u4e1c)','\u6db2\u51b7\u7cfb\u7edf\u96c6\u6210','\u65e9\u671f','\u6a21\u5757\u5316\u6db2\u51b7\u65b9\u6848'],
     ['\u67ef\u8bfa\u5a01(\u65e0\u9521)','\u65b0\u80fd\u6e90\u6db2\u51b7\u6e29\u63a7','\u6210\u957f\u671f','\u50a8\u80fd+DC\u53cc\u5e94\u7528'],
     ['\u9ed1\u76fe\u80a1\u4efd(\u82cf\u5dde)','\u6e29\u63a7\u89e3\u51b3\u65b9\u6848','\u65b0\u4e09\u677f','\u8f6c\u677f\u673a\u4f1a'],
     ['\u7965\u535a\u4f20\u70ed(\u6d59\u6c5f)','\u6db2\u51b7\u677f/\u6563\u70ed','\u6210\u957f\u671f','\u51b7\u677f\u7cbe\u5bc6\u52a0\u5de5'],
     ['\u591a\u5bb6\u6c1f\u5316\u6db2\u521b\u59cb','\u7535\u5b50\u6c1f\u5316\u6db2','\u79cd\u5b50/\u65e9\u671f','3M\u9000\u51fa\u66ff\u4ee3\u5e02\u573a']],
    [3.0,3.0,2.0,3.5])
tb(sl,0.8,4.5,11,0.5,'\u5173\u952e\u5224\u65ad: \u771f\u6b63\u7684\u4e00\u7ea7\u673a\u4f1a\u5728 (1)\u6c1f\u5316\u6db2\u6750\u6599 (2)\u6709\u539f\u521b\u6280\u672f\u7684\u65e9\u671f\u56e2\u961f (3)\u6db2\u51b7+\u7b97\u529b\u8fd0\u8425\u65b0\u6a21\u5f0f',13,BLU,True)

# S9: PE Action Plan
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
tb(sl,0.8,0.3,11,0.7,'PE\u884c\u52a8\u8def\u5f84\u4e0e\u6295\u8d44\u7b56\u7565',28,W,True);bar(sl,0.8,1.0,2.5,0.04,BLU)
tbl(sl,0.8,1.3,11.5,2.0,6,5,
    [['\u6295\u8d44\u4e3b\u9898','\u5e02\u573a\u7a7a\u95f4','\u6280\u672f\u58c1\u5792','\u7ade\u4e89\u683c\u5c40','PE\u7efc\u5408\u8bc4\u5206'],
     ['\u56fd\u4ea7\u5316(\u6c1f\u5316\u6db2)','****','*****','***','4.3/5'],
     ['\u6280\u672f\u5dee\u5f02\u5316','*****','*****','***','4.5/5'],
     ['\u6db2\u51b7+\u7b97\u529b','*****','***','***','4.0/5'],
     ['\u51b7\u5374\u6db2\u6750\u6599','****','*****','****','4.5/5'],
     ['\u5e76\u8d2d\u6574\u5408','****','***','****','4.3/5']],
    [3.0,1.5,1.5,2.0,2.0])
phases=[
    ('Phase1 (0-6\u6708) \u884c\u4e1a\u6df1\u8015','\u9009\u62e91-2\u4e2a\u6280\u672f\u8def\u7ebf, \u4e0e\u6218\u7565\u4e70\u65b9CVC\u56e2\u961f\u5efa\u7acb\u5173\u7cfb, \u6df1\u5ea6\u5c3d\u8c035\u5bb6\u5168\u7403\u6807\u7684',BLU),
    ('Phase2 (6-18\u6708) \u6295\u8d44\u5e03\u5c40','\u4f18\u5148\u6295\u8d44\u51b7\u5374\u6db2\u6750\u6599(\u8017\u6750+\u9ad8\u58c1\u5792), \u6210\u957f\u671f\u6295\u6709OEM\u5408\u4f5c\u7684\u516c\u53f8, \u65e9\u671f\u5e03\u5c40\u5fae\u6d41\u4f53\u82af\u7247\u7ea7',GRN),
    ('Phase3 (18-36\u6708) \u4ea7\u4e1a\u6574\u5408','\u6574\u5408\u4e3a\u4e2d\u56fd/\u4e9a\u6d32\u8fd0\u8425\u5e73\u53f0, \u5f15\u5165\u6218\u7565\u6295\u8d44\u8005, \u51c6\u5907\u9000\u51fa(\u5de5\u4e1a\u96c6\u56e2\u6536\u8d2d/SPAC/\u6218\u7565\u8f6c\u8ba9)',ORG),
]
for i,(p,d,c) in enumerate(phases):
    y=3.6+i*1.0
    bar(sl,0.8,y,0.04,0.7,c)
    tb(sl,1.1,y,11,0.4,p,15,W,True)
    tb(sl,1.1,y+0.4,11,0.4,d,11,LG)

# S10: Risk
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl,W)
tb(sl,0.8,0.3,11,0.7,'\u98ce\u9669\u63d0\u793a',28,DT,True);bar(sl,0.8,1.0,2.0,0.04,ORG)
for i,r in enumerate([
    '\u6280\u672f\u8def\u7ebf\u6807\u51c6\u672a\u5b9a: \u5355\u76f8vs\u53cc\u76f8\u3001\u51b7\u677fvs\u6d78\u6ca1\u2014\u2014\u62bc\u9519\u8def\u7ebf\u5f52\u96f6',
    '\u4f30\u503c\u6ce1\u6cab: 2025-2026\u5e74\u8d44\u672c\u6d8c\u5165\u53ef\u80fd\u63a8\u9ad8\u4f30\u503c',
    '\u4ea7\u80fd\u8fc7\u5269: \u4e2d\u56fd\u5382\u5546\u8702\u62e5\u800c\u4e0a, \u7ade\u4e89\u622a\u538b\u7f29\u5229\u6da6\u7a7a\u95f4',
    '\u6280\u672f\u8fed\u4ee3\u592a\u5feb: \u82af\u7247\u67b6\u6784\u53d8\u5316, \u6db2\u51b7\u65b9\u6848\u9700\u4e0d\u65ad\u8fed\u4ee3',
    '\u5730\u7f18\u653f\u6cbb: \u4e2d\u7f8e\u79d1\u6280\u8131\u94a9\u5f71\u54cd\u6280\u672f\u83b7\u53d6',
]):
    tb(sl,1.2,1.3+i*0.8,11,0.6,'\u26a0 '+r,13,DT)

# S11: Summary
sl=prs.slides.add_slide(prs.slide_layouts[6]);bg(sl)
bar(sl,0.8,2.0,0.08,2.5,BLU)
tb(sl,1.2,2.0,10,0.8,'\u603b\u7ed3: \u6db2\u51b7\u8d5b\u9053\u7684PE\u6838\u5fc3\u7b56\u7565',36,W,True)
tb(sl,1.2,3.0,10,0.5,'1. \u8d4c\u6280\u672f\u8def\u7ebf - \u9009\u62e9\u6700\u6709\u5e0c\u671b\u6210\u4e3a\u6807\u51c6\u7684\u8def\u5f84',20,BLU,True)
tb(sl,1.2,3.7,10,0.5,'2. \u8d4c\u8017\u6750 - \u51b7\u5374\u6db2\u7684\u590d\u8d2d\u5c5e\u6027',20,GRN,True)
tb(sl,1.2,4.4,10,0.5,'3. \u8d4c\u6574\u5408 - \u5e76\u8d2d1+1>3\u7684\u4ea7\u4e1a\u903b\u8f91',20,ORG,True)
tb(sl,1.2,5.5,10,0.5,'\u6570\u636e\u6765\u6e90: Net Zero Insights/Crunchbase/SEC EDGAR/GM Insights/\u534e\u7ecf\u4ea7\u4e1a\u7814\u7a76\u9662/\u5238\u5546\u7814\u62a5',11,MG)

prs.save(OUTPUT)
print('PPT 02 saved:', OUTPUT)
