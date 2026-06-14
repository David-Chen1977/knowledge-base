#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUTPUT = '/Users/Admin/三件套输出/01_深挖三个投资方向_PPT.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x00, 0xB3, 0x6F)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LG = RGBColor(0xD0, 0xD0, 0xE0)
MG = RGBColor(0x60, 0x60, 0x80)
DT = RGBColor(0x2D, 0x2D, 0x3F)
TH = RGBColor(0x2D, 0x2D, 0x3F)
TA = RGBColor(0xF0, 0xF4, 0xF8)

def bg(slide, c=DARK_BG):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    sp = s._element; sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def tb(slide, l, t, w, h, txt, sz=18, c=WHITE, b=False, al=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = txt; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.font.name = 'Microsoft YaHei'; p.alignment = al
    return bx

def bar(slide, l, t, w=0.06, h=0.5, c=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def tbl(slide, l, t, w, h, rows, cols, data, cw=None):
    ts = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = ts.table
    if cw:
        for i, x in enumerate(cw):
            tbl.columns[i].width = Inches(x)
    for ri, rd in enumerate(data):
        for ci, ct in enumerate(rd):
            cell = tbl.cell(ri, ci); cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11); p.font.name = 'Microsoft YaHei'
                p.font.color.rgb = WHITE if ri == 0 else DT
                p.font.bold = ri == 0
            cell.fill.solid()
            cell.fill.fore_color.rgb = TH if ri == 0 else (TA if ri % 2 == 0 else WHITE)

# S1: Cover
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
bar(sl, 0.8, 2.0, 0.08, 2.5, BLUE)
tb(sl, 1.2, 2.0, 10, 1.2, 'AI\u57fa\u7840\u8bbe\u65bd\u6295\u8d44\u7684\u4e09\u6761\u4e3b\u7ebf', 40, WHITE, True)
tb(sl, 1.2, 3.3, 10, 0.8, '\u7b97\u7535\u534f\u540c \u00b7 \u94a0\u7535\u50a8\u80fd \u00b7 \u6db2\u51b7\u7cfb\u7edf', 24, BLUE)
tb(sl, 1.2, 4.5, 10, 0.5, 'PE/VC \u6295\u8d44\u6846\u67b6\u4e0e\u8d5b\u9053\u626b\u63cf  |  2026.06', 16, LG)
tb(sl, 1.2, 5.5, 10, 0.5, 'AI Inference at Scale -> Energy Infrastructure at Scale', 14, MG)

# S2: Core Thesis
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, WHITE)
tb(sl, 0.8, 0.3, 11, 0.7, '\u6838\u5fc3\u6295\u8d44\u903b\u8f91', 32, DT, True)
bar(sl, 0.8, 1.0, 2.0, 0.04, BLUE)
tb(sl, 0.8, 1.3, 11.5, 0.5, 'AI\u7b97\u529b\u7206\u53d1 -> \u7535\u529b\u9700\u6c42\u66b4\u589e -> \u4e09\u6761\u57fa\u7840\u8bbe\u65bd\u8d5b\u9053\u540c\u65f6\u53d7\u76ca', 20, DT, True)

thesis = [
    ('AI\u7535\u529b\u9700\u6c42', '\u4e2d\u56fdDC\u7528\u75352030E 4,500\u4ebfkWh (+350%)\n\u5168\u7403DC\u7528\u75352030E 9,450\u4ebfkWh', BLUE),
    ('\u4e09\u6761\u8d5b\u9053\u8054\u52a8', '\u7b97\u7535\u534f\u540c: \u89e3\u51b3\u7535\u4ece\u54ea\u6765\n\u94a0\u7535\u50a8\u80fd: \u89e3\u51b3\u600e\u4e48\u5b58\n\u6db2\u51b7\u7cfb\u7edf: \u89e3\u51b3\u600e\u4e48\u6563', GREEN),
    ('PE\u6295\u8d44\u7a97\u53e3', '\u4e09\u6761\u8d5b\u9053\u5747\u5904\u4e8e\u6280\u672f\u2192\u5546\u4e1a\u5316\u62d0\u70b9\u671f\n\u786e\u5b9a\u6027\u6392\u5e8f: \u6db2\u51b7\u7cfb\u7edf > \u94a0\u7535\u50a8\u80fd > \u7b97\u7535\u534f\u540c', ORANGE),
]
for i, (t, c, cl) in enumerate(thesis):
    y = 2.0 + i * 1.5
    bar(sl, 0.8, y, 0.04, 0.8, cl)
    tb(sl, 1.1, y, 3.0, 0.4, t, 18, DT, True)
    tb(sl, 1.1, y + 0.4, 10.5, 0.6, c, 13, DT)

# S3: AI Power Demand Data
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, WHITE)
tb(sl, 0.8, 0.3, 11, 0.7, 'AI\u7b97\u529b\u9a71\u52a8\u7535\u529b\u9700\u6c42\u66b4\u589e', 30, DT, True)
bar(sl, 0.8, 1.0, 2.5, 0.04, BLUE)
tbl(sl, 0.8, 1.3, 11.5, 2.8, 6, 5,
    [['\u6307\u6807', '2024', '2026E', '2030E', '\u589e\u901f'],
     ['\u4e2d\u56fdDC\u7528\u7535(\u4ebfkWh)', '1,000', '2,500', '4,500', '+350%'],
     ['\u5168\u7403DC\u7528\u7535(\u4ebfkWh)', '4,150', '6,500', '9,450', '+128%'],
     ['\u5355GPU\u529f\u8017(NVIDIA)', 'H100: 700W', 'B200: 1000W', 'R200: 2300W', '+229%'],
     ['\u5355\u673a\u6a71\u5bc6\u5ea6', '\u98ce\u51b7\u6781\u9650 20kW', '\u98ce\u51b7 20kW', '\u6db2\u51b7 50-100kW+', '-'],
     ['AI\u7528\u7535(2027E)', '-', '1,340\u4ebfkWh', '=\u8377\u5170\u5168\u5e74', '-']],
    [3.5, 2.5, 2.5, 2.5, 1.5])
tb(sl, 0.8, 4.5, 11, 0.8,
    '\u5173\u952e\u5224\u65ad: AI\u82af\u7247\u529f\u80173\u5e74\u589e\u957f3\u500d, \u6570\u636e\u4e2d\u5fc3\u7528\u75355\u5e74\u589e3.5\u500d\n\u8fd9\u662f\u80fd\u6e90\u57fa\u7840\u8bbe\u65bd\u91cd\u6784\u7684\u4ea7\u4e1a\u7ea7\u6295\u8d44\u673a\u4f1a', 14, DT)

# S4: Three Directions Overview
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, 0.8, 0.3, 11, 0.7, '\u4e09\u6761\u8d5b\u9053\u7684\u5b9a\u4f4d\u4e0e\u5173\u8054', 30, WHITE, True)
bar(sl, 0.8, 1.0, 2.5, 0.04, BLUE)

dirs = [
    ('\u65b9\u5411\u4e00: \u7b97\u7535\u534f\u540c', BLUE,
     '\u7eff\u7535\u76f4\u4f9b+\u653f\u7b56\u4fdd\u969c\n\u5e02\u573a: \u5343\u4ebf\u7ea7\n\u6838\u5fc3\u53d8\u91cf: \u7535\u6539\u901f\u5ea6',
     '\u2022 \u56fd\u5bb6\u80fd\u6e90\u5c402026.05\u53d1\u6587\n\u2022 PUE\u22641.15\u786c\u7ea6\u675f\n\u2022 \u516b\u5927\u67a2\u7ebd\u8282\u70b9\n\u2022 \u7eff\u8bc1+\u78b3\u51cf\u6392\u4e92\u8ba4'),
    ('\u65b9\u5411\u4e8c: \u94a0\u7535\u50a8\u80fd', GREEN,
     '\u91cf\u4ea7\u5143\u5e74+\u6210\u672c\u62d0\u70b9\n\u5e02\u573a: 500-1051GWh(2030E)\n\u6838\u5fc3\u53d8\u91cf: \u78b3\u9178\u94fe\u4ef7\u683c',
     '\u2022 \u5168\u7403\u6700\u5927\u94a0\u7535\u8ba2\u535560GWh\n\u2022 CATL 175Wh/kg\u9886\u8dd1\n\u2022 \u6210\u672c0.47\u21920.35\u5143/Wh\n\u2022 AIDC\u914d\u50a8\u5929\u7136\u4f18\u52bf'),
    ('\u65b9\u5411\u4e09: \u6db2\u51b7\u7cfb\u7edf', ORANGE,
     '\u82af\u7247\u6563\u70ed\u521a\u6027\u9700\u6c42\n\u5e02\u573a: \u00a5920\u4ebf(2026\u4e2d\u56fd)\n\u6838\u5fc3\u53d8\u91cf: \u6e17\u900f\u7387\u63d0\u5347',
     '\u2022 \u82af\u7247>1000W\u5fc5\u9009\u6db2\u51b7\n\u2022 \u51b7\u677f\u5f0f\u5346~70%\u4e3b\u6d41\n\u2022 \u6d78\u6ca1\u5f0f\u589e\u901f30%+\n\u2022 3M\u9000\u51fa\u91cd\u5851\u4f9b\u5e94\u683c\u5c40'),
]
for i, (t, cl, lc, rc) in enumerate(dirs):
    x = 0.8 + i * 4.1
    bar(sl, x, 1.5, 3.7, 0.06, cl)
    tb(sl, x, 1.7, 3.7, 0.4, t, 20, WHITE, True)
    tb(sl, x, 2.2, 3.7, 1.2, lc, 12, LG)
    tb(sl, x, 3.6, 3.7, 2.0, rc, 11, LG)

tb(sl, 0.8, 6.2, 11.5, 0.5,
    '\u4e09\u6761\u8d5b\u9053\u9ad8\u5ea6\u8054\u52a8 -> \u5171\u540c\u670d\u52a1AI\u7b97\u529b\u80fd\u6e90\u57fa\u7840\u8bbe\u65bd\u9700\u6c42', 14, BLUE, True)

# S5: Computing-Electricity
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, WHITE)
tb(sl, 0.8, 0.3, 11, 0.7, '\u7b97\u7535\u534f\u540c: \u7eff\u8272\u74e6\u7279\u652f\u6491\u89c4\u6a21\u6bd4\u7279', 28, BLUE, True)
bar(sl, 0.8, 1.0, 2.5, 0.04, BLUE)
tbl(sl, 0.8, 1.3, 4.5, 2.5, 5, 2,
    [['\u5173\u952e\u6307\u6807', '\u6570\u636e'],
     ['DC\u7528\u75352030E', '4,500\u4ebfkWh (+350%)'],
     ['DC\u7528\u7535\u5e74\u589e\u901f', '~25% vs \u5168\u793e\u4f1a 4.6%'],
     ['AI\u7528\u75352027E', '1,340\u4ebfkWh'],
     ['PUE\u786c\u7ea6\u675f', '\u22641.15 (2026\u751f\u6548)']],
    [2.0, 2.5])
tb(sl, 6.0, 1.3, 6.5, 0.5, '\u653f\u7b56\u6846\u67b6: \u4e09\u534f\u540c\u4f53\u7cfb', 18, DT, True)
tb(sl, 6.0, 1.9, 6.5, 0.5, '\u2022 \u65f6\u7a7a\u534f\u540c: \u4e1c\u6570\u897f\u7b97+\u6c99\u6208\u8352\u65b0\u80fd\u6e90\u5927\u57fa\u5730', 12, DT)
tb(sl, 6.0, 2.3, 6.5, 0.5, '\u2022 \u6280\u672f\u534f\u540c: \u667a\u80fd\u8c03\u5ea6+\u67d4\u6027\u8d2f\u8377+\u50a8\u80fd', 12, DT)
tb(sl, 6.0, 2.7, 6.5, 0.5, '\u2022 \u673a\u5236\u534f\u540c: \u7b97\u529b-\u7535\u529b\u4ef7\u683c\u8054\u52a8+\u7eff\u8bc1\u4e92\u8ba4', 12, DT)
tb(sl, 0.8, 4.2, 11, 0.5, '\u843d\u5730\u6a21\u5f0f', 18, DT)
tbl(sl, 0.8, 4.7, 11.5, 1.8, 5, 3,
    [['\u6a21\u5f0f', '\u4ee3\u8868\u5730\u533a', '\u6838\u5fc3\u4ef7\u503c'],
     ['\u7eff\u7535\u76f4\u8fde', '\u5185\u8499\u53e4\u4e4c\u5170\u5bdf\u5e03', '100%\u7eff\u7535\u53ef\u6eaf\u6e90'],
     ['\u7eff\u7535\u805a\u5408\u4f9b\u5e94', '\u7518\u8083/\u5f20\u5bb6\u53e3/\u5b81\u590f', '\u6e90\u7f51\u8377\u50a8\u4e00\u4f53\u5316'],
     ['\u8de8\u533a\u7535\u7b97\u8054\u5408\u8c03\u5ea6', '\u4e0a\u6d77\u7535\u4fe1\u8bd5\u70b9', '\u8c03\u5cf0\u6536\u76ca'],
     ['\u5206\u5e03\u5f0f\u7eff\u7535\u81ea\u5efa', '\u817e\u8baf\u4eea\u5f81', '\u5c4b\u9876\u5149\u4f0f+\u50a8\u80fd']],
    [3.0, 4.0, 4.0])

# S6: Na-ion Battery
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, WHITE)
tb(sl, 0.8, 0.3, 11, 0.7, '\u94a0\u7535\u50a8\u80fd: \u957f\u65f6\u50a8\u80fd\u5143\u5e74\u5df2\u81f3', 28, GREEN, True)
bar(sl, 0.8, 1.0, 2.5, 0.04, GREEN)
tb(sl, 0.8, 1.2, 11, 0.4, '\u91cc\u7a0b\u7891: 2026.04 CATL\u00d7\u6d77\u535a\u601d\u521b 3\u5e7460GWh\u94a0\u7535\u50a8\u80fd\u8ba2\u5355(\u5168\u7403\u6700\u5927)', 16, GREEN)
tbl(sl, 0.8, 1.8, 11.5, 2.3, 5, 5,
    [['\u4f01\u4e1a', '\u4ea7\u80fd\u89c4\u5212(2026)', '\u80fd\u91cf\u5bc6\u5ea6', '\u5e94\u7528\u573a\u666f', '\u6838\u5fc3\u4f18\u52bf'],
     ['\u5b81\u5fb7\u65f6\u4ee3', '25-30GWh', '175Wh/kg', '\u50a8\u80fd/\u4e58\u7528/\u6362\u7535', '\u5168\u7403\u552f\u4e00\u5168\u9762\u91cf\u4ea7'],
     ['\u6bd4\u4e9a\u8fea', '\u91cf\u4ea7\u4e2d', '160Wh/kg', '\u4e58\u7528\u8f66/\u50a8\u80fd', '\u81ea\u7814\u5c42\u6c27\u5316\u7269'],
     ['\u4e2d\u79d1\u6d77\u94a0/\u534e\u9633', '>10GWh', '-', '\u50a8\u80fd/\u91cd\u5361', '\u83b7\u5b81\u5fb7\u4e13\u5229\u6388\u6743'],
     ['\u4ebf\u7ef4\u9502\u80fd', '8-10GWh', '-', '\u50a8\u80fd', '\u9996\u5957\u5927\u5bb9\u91cf\u5e76\u7f51']],
    [2.5, 2.5, 2.0, 2.5, 2.5])
tb(sl, 0.8, 4.5, 11, 0.5, '\u6210\u672c\u66f2\u7ebf\u4e0e\u62d0\u70b9', 18, DT)
tbl(sl, 0.8, 5.0, 8.0, 1.3, 3, 4,
    [['\u6307\u6807', '\u5f53\u524d(2026H1)', '2026E\u5e95\u76ee\u6807', '2027E'],
     ['LFP\u7535\u82af\u6210\u672c', '~0.38\u5143/Wh', '~0.35\u5143/Wh', '~0.30\u5143/Wh'],
     ['\u94a0\u7535\u7535\u82af\u6210\u672c', '~0.47\u5143/Wh', '~0.35\u5143/Wh', '<0.30\u5143/Wh']],
    [2.5, 2.0, 2.0, 1.5])

# S7: Liquid Cooling
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, WHITE)
tb(sl, 0.8, 0.3, 11, 0.7, '\u6db2\u51b7\u7cfb\u7edf: \u5343\u4ebf\u5e02\u573a\u7684\u653e\u91cf\u8f6c\u6298\u70b9', 28, ORANGE, True)
bar(sl, 0.8, 1.0, 2.5, 0.04, ORANGE)
tbl(sl, 0.8, 1.3, 9.0, 1.8, 4, 4,
    [['\u7ec6\u5206\u5e02\u573a', '2025/2026E', '2030/2033E', 'CAGR'],
     ['\u5168\u7403DC\u6db2\u51b7', '$4.9B(2025)', '$17.0B(2032)', '19.5%'],
     ['\u5168\u7403\u6d78\u6ca1\u5f0f\u6db2\u51b7', '$1.7B(2025)', '$10.9B(2035)', '19.8%'],
     ['\u4e2d\u56fd\u667a\u7b97\u6db2\u51b7', '\u00a5765\u4ebf(2025)', '\u00a51,300\u4ebf(2029)', '21%+']],
    [3.5, 2.5, 2.5, 1.5])
tb(sl, 0.8, 3.4, 11, 0.5, '\u6280\u672f\u8def\u7ebf\u5bf9\u6bd4', 18, DT)
tbl(sl, 0.8, 3.8, 11.5, 1.5, 3, 5,
    [['\u6280\u672f\u8def\u7ebf', '\u4efd\u989d', 'PUE', '\u4ee3\u8868\u4f01\u4e1a', '\u6295\u8d44\u5438\u5f15\u529b'],
     ['\u51b7\u677f\u5f0f(\u4e3b\u6d41)', '~70%', '1.10-1.15', '\u82f1\u7ef4\u514b/\u6d6a\u6f6e/\u8d85\u805a\u53d8', '***'],
     ['\u6d78\u6ca1\u5f0f(\u9ad8\u589e\u957f)', '~25%', '1.04-1.10', '\u66d9\u5149\u6570\u521b/\u9ad8\u6f9c/\u963f\u91cc', '****']],
    [3.0, 1.5, 1.5, 3.5, 2.0])
tb(sl, 0.8, 5.6, 11, 0.5,
    '\u5173\u952e\u4fe1\u53f7: NVIDIA B200(1000W)->B300(1400W)->R200(2300W), \u98ce\u51b7\u7269\u7406\u6781\u9650\u5df2\u5230', 14, ORANGE, True)

# S8: PE Strategy
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
tb(sl, 0.8, 0.3, 11, 0.7, '\u4ea7\u4e1a\u94fe\u4ef7\u503c\u5206\u5e03\u4e0ePE\u8fdb\u653b\u7b56\u7565', 28, WHITE, True)
bar(sl, 0.8, 1.0, 2.5, 0.04, BLUE)
tbl(sl, 0.8, 1.3, 8.0, 2.5, 6, 4,
    [['\u73af\u8282', '\u4ef7\u503c\u5360\u6bd4', '\u6280\u672f\u58c1\u5792', 'PE\u7b56\u7565'],
     ['\u82af\u7247\u7ea7\u51b7\u677f', '~15%', '\u9ad8(\u5fae\u901a\u9053)', '\u6295\u8d44\u7cbe\u5bc6\u5236\u9020\u56e2\u961f'],
     ['CDU(\u51b7\u5374\u6db2\u5206\u914d)', '~25%', '\u4e2d\u9ad8', '\u4ece\u96f6\u4ef6\u2192\u5168\u6808'],
     ['\u51b7\u5374\u6db2(\u8017\u6750)', '~10%', '\u9ad8', '\u2b50\u8017\u6750\u5c5e\u6027 \u590d\u8d2d\u7387\u9ad8'],
     ['\u7ba1\u8def/\u63a5\u5934/\u673a\u6a71', '~20%', '\u4e2d', '\u968f\u89c4\u6a21\u653e\u91cf\u964d\u672c'],
     ['\u7cfb\u7edf\u96c6\u6210/\u8fd0\u7ef4', '~30%', '\u4e2d', '\u5168\u6808\u4e00\u4f53\u5316\u8d8b\u52bf']],
    [2.5, 2.0, 2.0, 2.0])
tb(sl, 0.8, 4.2, 11, 0.5, 'PE/VC \u6838\u5fc3\u884c\u52a8\u8def\u5f84', 20, WHITE, True)
phases = [
    ('Phase1 (0-6\u6708)', '\u884c\u4e1a\u6df1\u8015', BLUE),
    ('Phase2 (6-18\u6708)', '\u6295\u8d44\u5e03\u5c40', GREEN),
    ('Phase3 (18-36\u6708)', '\u4ea7\u4e1a\u6574\u5408', ORANGE),
]
for i, (ph, desc, cl) in enumerate(phases):
    y = 4.9 + i * 0.7
    bar(sl, 0.8, y, 0.04, 0.5, cl)
    tb(sl, 1.1, y, 11, 0.5, ph + ' ' + desc, 13, LG)

# S9: Risk
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl, WHITE)
tb(sl, 0.8, 0.3, 11, 0.7, '\u98ce\u9669\u63d0\u793a', 28, DT, True)
bar(sl, 0.8, 1.0, 2.0, 0.04, ORANGE)
for i, r in enumerate([
    'AI\u7b97\u529b\u589e\u901f\u4e0d\u53ca\u9884\u671f -> Token\u7ecf\u6d4e\u6ce1\u6cab\u7834\u88c2',
    '\u6280\u672f\u8def\u7ebf\u6807\u51c6\u672a\u5b9a -> \u62bc\u9519\u8def\u7ebf\u5f52\u96f6\u98ce\u9669',
    '\u78b3\u9178\u94fe\u4ef7\u683c\u6301\u7eed\u4f4e\u8ff7 -> \u94a0\u7535\u5e73\u4ef7\u65f6\u95f4\u70b9\u5ef6\u540e',
    '\u4ea7\u80fd\u8fc7\u5269 -> \u4ef7\u683c\u6218\u538b\u7f29\u5229\u6da6\u7a7a\u95f4',
    '\u4f30\u503c\u6ce1\u6cab -> \u8d44\u672c\u6d8c\u5165\u63a8\u9ad8\u4f30\u503c',
    '\u5730\u7f18\u653f\u6cbb -> \u4e2d\u7f8e\u79d1\u6280\u8131\u94a9',
]):
    tb(sl, 1.2, 1.3 + i * 0.8, 11, 0.6, '\u26a0 ' + r, 13, DT)

# S10: Summary
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
bar(sl, 0.8, 2.0, 0.08, 2.5, BLUE)
tb(sl, 1.2, 2.0, 11, 0.8, '\u603b\u7ed3: AI\u57fa\u7840\u8bbe\u65bd\u6295\u8d44\u7684\u9ec4\u91d1\u7a97\u53e3', 36, WHITE, True)
tb(sl, 1.2, 3.0, 11, 0.5, '\u4e09\u6761\u8d5b\u9053 -> \u4e00\u4e2a\u4e3b\u9898: AI\u65f6\u4ee3\u7684\u80fd\u6e90\u57fa\u7840\u8bbe\u65bd\u91cd\u6784', 22, BLUE, True)
tb(sl, 1.2, 4.0, 11, 0.6, '1. \u6db2\u51b7\u7cfb\u7edf: \u9700\u6c42\u6700\u521a\u6027, 2026\u653e\u91cf\u786e\u5b9a\u6027\u6700\u9ad8 - \u4f18\u5148\u5e03\u5c40', 14, LG)
tb(sl, 1.2, 4.7, 11, 0.6, '2. \u94a0\u7535\u50a8\u80fd: 2026\u91cf\u4ea7\u5143\u5e74, \u62d0\u70b9\u5df2\u81f3 - \u5bc6\u5207\u5173\u6ce8\u4ea7\u80fd\u5151\u73b0', 14, LG)
tb(sl, 1.2, 5.4, 11, 0.6, '3. \u7b97\u7535\u534f\u540c: \u653f\u7b56\u9a71\u52a8, \u957f\u671f\u8d8b\u52bf\u786e\u5b9a - \u8010\u5fc3\u7b49\u5f85\u7535\u7f51\u6539\u9769\u8282\u594f', 14, LG)
tb(sl, 1.2, 6.3, 11, 0.5, '\u6570\u636e\u6765\u6e90: \u56fd\u5bb6\u80fd\u6e90\u5c40/\u4e2d\u56fd\u7535\u529b\u62a5/ESIE2026/SEC EDGAR/\u5238\u5546\u7814\u62a5', 11, MG)

prs.save(OUTPUT)
print('PPT 01 saved:', OUTPUT)
