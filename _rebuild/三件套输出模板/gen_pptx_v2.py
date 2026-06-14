#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Enhanced PPTX generator v2 — Professional charts, tables, layout.
Generates all 3 PPTs with consulting-firm quality visuals.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import copy

# ── Palette (McKinsey / BCG inspired) ──────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x2A)   # dark backgrounds
BLUE    = RGBColor(0x14, 0x5A, 0x9E)   # primary accent
LBLUE   = RGBColor(0x34, 0x7F, 0xC3)   # lighter blue
TEAL    = RGBColor(0x00, 0x7B, 0x7F)   # secondary
GOLD    = RGBColor(0xD4, 0x8B, 0x28)   # highlight
ORANGE  = RGBColor(0xE8, 0x6C, 0x00)   # warning / emphasis
RED     = RGBColor(0xC0, 0x39, 0x2B)   # risk
GREEN   = RGBColor(0x1E, 0x8C, 0x3E)   # positive
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xE8, 0xEC, 0xF0)   # light bg
MGRAY   = RGBColor(0x8A, 0x94, 0xA6)   # muted text
DGRAY   = RGBColor(0x2C, 0x3E, 0x50)   # dark text
BLACK   = RGBColor(0x1A, 0x1A, 0x1A)

# Chart accent palette (for multi-series)
CPALETTE = [BLUE, TEAL, GOLD, ORANGE, GREEN, RGBColor(0x8E, 0x44, 0xAD)]

FONT = 'Microsoft YaHei'

# ── Shared Helpers ──────────────────────────────────────────────

def new_prs(w=13.333, h=7.5):
    prs = Presentation()
    prs.slide_width = Inches(w)
    prs.slide_height = Inches(h)
    return prs

def bg(slide, color=WHITE):
    """Full-slide background rectangle (deepest layer)."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(13.333), Inches(7.5))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    sp = s._element; sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def accent_bar(slide, x, y, w=0.06, h=0.5, color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()

def tb(slide, x, y, w, h, txt, sz=16, color=DGRAY, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = str(txt)
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return bx

def multi_tb(slide, x, y, w, h, lines, default_sz=14, default_color=DGRAY, spacing=Pt(4)):
    """Textbox with multiple styled lines. Each line = (text, sz, color, bold) or str."""
    bx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    bx.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = bx.text_frame.paragraphs[0]
        else:
            p = bx.text_frame.add_paragraph()
        if isinstance(line, str):
            p.text = line
            p.font.size = Pt(default_sz); p.font.color.rgb = default_color
        else:
            txt, sz, clr, bld = line if len(line) == 4 else (line[0], default_sz, default_color, False)
            p.text = str(txt)
            p.font.size = Pt(sz); p.font.color.rgb = clr; p.font.bold = bld
        p.font.name = FONT
        if spacing: p.space_after = spacing
    return bx

def add_table(slide, x, y, w, h, data, col_widths=None, header_color=NAVY, alt_color=LGRAY):
    """Add a styled table. data[0] = header row."""
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            if cw: tbl.columns[i].width = Inches(cw)
    for ri, row_data in enumerate(data):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(val) if val is not None else ''
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.name = FONT
                p.font.bold = (ri == 0)
                p.font.color.rgb = WHITE if ri == 0 else DGRAY
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if ri == 0 else (alt_color if ri % 2 == 0 else WHITE)
    return ts

def add_bar_chart(slide, x, y, w, h, categories, series_data, chart_title=''):
    """Add a clustered bar chart. series_data = [(name, values), ...]"""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, vals in series_data:
        chart_data.add_series(name, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    ).chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    if chart_title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = chart_title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    # Color each series
    for i, s in enumerate(chart.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = CPALETTE[i % len(CPALETTE)]
    return chart

def add_pie_chart(slide, x, y, w, h, labels, values, title=''):
    chart_data = CategoryChartData()
    chart_data.categories = labels
    chart_data.add_series('', values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_percentage = True
    data_labels.show_category_name = False
    data_labels.font.size = Pt(9)
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = CPALETTE[i % len(CPALETTE)]
    return chart

def add_line_chart(slide, x, y, w, h, categories, series_data, title=''):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, vals in series_data:
        chart_data.add_series(name, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS, Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data
    ).chart
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title
    for i, s in enumerate(chart.series):
        s.format.line.color.rgb = CPALETTE[i % len(CPALETTE)]
    return chart

def slide_header(slide, title, subtitle=None, accent=BLUE):
    """Standard slide header with accent bar."""
    tb(slide, 0.8, 0.3, 11, 0.7, title, 30, NAVY, True)
    accent_bar(slide, 0.8, 1.05, 2.5, 0.04, accent)
    if subtitle:
        tb(slide, 0.8, 1.2, 11, 0.4, subtitle, 14, MGRAY)

def section_slide(prs, title, subtitle, accent=BLUE):
    """Full-bleed section divider."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.5, 0.08, 0.04, accent)
    tb(sl, 1.2, 1.5, 10, 0.6, title, 40, WHITE, True)
    tb(sl, 1.2, 2.2, 10, 0.5, '', 12, WHITE)
    tb(sl, 1.2, 2.8, 10, 2.0, subtitle, 18, LBLUE)
    return sl


# ══════════════════════════════════════════════════════════════════
# PPT 01: 深挖三个投资方向
# ══════════════════════════════════════════════════════════════════
def build_ppt01():
    prs = new_prs()
    LAYOUT = prs.slide_layouts[6]  # blank

    # ── S1: Cover ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.8, 0.08, 0.04, GOLD)
    tb(sl, 1.2, 1.5, 10, 1.0, 'AI基础设施投资的三条主线', 44, WHITE, True)
    multi_tb(sl, 1.2, 2.6, 10, 0.4, [
        ('算电协同  ·  钠电储能  ·  液冷系统', 20, GOLD, True),
    ])
    tb(sl, 1.2, 3.3, 10, 0.4, 'PE/VC 投资框架与赛道扫描  |  2026.06', 16, MGRAY)
    tb(sl, 1.2, 4.0, 10, 0.4, '"当单芯片功耗突破2300W，AI的瓶颈从算力转向能源基础设施"', 13, LBLUE)
    # Bottom band
    accent_bar(sl, 0, 7.0, 13.333, 0.5, BLUE)

    # ── S2: AI Power Demand ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, 'AI算力暴涨正在制造一场能源危机',
                 '数据中心用电量 3 年增长 3.5 倍，芯片功耗 3 年翻 3 倍', BLUE)
    # Bar chart: DC power consumption
    add_bar_chart(sl, 0.8, 1.6, 5.5, 4.0,
        ['2024', '2026E', '2030E'],
        [('中国DC用电 (亿kWh)', [1000, 2500, 4500]),
         ('全球DC用电 (亿kWh)', [4150, 6500, 9450])],
        '数据中心用电量增长')
    # Key metrics sidebar
    multi_tb(sl, 7.0, 1.6, 5.5, 4.5, [
        ('关键数据', 18, NAVY, True),
        ('', 6, WHITE, False),
        ('芯片功耗 (NVIDIA)', 13, MGRAY, True),
        ('  H100 (2024):  700W', 13, DGRAY, False),
        ('  B200 (2025): 1,000W', 13, DGRAY, False),
        ('  R200 (2027): 2,300W  ⬆ +229%', 14, ORANGE, True),
        ('', 6, WHITE, False),
        ('AI用电 (2027E)', 13, MGRAY, True),
        ('  1,340 亿 kWh = 荷兰全年用电量', 14, DGRAY, False),
        ('', 6, WHITE, False),
        ('中国DC用电增速 ~25%/年', 13, MGRAY, True),
        ('  是全社会用电增速的 5 倍', 14, DGRAY, False),
    ])

    # ── S3: Three Tracks Overview ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    slide_header(sl, '三条赛道共同服务AI能源基础设施', accent=GOLD)
    cards = [
        ('方向一: 算电协同', BLUE,
         '解决"电从哪来"\n绿电直供 + 政策框架\n市场: 千亿级\n核心变量: 电改速度'),
        ('方向二: 钠电储能', GREEN,
         '解决"怎么存电"\n量产元年 + 成本拐点\n市场: 500-1051GWh(2030E)\n核心变量: 碳酸锂价格'),
        ('方向三: 液冷系统', ORANGE,
         '解决"怎么散热"\n芯片散热刚性需求\n市场: ¥920亿(2026中国)\n核心变量: 渗透率提升'),
    ]
    for i, (title, clr, body) in enumerate(cards):
        cx = 0.8 + i * 4.1
        # card bg
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(1.8),
                                     Inches(3.7), Inches(3.8))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x38)
        card.line.color.rgb = clr; card.line.width = Pt(2)
        # accent top bar
        accent_bar(sl, cx + 0.2, 2.0, 3.3, 0.04, clr)
        tb(sl, cx + 0.3, 2.2, 3.2, 0.5, title, 22, clr, True)
        tb(sl, cx + 0.3, 2.8, 3.2, 2.5, body, 13, LGRAY)
    tb(sl, 0.8, 6.2, 11, 0.4,
       '三条赛道高度联动 → 共同服务 AI 算力能源基础设施需求', 14, GOLD, True,
       align=PP_ALIGN.CENTER)

    # ── S4: Computing-Electricity Synergy ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '算电协同: "绿色瓦特"支撑"规模比特"',
                 '国家能源局 2026.05 发文推动 "三协同" 体系', BLUE)
    # Policy framework - 3 columns
    policies = [
        ('时空协同', '东数西算 + 沙戈荒\n新能源大基地\n跨区电力调度'),
        ('技术协同', '智能调度 + 柔性\n负载 + 储能调峰\n提升绿电消纳'),
        ('机制协同', '算力-电力价格\n联动 + 绿证互认\n碳减排认证'),
    ]
    for i, (t, d) in enumerate(policies):
        cx = 0.8 + i * 4.1
        accent_bar(sl, cx, 1.6, 3.7, 0.04, BLUE)
        tb(sl, cx, 1.8, 3.7, 0.4, t, 20, BLUE, True)
        tb(sl, cx, 2.3, 3.7, 1.5, d, 13, DGRAY)
    # Implementation models table
    tb(sl, 0.8, 3.8, 11, 0.4, '四种落地模式', 18, NAVY, True)
    add_table(sl, 0.8, 4.2, 11.5, 2.0,
        [['模式', '代表地区', '核心价值', '适合投资主体'],
         ['绿电直连', '内蒙古乌兰察布', '100%绿电可溯源', '超大规模DC运营商'],
         ['绿电聚合供应', '甘肃/张家口/宁夏', '源网荷储一体化', '园区级DC'],
         ['跨区算电联合调度', '上海电信试点', '调峰收益', '存量DC改造'],
         ['分布式绿电自建', '腾讯仪征', '屋顶光伏+储能', '大型互联网企业']],
        [2.5, 2.5, 3.5, 3.0])

    # ── S5: Na-ion Battery ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '钠电储能: 2026年最大结构性机会',
                 '标志事件: 2026.04 CATL × 海博思创 三年 60GWh 全球最大钠电协议', GREEN)
    # Key metrics - big number display
    metrics = [
        ('26.8 GWh', '2026E全球出货\n同比+200%'),
        ('500-1051 GWh', '2030E市场规模'),
        ('0.47→0.35', '成本(元/Wh)\n2026H1→年底目标'),
        ('15,000次', '循环寿命\nvs LFP ~8,000'),
        ('-40°C', '低温保持率>90%\n本质安全不燃'),
    ]
    for i, (num, desc) in enumerate(metrics):
        cx = 0.5 + i * 2.5
        accent_bar(sl, cx, 1.6, 2.2, 0.04, GREEN)
        tb(sl, cx, 1.8, 2.2, 0.5, num, 26, GREEN, True)
        tb(sl, cx, 2.3, 2.2, 0.8, desc, 11, DGRAY)
    # Cost comparison chart
    add_line_chart(sl, 0.8, 3.2, 6.0, 3.5,
        ['2024', '2025', '2026H1', '2026E', '2027E'],
        [('钠电成本 (元/Wh)', [0.58, 0.52, 0.47, 0.35, 0.28]),
         ('LFP成本 (元/Wh)', [0.45, 0.40, 0.38, 0.35, 0.30])],
        '钠电 vs LFP 成本曲线')
    # Label the crossover
    multi_tb(sl, 7.0, 3.2, 5.5, 4.0, [
        ('成本曲线分析', 18, NAVY, True),
        ('', 4, WHITE, False),
        ('• 2024-25: 钠电溢价 ~0.10 元/Wh', 13, DGRAY, False),
        ('• 2026E年底: 成本打平 LFP', 14, GREEN, True),
        ('  目标 0.35 元/Wh', 13, DGRAY, False),
        ('• 2027E: 钠电成本反超', 14, ORANGE, True),
        ('  目标 < 0.30 元/Wh', 13, DGRAY, False),
        ('', 4, WHITE, False),
        ('催化剂:', 14, NAVY, True),
        ('  头部企业产能爬坡进度', 13, DGRAY, False),
        ('  是下半年关键跟踪指标', 13, DGRAY, False),
        ('', 4, WHITE, False),
        ('AIDC 匹配度: ⭐⭐⭐⭐⭐', 14, GREEN, True),
    ])

    # ── S6: Liquid Cooling Market ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '液冷系统: 千亿市场的放量转折点',
                 'GPU功耗突破风冷极限 → 液冷成为刚性需求', ORANGE)
    # Market size bar chart
    add_bar_chart(sl, 0.8, 1.6, 5.5, 3.0,
        ['2025', '2026E', '2030E'],
        [('全球DC液冷 ($B)', [4.9, 6.8, 17.0]),
         ('中国智算液冷 (¥B)', [76.5, 92.0, 130.0])],
        '市场规模')
    # Technology comparison
    add_table(sl, 0.8, 4.8, 11.5, 2.0,
        [['技术路线', '份额', 'PUE', '单机柜散热', '代表企业', '投资吸引力'],
         ['冷板式 (主流)', '~70%', '1.10-1.15', '40-100kW', '英维克/浪潮/超聚变', '⭐⭐⭐'],
         ['浸没式 (高增长)', '~25%', '1.04-1.10', '100kW+', '曙光数创/高澜/阿里', '⭐⭐⭐⭐'],
         ['两者混合', '~5%', '1.08-1.12', '80-120kW', '早期探索', '⭐⭐']],
        [2.0, 1.2, 1.3, 1.8, 2.5, 1.5])
    # Key message
    multi_tb(sl, 7.0, 1.6, 5.5, 3.0, [
        ('为什么是现在?', 18, NAVY, True),
        ('', 4, WHITE, False),
        ('⚡ B300 (1400W)', 16, ORANGE, True),
        ('  风冷物理极限已到 (~800W)', 13, DGRAY, False),
        ('⚡ R200 (2300W)', 16, RED, True),
        ('  必选液冷，无替代方案', 13, DGRAY, False),
        ('', 4, WHITE, False),
        ('⚡ 中国 PUE ≤ 1.15 硬约束', 16, BLUE, True),
        ('  2026 年生效，新建 DC 必须液冷', 13, DGRAY, False),
        ('', 4, WHITE, False),
        ('⚡ 2026-2027 催化剂密集', 16, GOLD, True),
        ('  NVIDIA 出货 → PUE 生效 → 华为超节点', 13, DGRAY, False),
    ])

    # ── S7: Value Chain ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '液冷产业链价值分布', '系统集成 + CDU 占据价值高地', BLUE)
    # Pie chart for value distribution
    add_pie_chart(sl, 0.8, 1.4, 5.0, 4.5,
        ['系统集成/运维\n30%', 'CDU\n25%', '管路/机柜\n20%',
         '冷板\n15%', '冷却液\n10%'],
        [30, 25, 20, 15, 10],
        '产业链各环节价值占比')
    # Key insights
    multi_tb(sl, 6.2, 1.4, 6.5, 5.0, [
        ('价值分析与投资启示', 20, NAVY, True),
        ('', 6, WHITE, False),
        ('🥇 系统集成 (30%)', 16, BLUE, True),
        ('  全栈化趋势，附加值最高', 12, DGRAY, False),
        ('  关注: 具备垂直整合能力的平台', 12, DGRAY, False),
        ('', 4, WHITE, False),
        ('🥈 CDU (25%)', 16, TEAL, True),
        ('  核心系统，从部件 → 全栈', 12, DGRAY, False),
        ('  关注: 技术壁垒+客户粘性', 12, DGRAY, False),
        ('', 4, WHITE, False),
        ('🥉 冷却液 (10%)', 16, ORANGE, True),
        ('  耗材属性，高复购率 ⭐', 12, DGRAY, False),
        ('  关注: 3M退出后的国产替代', 12, DGRAY, False),
        ('', 4, WHITE, False),
        ('', 4, WHITE, False),
        ('全球格局: Vertiv/Schneider 主导系统层', 13, MGRAY, False),
        ('中国格局: 英维克/曙光数创/高澜 领先', 13, MGRAY, False),
    ])

    # ── S8: Investment Strategy ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    slide_header(sl, 'PE布局框架: "一体两翼" 策略', accent=GOLD)
    # Strategy framework
    tb(sl, 0.8, 1.8, 11, 0.4, '确定性排序:  液冷系统 > 钠电储能 > 算电协同', 20, GOLD, True)
    # Three strategy boxes
    strategies = [
        ('一体（核心仓位）', '液冷系统', BLUE,
         '高确定性 + 高增速\n\n• 2026-2027放量确定性最高\n• 全球DC液冷 CAGR 31.5%\n• 关注: 集成平台+核心部件'),
        ('左翼（高弹性）', '钠电储能', GREEN,
         '拐点已至，关注产能兑现\n\n• 2026量产元年+成本拐点\n• 60GWh订单验证需求\n• 关注: 产能爬坡+良率'),
        ('右翼（中长期）', '算电协同', ORANGE,
         '政策驱动，关注电网改革\n\n• 政策框架已完备\n• 执行力取决于电改进度\n• 关注: 跨区调度+直连模式'),
    ]
    for i, (role, name, clr, desc) in enumerate(strategies):
        cx = 0.5 + i * 4.2
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(2.5),
                                     Inches(3.8), Inches(3.5))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x38)
        card.line.color.rgb = clr; card.line.width = Pt(2)
        accent_bar(sl, cx + 0.2, 2.7, 3.4, 0.04, clr)
        tb(sl, cx + 0.3, 1.5, 3.5, 0.3, role, 12, MGRAY)
        tb(sl, cx + 0.3, 2.0, 3.5, 0.5, name, 28, clr, True)
        tb(sl, cx + 0.3, 3.0, 3.4, 2.5, desc, 13, LGRAY)

    tb(sl, 0.8, 6.5, 11, 0.4,
       '交叉机会: "绿电+储能+液冷" 一体化解决方案 — 市场空白待填补', 14, GOLD, True,
       align=PP_ALIGN.CENTER)

    # ── S9: Risk ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '风险因素', '四条核心风险与应对思路', ORANGE)
    risks = [
        ('政策落地不及预期', '算电协同依赖电网改革与地方执行', '中', ORANGE),
        ('钠电产能兑现风险', '60GWh订单交付、良率爬坡、成本曲线', '高', RED),
        ('液冷技术路径切换', '冷板→浸没式可能重塑竞争格局', '中', ORANGE),
        ('AI资本开支波动', '全球AI capex放缓或出口管制收紧', '中高', RED),
    ]
    for i, (title, desc, level, clr) in enumerate(risks):
        cy = 1.6 + i * 1.3
        # Risk indicator bar
        accent_bar(sl, 0.8, cy + 0.1, 0.06, 0.7, clr)
        # Level badge
        badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(11.0), Inches(cy + 0.1), Inches(1.3), Inches(0.4))
        badge.fill.solid(); badge.fill.fore_color.rgb = clr
        badge.line.fill.background()
        badge_tf = badge.text_frame; badge_tf.paragraphs[0].text = level
        badge_tf.paragraphs[0].font.size = Pt(11)
        badge_tf.paragraphs[0].font.color.rgb = WHITE
        badge_tf.paragraphs[0].font.bold = True; badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(sl, 1.1, cy, 9.5, 0.4, title, 18, NAVY, True)
        tb(sl, 1.1, cy + 0.4, 9.5, 0.4, desc, 13, MGRAY)

    tb(sl, 0.8, 7.0, 11, 0.3,
       '风险不是否定，而是投资框架的必要组成部分 | 数据来源: 国家能源局/ESIE2026/SEC EDGAR', 10, MGRAY)

    # ── S10: Summary ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.8, 0.08, 0.04, GOLD)
    tb(sl, 1.2, 1.8, 11, 0.8, '总结: AI基础设施投资的黄金窗口', 40, WHITE, True)
    tb(sl, 1.2, 2.6, 11, 0.4, '三条赛道 → 一个主题: AI时代的能源基础设施重构', 22, GOLD, True)
    multi_tb(sl, 1.2, 3.5, 11, 2.5, [
        ('1. 液冷系统: 需求最刚性, 2026放量确定性最高 — 优先布局', 18, WHITE, False),
        ('   全球DC液冷 CAGR 31.5% | 中国智算液冷 ¥1300亿(2029E)', 14, MGRAY, False),
        ('', 6, WHITE, False),
        ('2. 钠电储能: 2026量产元年, 拐点已至 — 密切跟踪产能兑现', 18, WHITE, False),
        ('   全球最大订单60GWh | 成本年底打平LFP | AIDC天然匹配', 14, MGRAY, False),
        ('', 6, WHITE, False),
        ('3. 算电协同: 政策驱动, 长期趋势确定 — 耐心等待电网改革节奏', 18, WHITE, False),
        ('   国家战略级 | PUE≤1.15硬约束 | 四大落地模式验证中', 14, MGRAY, False),
    ])
    accent_bar(sl, 0, 7.0, 13.333, 0.5, BLUE)

    prs.save('/Users/Admin/三件套输出/01_深挖三个投资方向_PPT_v2.pptx')
    print('✅ PPT 01 saved')


# ══════════════════════════════════════════════════════════════════
# PPT 02: 液冷赛道PE投资视角
# ══════════════════════════════════════════════════════════════════
def build_ppt02():
    prs = new_prs()
    LAYOUT = prs.slide_layouts[6]

    # ── S1: Cover ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.8, 0.08, 0.04, TEAL)
    tb(sl, 1.2, 1.5, 10, 1.0, '液冷赛道 PE 投资视角', 44, WHITE, True)
    tb(sl, 1.2, 2.6, 10, 0.4, '千亿市场的放量逻辑与投资路径  |  2026.06', 16, MGRAY)
    tb(sl, 1.2, 3.3, 10, 0.4,
       '"当单芯片功耗突破 2300W，液冷不再是备选项而是必选项"', 14, TEAL)
    # Slide count + data source
    tb(sl, 1.2, 4.2, 10, 0.3, 'GPU液冷市场 690-970亿元 (2026E)  ·  CAGR 31.5%', 16, GOLD, True)
    accent_bar(sl, 0, 7.0, 13.333, 0.5, TEAL)

    # ── S2: Why Liquid Cooling ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '为什么液冷? — 芯片功耗突破风冷物理极限',
                 'NVIDIA GPU功耗3年增长229%，风冷天花板已到', TEAL)
    # GPU power bar chart
    add_bar_chart(sl, 0.8, 1.5, 6.0, 3.5,
        ['H100\n2024', 'B200\n2025', 'B300\n2026', 'R200\n2027'],
        [('GPU功耗 (W)', [700, 1000, 1400, 2300])],
        'NVIDIA GPU 功耗演进 (W)')
    multi_tb(sl, 7.5, 1.5, 5.0, 4.0, [
        ('风冷极限 ~800W', 16, RED, True),
        ('  超过此阈值必须上液冷', 13, DGRAY, False),
        ('', 4, WHITE, False),
        ('B300 (1400W)', 16, ORANGE, True),
        ('  超过风冷极限 75%', 13, DGRAY, False),
        ('', 4, WHITE, False),
        ('R200 (2300W)', 16, RED, True),
        ('  超过风冷极限 188%', 13, DGRAY, False),
        ('  无替代方案，必选液冷', 12, RED, True),
        ('', 6, WHITE, False),
        ('⚡ 结论', 18, TEAL, True),
        ('  液冷需求 = GPU出货量 × 功耗增速', 14, DGRAY, False),
        ('  2026-2027 需求爆发确定性强', 14, DGRAY, False),
    ])

    # ── S3: Market Size ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '液冷市场规模: 全球与中国共振',
                 '三重视角分别验证千亿级市场空间', TEAL)
    # Multi-chart: Global DC liquid cooling
    add_bar_chart(sl, 0.8, 1.5, 5.5, 2.5,
        ['2025', '2026E', '2027E', '2030E'],
        [('全球DC液冷 ($B)', [4.9, 6.8, 9.2, 17.0]),
         ('全球浸没式 ($B)', [1.7, 2.3, 3.1, 5.8])],
        '全球液冷市场规模')
    add_bar_chart(sl, 7.0, 1.5, 5.5, 2.5,
        ['2025', '2026E', '2027E', '2029E'],
        [('中国智算液冷 (¥B)', [76.5, 92.0, 110.0, 130.0]),
         ('GPU液冷 (¥B)', [50.0, 69.0, 85.0, 97.0])],
        '中国液冷市场规模')
    # Growth rate table
    add_table(sl, 0.8, 4.3, 11.5, 2.0,
        [['细分市场', '2025', '2026E', '2027E', '2030E/2033E', 'CAGR'],
         ['全球DC液冷', '$4.9B', '$6.8B', '$9.2B', '$17.0B (2032)', '19.5%'],
         ['全球浸没式液冷', '$1.7B', '$2.3B', '$3.1B', '$10.9B (2035)', '19.8%'],
         ['中国智算中心液冷', '¥765亿', '¥920亿', '¥1,100亿', '¥1,300亿 (2029)', '21%+'],
         ['中国GPU液冷 (英伟达链)', '¥500亿', '¥690亿', '¥850亿', '-', '~30%']],
        [3.5, 2.0, 2.0, 2.0, 2.0, 1.5])
    tb(sl, 0.8, 6.6, 11, 0.3,
       '数据来源: MarketsandMarkets, GM Insights, PS Market, 券商研报', 10, MGRAY)

    # ── S4: Technology Routes ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '技术路线: 冷板为主, 浸没为势',
                 '冷板式 ~70% 份额，浸没式增速 30%+', TEAL)
    # Pie chart
    add_pie_chart(sl, 0.8, 1.4, 4.5, 4.0,
        ['冷板式 70%', '浸没式 25%', '其他 5%'],
        [70, 25, 5], '技术路线份额')
    # Comparison table
    add_table(sl, 5.8, 1.4, 7.0, 3.5,
        [['对比维度', '冷板式', '浸没式'],
         ['PUE', '1.10 - 1.15', '1.04 - 1.10'],
         ['单机柜散热', '40 - 100 kW', '100 kW+ 无上限'],
         ['改造成本', '低 (可存量改造)', '高 (需新建/大幅改造)'],
         ['冷却液', '水基 / 氟化液', '电子氟化液 (单相)\nNovec类 (两相)'],
         ['运维复杂度', '中', '高'],
         ['代表企业', '英维克/浪潮/超聚变', '曙光数创/高澜/阿里'],
         ['适用场景', '存量改造+新建', '新建高密度']],
        [2.0, 2.5, 2.5])
    mt = multi_tb(sl, 5.8, 5.2, 7.0, 2.0, [
        ('趋势判断', 16, TEAL, True),
        ('• 2026-2027: 冷板式仍为主流', 13, DGRAY, False),
        ('• 2027-2028: 浸没式成本打平冷板式', 14, ORANGE, False),
        ('• 2028+: 浸没式份额加速提升至 40%+', 14, ORANGE, False),
        ('投资启示: 双线布局, 但当前优先冷板链', 14, NAVY, True),
    ])

    # ── S5: Value Chain Deep Dive ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '产业链价值分布与投资热度', '冷却液 (耗材属性) ⭐ 最值得关注的长周期环节', TEAL)
    # Horizontal bar chart for value distribution
    add_bar_chart(sl, 0.8, 1.5, 7.0, 4.0,
        ['系统集成', 'CDU', '管路/机柜', '冷板', '冷却液'],
        [('价值占比 (%)', [30, 25, 20, 15, 10])],
        '产业链价值分布')
    # Investment heat map
    add_table(sl, 8.2, 1.5, 4.5, 4.5,
        [['环节', '壁垒', 'PE热度', '推荐度'],
         ['系统集成', '中', '🔥🔥🔥', '⭐⭐⭐'],
         ['CDU', '中高', '🔥🔥🔥🔥', '⭐⭐⭐⭐'],
         ['管路/机柜', '低', '🔥🔥', '⭐⭐'],
         ['冷板', '高', '🔥🔥🔥', '⭐⭐⭐'],
         ['冷却液', '高', '🔥🔥🔥🔥🔥', '⭐⭐⭐⭐⭐']],
        [1.2, 0.8, 1.2, 1.0])
    multi_tb(sl, 8.2, 6.2, 4.5, 0.6, [
        ('🔥 PE 关注度 (来自一级市场调研)', 10, MGRAY, False),
    ])

    # ── S6: Cooling Liquid Deep Dive ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '冷却液: 耗材属性驱动的长周期机会',
                 '3M 退出 PFAS 市场 → 催生 ~¥50 亿国产替代空间', TEAL)
    add_table(sl, 0.8, 1.5, 11.5, 2.5,
        [['冷却液类型', '应用路线', '当前格局', '国产替代进度', '投资关注度'],
         ['电子氟化液 (单相)', '浸没式-单相', '3M主导→退出\n国产: 巨化/新安/多氟多', '试产验证中', '⭐⭐⭐⭐⭐'],
         ['Novec类 (两相)', '浸没式-两相', '3M独家→退出\n国产: 中化蓝天/巨化', '实验室阶段', '⭐⭐⭐'],
         ['水基冷却液', '冷板式', '成熟市场, 多供应商', '已完成国产化', '⭐⭐'],
         ['AEC密封材料', '冷板式辅材', '进口+国产并行', '部分国产', '⭐⭐⭐']],
        [2.5, 1.5, 3.0, 2.0, 2.0])
    multi_tb(sl, 0.8, 4.3, 11.5, 2.5, [
        ('3M 退出 PFAS 时间线', 18, TEAL, True),
        ('', 4, WHITE, False),
        ('• 2025.12: 3M 宣布 2026 年底全面停产含 PFAS 冷却液 (Novec / Fluorinert)', 13, DGRAY, False),
        ('• 2026E: 全球冷却液供应缺口开始显现', 13, DGRAY, False),
        ('• 2027E: 国产替代窗口完全打开 (~¥50 亿 TAM)', 14, ORANGE, True),
        ('', 4, WHITE, False),
        ('⚡ PE 启示:', 16, TEAL, True),
        ('  冷却液具有耗材属性 → 复购率极高 → 现金流稳定', 13, DGRAY, False),
        ('  具备 OEM 认证的冷却液生产商是最优标的', 13, DGRAY, False),
        ('  关注: 已进入英维克/曙光数创供应链的企业', 14, NAVY, True),
    ])

    # ── S7: Competitive Landscape ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '竞争格局: 国内龙头各占山头',
                 '英维克 (冷板) vs 曙光数创 (浸没) vs 高澜 (双线)', TEAL)
    add_table(sl, 0.8, 1.4, 11.5, 4.5,
        [['企业', '技术路线', '核心客户', '收入规模 (2025E)', '竞争优势', 'PE关注度'],
         ['英维克', '冷板式', '三大运营商/华为', '~¥40亿', '冷板龙头, 全栈能力', '🔥🔥🔥🔥'],
         ['曙光数创', '浸没式', '中科系/运营商', '~¥25亿', '浸没式市占率>40%', '🔥🔥🔥🔥'],
         ['高澜股份', '冷板+浸没', '电网/IDC', '~¥15亿', '双线布局, 性价比优势', '🔥🔥🔥'],
         ['浪潮信息', '冷板式(整机)', '互联网/运营商', '整机~¥800亿', '整机厂商, 液冷为增量', '🔥🔥'],
         ['超聚变', '冷板式(整机)', '互联网/运营商', '整机~¥200亿', '华为系, 技术实力强', '🔥🔥🔥'],
         ['维谛Vertiv', '系统层(全球)', '全球DC', '~$80亿', '全球系统集成龙头', '🔥🔥🔥'],
         ['Schneider', '系统层(全球)', '全球DC', '~$350亿', '电力+制冷全栈', '🔥🔥']],
        [2.0, 1.5, 2.0, 2.0, 2.5, 1.5])
    multi_tb(sl, 0.8, 6.2, 11, 0.6, [
        ('投资聚焦: 核心零部件 (CDU/冷板/冷却液) ＞ 整机组装 ＞ 系统集成 (中国市场)', 13, NAVY, True),
    ])

    # ── S8: PE Strategy ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    slide_header(sl, 'PE 投资策略: 三阶段布局路线图', accent=TEAL)
    phases = [
        ('Phase 1\n0-6月', '行业深耕', TEAL,
         '• 完成产业链 mapping\n• 拜访 20+ 目标公司\n• 锁定 3-5 个核心标的\n• 评估技术壁垒+客户粘性'),
        ('Phase 2\n6-18月', '投资布局', GOLD,
         '• 领投/跟投冷却液+CDU\n• 布局国产替代龙头\n• 关注 pre-IPO 轮次\n• 估值区间: 15-25x PE'),
        ('Phase 3\n18-36月', '产业整合', ORANGE,
         '• 打造"冷却液+模块"平台\n• 推动被投企业战略协同\n• 对接下游客户资源\n• 退出规划: IPO/并购'),
    ]
    for i, (phase, title, clr, desc) in enumerate(phases):
        cx = 0.5 + i * 4.2
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(1.8),
                                     Inches(3.8), Inches(4.0))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x38)
        card.line.color.rgb = clr; card.line.width = Pt(2)
        accent_bar(sl, cx + 0.2, 2.0, 3.4, 0.04, clr)
        tb(sl, cx + 0.3, 1.3, 3.5, 0.6, phase, 18, clr, True)
        multi_tb(sl, cx + 0.3, 1.8, 3.5, 0.4, [(title, 16, clr, True)])
        tb(sl, cx + 0.3, 2.6, 3.4, 2.5, desc, 13, LGRAY)

    tb(sl, 0.8, 6.3, 11, 0.4,
       '核心标的: 具备OEM认证的冷却液生产商 + 冷却液+模块集成平台', 14, GOLD, True,
       align=PP_ALIGN.CENTER)
    accent_bar(sl, 0, 7.0, 13.333, 0.5, TEAL)

    # ── S9: Catalysts ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '2026-2027 关键催化剂', '密集催化剂推动液冷渗透率从 20% → 50%+', TEAL)
    events = [
        ['时间', '事件', '影响', '确定性'],
        ['2026 Q2-Q4', 'NVIDIA GB300/R200 批量出货',
         '全球 GPU 液冷 $6.9B 市场启动', '高'],
        ['2026全年', '中国 PUE ≤ 1.15 硬约束生效',
         '新建 DC 必须考虑液冷方案', '高'],
        ['2026 H2', '华为 CloudMatrix 384 超节点',
         '国产算力链液冷增量', '中高'],
        ['2026-2027', '3M PFAS 全面停产',
         '催生 ~¥50 亿冷却液国产替代', '高'],
        ['2027', '浸没式成本打平冷板式',
         '加速技术迭代, 份额提升', '中'],
        ['2027', '中国液冷 DC 渗透率 20% → 50%+',
         '¥1,300 亿市场空间打开', '中高'],
    ]
    add_table(sl, 0.8, 1.4, 11.5, 3.5,
        events, [1.5, 3.5, 3.5, 1.5])
    multi_tb(sl, 0.8, 5.2, 11, 1.5, [
        ('催化剂密集度评估', 18, TEAL, True),
        ('', 4, WHITE, False),
        ('• 2026 年是催化剂最密集的一年 — NVIDIA 出货 + PUE 生效 + 3M 退出', 14, DGRAY, False),
        ('• 每个催化剂都是液冷渗透率的加速器', 14, DGRAY, False),
        ('• PE 应在上半年完成行业 mapping, 下半年进入重点谈判', 14, NAVY, True),
    ])

    # ── S10: Risk ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '风险因素', '三条核心风险与缓释措施', ORANGE)
    risks = [
        ('技术路线风险', '冷板→浸没式切换可能冲击现有冷板龙头', '缓释: 双线布局, 关注全栈能力'),
        ('估值过热风险', '液冷赛道融资火热导致一二级估值倒挂', '缓释: 聚焦技术壁垒高的环节'),
        ('3M退出不及预期', '若 3M 延迟退出, 国产替代节奏放缓', '缓释: 关注已获认证的替代厂商'),
        ('下游需求波动', 'AI capex 放缓或影响液冷出货节奏', '缓释: 锁定运营商+互联网长期合约'),
    ]
    for i, (title, desc, miti) in enumerate(risks):
        cy = 1.5 + i * 1.3
        risk_bar = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.8), Inches(cy), Inches(0.06), Inches(0.8))
        risk_bar.fill.solid(); risk_bar.fill.fore_color.rgb = ORANGE if i < 2 else RED
        risk_bar.line.fill.background()
        tb(sl, 1.1, cy, 11, 0.35, title, 18, NAVY, True)
        tb(sl, 1.1, cy + 0.35, 11, 0.35, desc, 13, MGRAY)
        tb(sl, 11.0, cy + 0.35, 2.0, 0.35, miti, 10, TEAL)
    tb(sl, 0.8, 6.8, 11, 0.3, '数据来源: 国家能源局/M&M/GM Insights/SEC EDGAR/券商研报', 10, MGRAY)

    # ── S11: Summary ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.5, 0.08, 0.04, TEAL)
    tb(sl, 1.2, 1.2, 11, 0.8, '液冷赛道投资总结', 40, WHITE, True)
    multi_tb(sl, 1.2, 2.3, 11, 0.4, [
        ('GPU 功耗突破 + PUE 硬约束 = 液冷需求的刚性逻辑', 20, TEAL, True),
    ])
    multi_tb(sl, 1.2, 3.2, 11, 3.0, [
        ('核心结论', 22, WHITE, True),
        ('', 6, WHITE, False),
        ('1️⃣  液冷是目前 AI 基础设施赛道中确定性最高的方向', 18, WHITE, False),
        ('    CAGR 31.5% | 催化剂密度最高 | 国产替代空间明确', 14, MGRAY, False),
        ('', 6, WHITE, False),
        ('2️⃣  产业链上, 冷却液 (耗材属性) + CDU (核心系统) 价值最高', 18, WHITE, False),
        ('    冷却液: 复购率高 + 国产替代 ¥50 亿市场', 14, MGRAY, False),
        ('', 6, WHITE, False),
        ('3️⃣  PE 策略: 三阶段布局 — 深耕 → 投资 → 整合', 18, WHITE, False),
        ('    当前 Phase 1: 完成产业链 mapping, 锁定核心标的', 14, MGRAY, False),
        ('', 6, WHITE, False),
        ('4️⃣  建议关注: 英维克/曙光数创/高澜 + 冷却液新锐', 18, WHITE, False),
    ])
    accent_bar(sl, 0, 7.0, 13.333, 0.5, TEAL)

    prs.save('/Users/Admin/三件套输出/02_液冷赛道PE投资视角_PPT_v2.pptx')
    print('✅ PPT 02 saved')


# ══════════════════════════════════════════════════════════════════
# PPT 03: 算电协同公司深度对比
# ══════════════════════════════════════════════════════════════════
def build_ppt03():
    prs = new_prs()
    LAYOUT = prs.slide_layouts[6]

    # ── S1: Cover ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.8, 0.08, 0.04, GOLD)
    tb(sl, 1.2, 1.5, 10, 1.0, '算电协同公司深度对比', 44, WHITE, True)
    tb(sl, 1.2, 2.6, 10, 0.4, '绿电直连 · 源网荷储 · 跨区调度  |  2026.06', 16, MGRAY)
    tb(sl, 1.2, 3.3, 10, 0.4,
       '"从"绿色瓦特"到"规模比特" — 算力和电力正在成为一个市场"', 14, GOLD)
    tb(sl, 1.2, 4.2, 10, 0.3, '覆盖: 8家核心标的  |  3种商业模式  |  4个估值维度', 14, LBLUE)
    accent_bar(sl, 0, 7.0, 13.333, 0.5, BLUE)

    # ── S2: Policy Landscape ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '政策全景: 算电协同从概念走向落地',
                 '2025-2026 政策密集出台，框架已完备', BLUE)
    policies = [
        ('2025.12', '国家发改委\n算力基础设施\n高质量发展意见', '首次将算力与电力\n协同发展写入国家文件'),
        ('2026.03', '国家数据局\n数字中国建设\n整体规划', '算力调度+绿电\n消费比例要求'),
        ('2026.05', '国家能源局\n以算电协同赋能\n新型能源体系', '"三协同"框架\n四大落地模式'),
        ('2026 Q3\n(预期)', '算电协同\n试点示范\n管理办法', '地方试点申报\n补贴政策落地'),
    ]
    for i, (time, title, impact) in enumerate(policies):
        cx = 0.5 + i * 3.2
        # Timeline connector
        if i > 0:
            conn = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0.5 + i * 3.2 - 0.1), Inches(2.3),
                                         Inches(0.1), Inches(0.04))
            conn.fill.solid(); conn.fill.fore_color.rgb = BLUE; conn.line.fill.background()
        # Time badge
        badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(cx), Inches(1.5), Inches(1.8), Inches(0.35))
        badge.fill.solid(); badge.fill.fore_color.rgb = BLUE
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = time
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Content card
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(cx), Inches(2.0), Inches(2.8), Inches(2.5))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        card.line.color.rgb = BLUE; card.line.width = Pt(1)
        tb(sl, cx + 0.2, 2.1, 2.4, 0.8, title, 12, NAVY, True)
        tb(sl, cx + 0.2, 3.0, 2.4, 1.0, impact, 11, DGRAY)
    tb(sl, 0.8, 5.0, 11, 0.3, '关键判断: 政策框架已完备，执行力取决于电网改革进度', 15, BLUE, True)

    # ── S3: Business Models ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '三种商业模式: 各有利弊',
                 '算电协同的四种落地模式可归纳为三大类', BLUE)
    models = [
        ('绿电直连\n(资产型)', BLUE,
         '• 自建/购买绿电源\n• 长期购电协议 (PPA)\n• 100% 绿电可溯源\n• 适合: 超大规模 DC',
         '• 最高绿电确定性\n• 长期锁价降本\n• 碳减排认证溢价',
         '• 资本开支重\n• 选址受限\n• 并网周期长'),
        ('聚合供应\n(平台型)', LBLUE,
         '• 源网荷储一体化\n• 区域绿电聚合\n• 智能调度+储能\n• 适合: 园区级DC',
         '• 轻资产运营\n• 灵活性高\n• 可复制性强',
         '• 依赖电改进度\n• 调度复杂度高\n• 收益分成不确定'),
        ('跨区调度\n(服务型)', TEAL,
         '• 算电联合调度\n• 调峰收益分享\n• 绿色电力证书\n• 适合: 存量DC改造',
         '• 最低资本开支\n• 存量改造易实施\n• 可快速验证',
         '• 收益天花板低\n• 依赖电网配合\n• 竞争加剧'),
    ]
    for i, (name, clr, how, pro, con) in enumerate(models):
        cx = 0.5 + i * 4.2
        accent_bar(sl, cx, 1.6, 3.8, 0.04, clr)
        tb(sl, cx, 1.8, 3.8, 0.6, name, 20, clr, True)
        # How
        card1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(cx), Inches(2.6), Inches(3.8), Inches(1.3))
        card1.fill.solid(); card1.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        card1.line.fill.background()
        tb(sl, cx + 0.2, 2.7, 3.4, 1.0, how, 11, DGRAY)
        # Pro
        tb(sl, cx + 0.2, 4.1, 3.4, 0.2, '✅ 优势', 11, GREEN, True)
        tb(sl, cx + 0.2, 4.3, 3.4, 0.6, pro, 10, DGRAY)
        # Con
        tb(sl, cx + 0.2, 5.1, 3.4, 0.2, '⚠️ 劣势', 11, ORANGE, True)
        tb(sl, cx + 0.2, 5.3, 3.4, 0.6, con, 10, DGRAY)

    # ── S4: Company Comparison ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '核心标的对比: 8家公司全景扫描',
                 '覆盖 IDC 运营商 / 能源企业 / 科技巨头', BLUE)
    add_table(sl, 0.5, 1.4, 12.3, 4.5,
        [['公司', '商业模式', '绿电规模', '核心优势', '估值参考', 'PE关注度'],
         ['万国数据', '绿电直连+PPA', '300MW+', '全球布局+ESG', 'EV/EBITDA 18x', '🔥🔥🔥🔥'],
         ['世纪互联', '绿电直连+采购', '150MW+', '存量改造经验', 'EV/EBITDA 12x', '🔥🔥🔥'],
         ['秦淮数据', '源网荷储', '250MW+', '字节跳动生态', '私有', '🔥🔥🔥🔥🔥'],
         ['光环新网', '绿电直连', '100MW+', 'AWS合作+稳健', 'EV/EBITDA 14x', '🔥🔥🔥'],
         ['奥飞数据', '绿电直连+调度', '80MW+', '高弹性+成长期', 'EV/EBITDA 16x', '🔥🔥🔥'],
         ['数据港', '源网荷储', '120MW+', '阿里生态链', 'EV/EBITDA 15x', '🔥🔥🔥'],
         ['中国电力(国家电投)', '综合能源服务', 'GW级', '电源侧资源', '1.2x PB', '🔥🔥'],
         ['三峡能源', '绿电直供', 'GW级', '清洁能源龙头', '25x PE', '🔥🔥🔥']],
        [2.0, 2.0, 1.5, 2.5, 2.0, 1.5])
    multi_tb(sl, 0.5, 6.2, 12, 0.5, [
        ('投资启示: 秦淮数据 (字节生态+源网荷储) 和 万国数据 (全球布局+绿电) 综合评分最高', 13, NAVY, True),
    ])

    # ── S5: Financial Comparison ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '财务对比: 收入与盈利能力',
                 '选取上市/可获取数据的可比公司', BLUE)
    add_bar_chart(sl, 0.8, 1.4, 6.0, 3.5,
        ['万国数据', '世纪互联', '光环新网', '数据港', '奥飞数据'],
        [('营收 (¥亿, 2025E)', [95, 65, 80, 35, 18]),
         ('EBITDA (¥亿)', [42, 22, 28, 14, 7])],
        '营收 vs EBITDA 对比')
    # Key metrics
    add_table(sl, 7.2, 1.4, 5.5, 3.5,
        [['指标', '万国数据', '世纪互联', '光环新网'],
         ['营收 2025E', '¥95亿', '¥65亿', '¥80亿'],
         ['营收 2026E', '¥115亿', '¥75亿', '¥90亿'],
         ['EBITDA 利润率', '44%', '34%', '35%'],
         ['CAPEX/营收', '65%', '45%', '30%'],
         ['平均PUE', '1.25', '1.30', '1.28'],
         ['绿电比例', '38%', '22%', '18%']],
        [1.5, 1.3, 1.3, 1.3])
    multi_tb(sl, 0.8, 5.2, 11, 1.5, [
        ('财务趋势', 18, NAVY, True),
        ('• 头部效应明显: 万国数据 EBITDA 利润率 44% 显著领先', 13, DGRAY, False),
        ('• CAPEX 占比高: 行业都在投资扩张期, 现金流承压', 13, DGRAY, False),
        ('• 绿电比例提升正在降低用电成本, 改善利润率', 13, DGRAY, False),
        ('• 关注: 绿电占比增速 vs CAPEX 回报周期', 14, NAVY, True),
    ])

    # ── S6: Valuation Analysis ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '估值分析: 横向对比与投资锚点',
                 'EV/EBITDA 为主要估值框架', BLUE)
    add_bar_chart(sl, 0.8, 1.5, 6.0, 3.5,
        ['万国数据', '世纪互联', '光环新网', '数据港', '奥飞数据'],
        [('EV/EBITDA (当前)', [18, 12, 14, 15, 16])],
        '可比公司 EV/EBITDA (x)')
    multi_tb(sl, 7.2, 1.5, 5.5, 4.0, [
        ('估值框架', 18, NAVY, True),
        ('', 4, WHITE, False),
        ('行业均值: 15x EV/EBITDA', 16, GOLD, True),
        ('', 4, WHITE, False),
        ('溢价标的:', 14, DGRAY, True),
        ('  万国数据 18x', 13, DGRAY, False),
        ('  理由: 全球布局 + 高绿电比例', 11, MGRAY, False),
        ('', 4, WHITE, False),
        ('折价标的:', 14, DGRAY, True),
        ('  世纪互联 12x', 13, DGRAY, False),
        ('  理由: 存量改造为主, 增长受限', 11, MGRAY, False),
        ('', 4, WHITE, False),
        ('PE 锚点:', 16, TEAL, True),
        ('  合理区间: 12-18x EV/EBITDA', 13, DGRAY, False),
        ('  成长型溢价: +2-3x (绿电转型)', 13, DGRAY, False),
        ('  资产折价: -2-3x (存量为主)', 13, DGRAY, False),
    ])

    # ── S7: Investment Thesis ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    slide_header(sl, '投资逻辑: 三条投资主线', accent=GOLD)
    theses = [
        ('主线一', '绿电比例提升', TEAL,
         '• PUE≤1.15 驱动绿电采购\n• 绿电占比 → 核心竞争力\n• 受益标的: 万国/秦淮/数据港\n• 指标: 绿电占比增速'),
        ('主线二', '源网荷储平台', GOLD,
         '• 调度能力 = 核心壁垒\n• 储能配套提升绿电消纳\n• 受益标的: 秦淮/数据港\n• 关注: 储能配置比例'),
        ('主线三', '跨区调度服务', ORANGE,
         '• 算电联合调度新市场\n• 调峰收益增量\n• 受益标的: 奥飞/世纪互联\n• 关注: 试点项目进展'),
    ]
    for i, (label, name, clr, desc) in enumerate(theses):
        cx = 0.5 + i * 4.2
        card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(cx), Inches(1.8), Inches(3.8), Inches(3.8))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x15, 0x25, 0x38)
        card.line.color.rgb = clr; card.line.width = Pt(2)
        accent_bar(sl, cx + 0.2, 2.0, 3.4, 0.04, clr)
        tb(sl, cx + 0.3, 1.3, 3.5, 0.3, label, 11, MGRAY)
        tb(sl, cx + 0.3, 1.7, 3.5, 0.5, name, 24, clr, True)
        tb(sl, cx + 0.3, 2.6, 3.4, 2.5, desc, 13, LGRAY)

    # ── S8: Risk ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '风险因素', '算电协同特有的政策+执行风险', ORANGE)
    risks = [
        ('电网改革进度不及预期', '算电协同的核心前提是电力市场化改革', '高'),
        ('绿电溢价侵蚀利润', '绿电成本高于火电, 短期影响利润率', '中'),
        ('政策执行偏差', '中央政策→地方执行可能存在落差', '中'),
        ('技术调度挑战', '跨区域算电联合调度的技术难点', '低中'),
    ]
    for i, (title, desc, level) in enumerate(risks):
        cy = 1.5 + i * 1.3
        level_color = RED if level == '高' else (ORANGE if '中' in level else GOLD)
        accent_bar(sl, 0.8, cy, 0.06, 0.7, level_color)
        badge = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(11.5), Inches(cy), Inches(1.0), Inches(0.35))
        badge.fill.solid(); badge.fill.fore_color.rgb = level_color
        badge.line.fill.background()
        badge.text_frame.paragraphs[0].text = level
        badge.text_frame.paragraphs[0].font.size = Pt(10)
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb(sl, 1.1, cy, 10, 0.35, title, 18, NAVY, True)
        tb(sl, 1.1, cy + 0.35, 10, 0.35, desc, 13, MGRAY)

    # ── S9: Catalysts ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, WHITE)
    slide_header(sl, '2026-2027 催化剂时间线',
                 '政策 + 项目 + 财报 三重催化剂', BLUE)
    add_table(sl, 0.8, 1.4, 11.5, 3.0,
        [['时间', '催化剂', '影响标的', '影响程度'],
         ['2026 Q2', '国家算电协同试点名单公布', '行业整体', '⭐⭐⭐⭐⭐'],
         ['2026 Q3', '绿电交易量创新高 (预期)', '万国/秦淮', '⭐⭐⭐⭐'],
         ['2026 Q3-Q4', '多家公司绿电占比突破30%', '万国/数据港', '⭐⭐⭐'],
         ['2026 H2', '各地算电协同补贴细则出台', '行业整体', '⭐⭐⭐⭐'],
         ['2027 H1', '秦淮数据/万国数据 ESG评级上调', '秦淮/万国', '⭐⭐⭐'],
         ['2027', '跨区算电调度试点商业化', '奥飞/世纪互联', '⭐⭐⭐⭐']],
        [1.5, 4.0, 3.0, 1.5])
    multi_tb(sl, 0.8, 4.8, 11, 1.8, [
        ('催化剂节奏分析', 18, NAVY, True),
        ('', 4, WHITE, False),
        ('• 2026 Q2-Q3 是最密集窗口期: 试点+绿电+补贴三重催化', 14, DGRAY, False),
        ('• 估值重塑可能发生在 ESG 评级上调后 (2027)', 14, DGRAY, False),
        ('• PE 建议: 2026 Q2 完成建仓, 等待 Q3-Q4 催化剂释放', 14, NAVY, True),
    ])

    # ── S10: Summary ──
    sl = prs.slides.add_slide(LAYOUT); bg(sl, NAVY)
    accent_bar(sl, 0.8, 2.5, 0.08, 0.04, GOLD)
    tb(sl, 1.2, 1.2, 11, 0.8, '算电协同投资总结', 40, WHITE, True)
    multi_tb(sl, 1.2, 2.3, 11, 4.0, [
        ('核心结论', 22, WHITE, True),
        ('', 6, WHITE, False),
        ('1️⃣  政策框架已完备 → 落地阶段开启', 18, WHITE, False),
        ('    "三协同" 框架 + PUE ≤ 1.15 硬约束 = 政策确定性高', 14, MGRAY, False),
        ('', 4, WHITE, False),
        ('2️⃣  绿电比例: 下一个关键竞争维度', 18, WHITE, False),
        ('    绿电占比高的 DC 运营商将获得估值溢价', 14, MGRAY, False),
        ('', 4, WHITE, False),
        ('3️⃣  首选标的: 秦淮数据 (字节生态+源网荷储)', 18, WHITE, False),
        ('    次选: 万国数据 (全球布局+高绿电比例)', 14, MGRAY, False),
        ('', 4, WHITE, False),
        ('4️⃣  估值锚点: EV/EBITDA 12-18x', 18, WHITE, False),
        ('    合理中枢 15x, 绿电转型溢价 +2-3x', 14, MGRAY, False),
    ])
    accent_bar(sl, 0, 7.0, 13.333, 0.5, BLUE)

    prs.save('/Users/Admin/三件套输出/03_算电协同公司深度对比_PPT_v2.pptx')
    print('✅ PPT 03 saved')


# ══════════════════════════════════════════════════════════════════
if __name__ == '__main__':
    build_ppt01()
    build_ppt02()
    build_ppt03()
    print('\n=== All 3 PPTs generated successfully ===')
