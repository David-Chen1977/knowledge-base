#!/usr/bin/env python3
"""
PPT 模板渲染器 — 以指定 PPTX 为母版，只换文字生成新报告

用法:
    python3 pptx_template_render.py <内容.json> [输出路径]

内容 JSON 格式参见 示例内容.json

工作方式:
    1. 打开定稿的汇竑模板 PPTX
    2. 清空所有幻灯片（保留母版/配色/版式）
    3. 按 JSON 内容生成新幻灯片
    4. 输出新 .pptx 文件
"""

import json, sys, os, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ── 汇竑配色 ──
H_RED    = RGBColor(0xC0,0x00,0x00)
H_DARK   = RGBColor(0x10,0x18,0x28)
H_BODY   = RGBColor(0x2C,0x2C,0x2C)
H_GOLD   = RGBColor(0xB8,0x9A,0x6A)
H_MAROON = RGBColor(0x8B,0x00,0x00)
H_WHITE  = RGBColor(0xFF,0xFF,0xFF)
H_GRAY   = RGBColor(0x36,0x41,0x53)
H_LSLATE = RGBColor(0x8D,0x99,0xAE)
H_CARD   = RGBColor(0x2C,0x3E,0x50)
H_BGRAY  = RGBColor(0xF5,0xF7,0xFA)
H_EMPH   = RGBColor(0xC4,0x1E,0x3A)
H_SLATE  = RGBColor(0x34,0x3A,0x40)

FONT_BODY = '\u5fae\u8f6f\u96c5\u9ed1'
FONT_TITLE = '\u5fae\u8f6f\u96c5\u9ed1'

TEMPLATE_PATH = '/Volumes/移动硬盘/基金投融资相关/内蒙古新质动能科创投资基金路演报告申报材料20260427v2.pptx'
LOGO_PATH = '/Users/Admin/三件套输出/huihong_logo_red.png'

# ── 辅助函数 ──

def tx(s, l, t, w, h, txt, sz=12, c=H_BODY, b=False, f=FONT_BODY, a=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.alignment = a
    p.font.name = f; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b

def rx(s, l, t, w, h, c):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = c; r.line.fill.background()

def pi(s, pth, l, t, w=None, h=None):
    if os.path.exists(pth):
        if w and h: s.shapes.add_picture(pth, Inches(l), Inches(t), Inches(w), Inches(h))

def ph(s, ttl, pn):
    sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.63), Inches(0.33), Inches(0.65), Inches(0.65))
    sq.fill.solid(); sq.fill.fore_color.rgb = H_WHITE
    sq.line.color.rgb = RGBColor(0x2C,0x2C,0x2C); sq.line.width = Pt(1.5)
    rx(s, 1.05, 0.77, 0.25, 0.25, H_RED)
    tx(s, 1.30, 0.37, 9.5, 0.65, ttl, 26, H_RED, True)
    rx(s, 1.30, 0.99, 11.12, 0.02, H_RED)
    pi(s, LOGO_PATH, 11.67, 0.12, 0.75, 0.75)
    tx(s, 12.0, 7.05, 0.8, 0.3, '\u7b2c%d\u9875' % pn, 10, H_LSLATE, a=PP_ALIGN.RIGHT)

def tb(s, x, y, w, h, data, cw=None):
    rows = len(data); cols = len(data[0]) if data else 0
    if cols == 0: return
    ts = s.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    t = ts.table
    if cw:
        for ci, cv in enumerate(cw): t.columns[ci].width = Inches(cv)
    for r in range(rows):
        for c in range(cols):
            cl = t.cell(r, c); cl.text = str(data[r][c])
            for p in cl.text_frame.paragraphs:
                p.font.name = FONT_BODY; p.font.size = Pt(10); p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = H_WHITE if r == 0 else H_BODY
                p.font.bold = (r == 0)
            cl.text_frame.word_wrap = True; cl.fill.solid()
            if r == 0: cl.fill.fore_color.rgb = H_DARK
            elif r % 2 == 0: cl.fill.fore_color.rgb = H_BGRAY
            else: cl.fill.fore_color.rgb = H_WHITE


def prepare_prs(template_path):
    """从模板创建新Presentation，保留母版"""
    import shutil, tempfile, os
    tmp = os.path.join(tempfile.gettempdir(), '_pptx_template_copy.pptx')
    shutil.copy2(template_path, tmp)
    prs = Presentation(tmp)
    # Delete all slides by manipulating XML
    ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    rns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    pres_el = prs.part._element
    sldIdLst = pres_el.find('{%s}sldIdLst' % ns)
    if sldIdLst is not None:
        for sldId in list(sldIdLst):
            rid = sldId.get('{%s}id' % rns)
            if rid:
                try: prs.part.drop_rel(rid)
                except: pass
            sldIdLst.remove(sldId)
        # Add empty element to avoid XML errors
        etree.SubElement(pres_el, '{%s}sldIdLst' % ns)
    os.remove(tmp)
    return prs


def build(template_path, content, output_path):
    if not os.path.exists(template_path):
        print('⚠️ 模板不存在，使用内置设计')
        from pptx import Presentation as P2
        prs = P2()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    else:
        prs = prepare_prs(template_path)

    pg = 0
    slides_def = content.get('slides', [])

    for slide_def in slides_def:
        stype = slide_def.get('type', 'content')
        pg += 1

        try:
            layout = prs.slide_layouts[6]  # Blank layout
        except:
            from pptx import Presentation as P2
            t = P2()
            layout = t.slide_layouts[6]

        s = prs.slides.add_slide(layout)
        # White background
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid(); bg.fill.fore_color.rgb = H_WHITE; bg.line.fill.background()
        sp = bg._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)

        if stype == 'cover':
            rx(s, 0, 0, 0.06, 3.5, H_RED)
            pi(s, LOGO_PATH, 11.5, 0.3, 1.2, 1.2)
            tx(s, 1.5, 1.8, 10, 0.8, slide_def.get('title', ''), 38, H_RED, True)
            tx(s, 1.5, 2.7, 10, 0.5, slide_def.get('subtitle', ''), 22, H_DARK)
            rx(s, 1.5, 3.5, 3.0, 0.03, H_GOLD)
            if slide_def.get('data_line'):
                tx(s, 1.5, 3.9, 10, 0.4, slide_def['data_line'], 18, H_GOLD, True)
            rx(s, 0, 7.0, 13.333, 0.02, H_RED)
            tx(s, 1.5, 7.1, 10, 0.3,
               slide_def.get('footer', '\u672c\u6750\u6599\u4ec5\u4f9b\u5185\u90e8\u8ba8\u8bba\uff0c\u4e0d\u6784\u6210\u6295\u8d44\u5efa\u8bae'),
               10, H_LSLATE, a=PP_ALIGN.CENTER)

        elif stype == 'content':
            ph(s, slide_def.get('title', ''), pg)
            sub = slide_def.get('subtitle', '')
            if sub: tx(s, 1.3, 1.2, 10, 0.3, sub, 13, H_LSLATE)
            y = 1.6
            for sec in slide_def.get('sections', []):
                hd = sec.get('heading', '')
                bd = sec.get('body', '')
                items = sec.get('items', [])
                if hd: tx(s, 1.3, y, 10, 0.35, hd, 16, H_DARK, True); y += 0.4
                if bd: tx(s, 1.3, y, 10, 0.4, bd, 12, H_BODY); y += 0.45
                for it in items: tx(s, 1.6, y, 10, 0.3, '\u00b7 ' + it, 11, H_BODY); y += 0.35
                y += 0.1
            tbl = slide_def.get('table')
            if tbl:
                hdrs = tbl.get('headers', []); rows = tbl.get('rows', [])
                if hdrs and rows:
                    data = [hdrs] + rows
                    cw = [1.5] * len(hdrs)
                    tb(s, 1.3, y, 10.0, 0.35 * len(data), data, cw)

        elif stype == 'cards':
            ph(s, slide_def.get('title', ''), pg)
            sub = slide_def.get('subtitle', '')
            if sub: tx(s, 1.3, 1.2, 10, 0.3, sub, 13, H_LSLATE)
            cards = slide_def.get('cards', [])
            CBG = RGBColor(0xF0, 0xF4, 0xF8)
            for i, card in enumerate(cards):
                cx = 0.5 + i * 4.0
                num = card.get('num', str(i+1))
                tit = card.get('title', '')
                items = card.get('items', [])
                cd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(1.8), Inches(3.6), Inches(4.0))
                cd.fill.solid(); cd.fill.fore_color.rgb = CBG; cd.line.fill.background()
                rx(s, cx, 1.8, 3.6, 0.035, H_RED)
                c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx+0.2), Inches(2.1), Inches(0.40), Inches(0.40))
                c.fill.solid(); c.fill.fore_color.rgb = H_MAROON; c.line.fill.background()
                tx(s, cx+0.2, 2.1, 0.40, 0.40, num, 15, H_WHITE, True, a=PP_ALIGN.CENTER)
                tx(s, cx+0.7, 2.15, 2.5, 0.3, tit, 20, H_DARK, True)
                for j, it in enumerate(items):
                    tx(s, cx+0.2, 2.9+j*0.45, 3.2, 0.4, '\u00b7 '+it, 11, H_BODY)

        elif stype == 'summary':
            ph(s, slide_def.get('title', ''), pg)
            conclusions = slide_def.get('conclusions', [])
            for i, c in enumerate(conclusions):
                cy = 1.4 + i * 1.3
                rx(s, 1.3, cy, 0.05, 0.8, H_RED)
                tx(s, 1.6, cy, 10, 0.45, c.get('main', ''), 15, H_DARK, True)
                sub = c.get('sub', '')
                if sub: tx(s, 1.6, cy+0.5, 10, 0.35, sub, 12, H_GRAY)

    prs.save(output_path)
    print(f'\u2705 \u5df2\u751f\u6210: {output_path} ({pg}\u9875)')
    return True


def main():
    if len(sys.argv) < 2:
        print('\u7528\u6cd5: python3 pptx_template_render.py <\u5185\u5bb9.json> [\u8f93\u51fa\u8def\u5f84]')
        print('\u793a\u4f8b: python3 pptx_template_render.py \u5185\u5bb9.json \u62a5\u544a.pptx')
        sys.exit(1)

    content_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    with open(content_path, 'r', encoding='utf-8') as f:
        content = json.load(f)

    output = output_path or content.get('output', 'output.pptx')

    build(TEMPLATE_PATH, content, output)


if __name__ == '__main__':
    main()
