#!/usr/bin/env python3
"""
PPT 翻页短视频管线 v2

模式：每页1条核心观点 + 3条要点 + 口语化配音
输出：.pptx（可编辑）+ .mp4（短视频）

用法：
    python3 ppt_to_video.py --bundle 三件套输出/bundle_01.json
    python3 ppt_to_video.py --article 文章.md --title "标题"
    python3 ppt_to_video.py --bundle 三件套输出/bundle_01.json --slides-only
"""

import os, sys, re, json, argparse, asyncio, tempfile
from pathlib import Path
from dataclasses import dataclass, field

from PIL import Image, ImageDraw, ImageFont

# ── 品牌配置 ──
W, H = 1080, 1920

FONT_BOLD = "/System/Library/Fonts/STHeiti Medium.ttc"
FONT_LIGHT = "/System/Library/Fonts/STHeiti Light.ttc"

C = {
    'bg':       (16, 30, 52),
    'card':     (24, 42, 68),
    'red':      (200, 30, 30),
    'gold':     (195, 165, 110),
    'white':    (255, 255, 255),
    'gray':     (155, 165, 180),
    'dark':     (10, 20, 38),
}

OUTPUT_DIR = "/Users/Admin/三件套输出"


# ═══════════════════════════════════════
#  数据结构
# ═══════════════════════════════════════

@dataclass
class Slide:
    type: str          # title / content / closing
    headline: str      # 核心观点（大字）
    bullets: list[str] = field(default_factory=list)
    highlight: str = ""  # 数字/金句（可选强调）
    note: str = ""       # 配音稿


# ═══════════════════════════════════════
#  内容解析
# ═══════════════════════════════════════

def extract_bullets(text: str, max_bullets: int = 3) -> list[str]:
    lines = []
    for line in text.strip().split('\n'):
        line = re.sub(r'\*\*(.+?)\*\*', r'\1', line).strip()
        line = re.sub(r'\[(.+?)\]\(.*?\)', r'\1', line)
        if len(line) < 10:
            continue
        if any(k in line for k in ['是', '不是', '在于', '关键', '核心', '本质', '意味', '证明']):
            lines.append(line)
    if len(lines) < max_bullets:
        for line in text.strip().split('\n'):
            line = re.sub(r'\*\*(.+?)\*\*', r'\1', line).strip()
            if len(line) > 15 and line not in lines:
                lines.append(line)
    return [l.strip() for l in lines[:max_bullets] if l.strip()]


def extract_highlight(text: str) -> str:
    nums = re.findall(r'[+-]?\d+[\.\d]*(?:%|万千百亿)?', text)
    if nums:
        return nums[0]
    for line in text.split('。'):
        line = line.strip()
        if 8 < len(line) < 30 and any(k in line for k in ['是', '在', '有', '将']):
            return line
    return ""


def shorten_bullets(bullets: list[str], max_len: int = 32) -> list[str]:
    """压缩要点到适合幻灯片显示的长度"""
    result = []
    for b in bullets:
        if len(b) <= max_len:
            result.append(b)
        else:
            result.append(b[:max_len-1] + '…')
    return result


def parse_article(md_text: str, title_hint: str = "") -> list[Slide]:
    lines = md_text.strip().split('\n')
    article_title = title_hint or "投资研究"
    for line in lines:
        m = re.match(r'^#\s+(.+)$', line)
        if m:
            article_title = m.group(1).strip()
            break

    slides = []

    # 封面
    slides.append(Slide(
        type="title",
        headline=article_title,
        note=f"大家好，欢迎来到道雷投资研究。今天我们来聊一个很有意思的话题——{article_title}。",
    ))

    # 按标题分节
    sections = []
    current_section = ""
    current_content = []
    for line in lines[2:]:
        stripped = line.strip()
        if not stripped or re.match(r'^#\s+', stripped):
            continue
        if re.match(r'^##+\s+', stripped):
            if current_content:
                sections.append((current_section, '\n'.join(current_content)))
            current_section = re.sub(r'^#+\s+', '', stripped).strip()
            current_content = []
        else:
            current_content.append(stripped)
    if current_content:
        sections.append((current_section, '\n'.join(current_content)))

    # 每节生成 1 页
    for section_title, section_text in sections:
        bullets = extract_bullets(section_text)
        highlight = extract_highlight(section_text)
        primary_bullet = bullets[0] if bullets else ""
        rest_bullets = bullets[1:] if len(bullets) > 1 else []

        if not section_title:
            first_line = section_text.strip().split('\n')[0][:15]
            section_title = f"事件背景：{first_line}…" if first_line else "核心内容"

        slides.append(Slide(
            type="content",
            headline=section_title,
            bullets=shorten_bullets(rest_bullets or bullets),
            highlight=highlight,
            note=primary_bullet,
        ))

    # 尾页
    slides.append(Slide(
        type="closing",
        headline="感谢观看",
        bullets=["关注道雷，获取更多深度研究", "公众号：道雷说道"],
        note="以上就是今天的分享。如果觉得有收获，欢迎关注我的公众号道雷说道。我们下期再见。",
    ))

    return slides


# ═══════════════════════════════════════
#  幻灯片渲染
# ═══════════════════════════════════════

def get_font(size: int, bold=False):
    path = FONT_BOLD if bold else FONT_LIGHT
    try:
        return ImageFont.truetype(path, size)
    except:
        try:
            return ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", size)
        except:
            return ImageFont.load_default()


def draw_left(draw, text, font, x, y, color=C['white'], max_w=None):
    if not text:
        return y
    if max_w is None:
        max_w = W - x - 80
    lines = []
    for word in list(text):
        if not lines:
            lines.append(word)
        else:
            test_line = lines[-1] + word
            bbox = font.getbbox(test_line)
            if bbox and (bbox[2] - bbox[0]) > max_w:
                lines.append(word)
            else:
                lines[-1] = test_line
    lh = font.getbbox("测")[3] - font.getbbox("测")[1] + 6 if font.getbbox("测") else font.size
    for line in lines:
        draw.text((x, y), line, font=font, fill=color)
        y += lh
    return y


def render_slide(slide: Slide) -> Image.Image:
    img = Image.new('RGB', (W, H), C['dark'] if slide.type == 'title' else C['bg'])
    draw = ImageDraw.Draw(img)
    m = 80  # margin

    if slide.type == "title":
        draw.rectangle([m, 340, m+100, 348], fill=C['red'])
        t = slide.headline
        ft = get_font(52, bold=True)
        if len(t) > 14:
            mid = len(t)//2
            for sep in ['的', '，', '：', '—']:
                idx = t.find(sep, len(t)//3, len(t)*2//3)
                if idx > 0:
                    mid = idx+1
                    break
            draw_left(draw, t[:mid], ft, m, 380, C['white'])
            draw_left(draw, t[mid:], ft, m, 455, C['white'])
        else:
            draw_left(draw, t, ft, m, 400, C['white'])
        draw_left(draw, "道雷 · PE/VC投资研究", get_font(22), m, 560, C['gold'])
        draw.rectangle([m, H-160, W-m, H-160], fill=C['red'])

    elif slide.type == "closing":
        ft = get_font(58, bold=True)
        bbox = ft.getbbox(slide.headline)
        tw = (bbox[2]-bbox[0]) if bbox else 0
        draw.text(((W-tw)//2, H//2-80), slide.headline, font=ft, fill=C['gold'])
        y = H//2+20
        for b in slide.bullets:
            fb = get_font(28)
            bbox2 = fb.getbbox(b)
            tw2 = (bbox2[2]-bbox2[0]) if bbox2 else 0
            draw.text(((W-tw2)//2, y), b, font=fb, fill=C['gray'])
            y += 45

    else:
        draw.rectangle([m, 60, m+60, 64], fill=C['red'])
        y = 120
        y = draw_left(draw, slide.headline, get_font(50, bold=True), m, y, C['white'], W-m*2) + 20
        draw.rectangle([m, y, m+80, y+2], fill=C['gold'])
        y += 30

        if slide.highlight:
            y = draw_left(draw, slide.highlight, get_font(72, bold=True), m, y+10, C['gold'], W-m*2) + 20

        if slide.bullets:
            for b in slide.bullets:
                y = draw_left(draw, f"●  {b}", get_font(30), m+10, y+14, C['white'], W-m*2-20) + 6

        draw_left(draw, "道雷投资研究", get_font(18), m, H-50, C['gray'])

    return img


# ═══════════════════════════════════════
#  PPTX 输出
# ═══════════════════════════════════════

def export_pptx(slides: list[Slide], output_path: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide in slides:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(10, 20, 38) if slide.type == 'title' else RGBColor(16, 30, 52)

        if slide.type == "title":
            tb = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            p = tb.text_frame.paragraphs[0]
            p.text = slide.headline
            p.font.size = Pt(40)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            tb2 = s.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.8))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = "道雷 · PE/VC 投资研究"
            p2.font.size = Pt(18)
            p2.font.color.rgb = RGBColor(195, 165, 110)
        elif slide.type == "closing":
            tb = s.shapes.add_textbox(Inches(3), Inches(3), Inches(7), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = slide.headline
            p.font.size = Pt(44)
            p.font.color.rgb = RGBColor(195, 165, 110)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        else:
            tb = s.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
            p = tb.text_frame.paragraphs[0]
            p.text = slide.headline
            p.font.size = Pt(36)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            y = 1.5
            if slide.highlight:
                tb = s.shapes.add_textbox(Inches(1), Inches(y), Inches(11), Inches(1))
                p = tb.text_frame.paragraphs[0]
                p.text = slide.highlight
                p.font.size = Pt(48)
                p.font.color.rgb = RGBColor(195, 165, 110)
                p.font.bold = True
                y += 1.2
            for b in slide.bullets:
                tb = s.shapes.add_textbox(Inches(1.2), Inches(y), Inches(11), Inches(0.6))
                p = tb.text_frame.paragraphs[0]
                p.text = f"•  {b}"
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(255, 255, 255)
                y += 0.6

    prs.save(output_path)
    print(f"  ✅ PPTX 已保存: {output_path}")


# ═══════════════════════════════════════
#  配音 & 视频
# ═══════════════════════════════════════

def build_script(slides: list[Slide]) -> str:
    parts = []
    for s in slides:
        if s.note:
            parts.append(s.note)
        elif s.bullets:
            parts.append(f"接下来看{s.headline}。{'，'.join(s.bullets[:2])}。")
    return "。" .join(parts)


async def gen_tts(text: str, path: str):
    import edge_tts
    com = edge_tts.Communicate(text, "zh-CN-XiaoxiaoNeural", rate="-5%", pitch="+0Hz")
    await com.save(path)
    print(f"  ✅ TTS 配音生成")


def compose_video(images: list[Image.Image], notes: list[str], audio_path: str, out: str):
    from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
    import numpy as np

    audio = AudioFileClip(audio_path)
    dur = audio.duration
    total = sum(len(n) for n in notes) or 1

    clips = []
    for i, img in enumerate(images):
        d = max(2.0, dur * len(notes[i]) / total)
        clips.append(ImageClip(np.array(img)).with_duration(d))

    final = concatenate_videoclips(clips, method="compose")
    if final.duration < dur:
        audio = audio.subclipped(0, final.duration)
    final = final.with_audio(audio)
    final.write_videofile(out, fps=24, codec="libx264", audio_codec="aac",
                          threads=4, preset="medium", bitrate="4000k")
    print(f"  ✅ 视频: {out}  |  {final.duration:.0f}秒 | {len(clips)} 页")
    audio.close()


# ═══════════════════════════════════════
#  CLI
# ═══════════════════════════════════════

async def main():
    p = argparse.ArgumentParser(description="PPT翻页短视频 v2")
    p.add_argument("--bundle")
    p.add_argument("--article")
    p.add_argument("--title")
    p.add_argument("--output")
    p.add_argument("--slides-only", action="store_true")
    p.add_argument("--list", action="store_true")
    args = p.parse_args()

    if args.list:
        for f in sorted(os.listdir(OUTPUT_DIR)):
            if f.startswith("bundle_") and f.endswith(".json"):
                with open(os.path.join(OUTPUT_DIR, f)) as fh:
                    d = json.load(fh)
                print(f"  {f}: {d.get('title', '?')[:50]}")
        return

    md, hint = "", args.title or ""
    if args.bundle:
        with open(args.bundle) as f:
            b = json.load(f)
        hint = hint or b.get("title", "")
        ap = os.path.join(OUTPUT_DIR, "01_算随电动_公众号文章.md")
        if os.path.exists(ap):
            with open(ap) as f:
                md = f.read()
        if not md:
            md = f"# {hint}\n\n{b.get('digest','')}"
    elif args.article:
        with open(args.article) as f:
            md = f.read()
        m = re.search(r"^#\s+(.+)$", md, re.MULTILINE)
        hint = hint or (m.group(1).strip() if m else "投资研究")
    else:
        p.print_help()
        return

    slides = parse_article(md, hint)
    base = re.sub(r'[^\w\u4e00-\u9fff]+', '_', (hint or "视频")[:30]).strip('_')
    print(f"📑 {len(slides)} 页")
    for i, s in enumerate(slides):
        print(f"  {i+1}. [{s.type}] {s.headline[:40]}  notes={len(s.note)}")

    pptx = os.path.join(OUTPUT_DIR, f"{base}.pptx")
    export_pptx(slides, pptx)

    if args.slides_only:
        print(f"\n✅ PPTX: {pptx}")
        return

    images = [render_slide(s) for s in slides]
    notes = [s.note or "、".join(s.bullets[:2]) for s in slides]

    print(f"🎤 TTS ({sum(len(n) for n in notes)}字)...")
    audio = tempfile.mktemp(suffix=".mp3")
    await gen_tts(build_script(slides), audio)

    out = args.output or os.path.join(OUTPUT_DIR, f"翻页视频_{base}.mp4")
    compose_video(images, notes, audio, out)
    os.remove(audio)
    print(f"\n🎉 PPTX: {pptx}  |  视频: {out}")


if __name__ == "__main__":
    asyncio.run(main())
