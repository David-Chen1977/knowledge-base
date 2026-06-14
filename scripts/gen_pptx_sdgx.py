#!/usr/bin/env python3
"""
生成PPT：算电协同进入落地期
Apple 极简设计风格
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os

OUTPUT_DIR = "/Users/Admin/OpencodeWorkspace/内容输出/算电协同2026"
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "算电协同进入落地期_Apple设计.pptx")

# ── Apple 色板 ──
C = {
    'blue':       RGBColor(0x00, 0x71, 0xE3),
    'blue_dark':  RGBColor(0x00, 0x56, 0xB3),
    'text_dark':  RGBColor(0x1D, 0x1D, 0x1F),
    'text_main':  RGBColor(0x3A, 0x3A, 0x3C),
    'text_sec':   RGBColor(0x86, 0x86, 0x8B),
    'white':      RGBColor(0xFF, 0xFF, 0xFF),
    'bg':         RGBColor(0xF5, 0xF5, 0xF7),
    'accent2':    RGBColor(0x34, 0xC7, 0x59),
    'accent3':    RGBColor(0xFF, 0x95, 0x00),
}

def set_cn_font(font_obj, name_cn='PingFang SC', name_en='Helvetica Neue'):
    font_obj.name = name_en
    try:
        rPr = font_obj._rPr
        if rPr is None:
            return
        for tag, typeface in [('a:latin', name_en), ('a:ea', name_cn)]:
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {}); rPr.append(el)
            el.set('typeface', typeface)
    except Exception:
        pass

def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_bg(slide, color=C['white']):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=C['text_main'], align=PP_ALIGN.LEFT, font_name='PingFang SC'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    set_cn_font(p.font)
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=14, color=C['text_main'],
                spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        set_cn_font(p.font)
    return txBox

def add_accent_bar(slide, left=0.6, top=0.55, width=0.08, height=0.05, color=C['blue']):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_footer(slide, text="陈道雷 | 汇竑资本 | 2026.06"):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = C['text_sec']
    set_cn_font(p.font)
    return txBox

def make_section_slide(prs, number, title, subtitle=""):
    slide = add_blank_slide(prs)
    add_bg(slide, C['text_dark'])
    # Number
    add_textbox(slide, 0.6, 2.0, 1, 0.6, f"{number:02d}", font_size=48,
                bold=True, color=C['blue'], font_name='Helvetica Neue')
    # Title
    add_textbox(slide, 0.6, 2.8, 8, 1.2, title, font_size=36,
                bold=True, color=C['white'])
    if subtitle:
        add_textbox(slide, 0.6, 4.2, 8, 0.8, subtitle, font_size=16, color=C['text_sec'])
    add_footer(slide)
    return slide

# ══════════════════════════════════════════
# 开始生成
# ══════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Slide 1: 封面 ──
slide = add_blank_slide(prs)
add_bg(slide, C['text_dark'])
add_accent_bar(slide, left=0.6, top=2.8, width=0.1, height=0.06)
add_textbox(slide, 0.6, 3.0, 11, 1.5, "算电协同进入落地期", font_size=44,
            bold=True, color=C['white'])
add_textbox(slide, 0.6, 4.5, 11, 0.8, "从算力用户的需求倒推，投资机会在哪里", font_size=22, color=C['text_sec'])
add_textbox(slide, 0.6, 5.5, 5, 0.5, "陈道雷 | 汇竑资本 | 2026年6月", font_size=14, color=C['text_sec'])

# ── Slide 2: 引子 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 10, 0.6, "核心问题", font_size=28, bold=True, color=C['text_dark'])
add_textbox(slide, 0.6, 1.4, 11, 1.2,
            "关于算电协同的讨论大多围绕供给侧——但很少有人问：算力用户到底需要什么？",
            font_size=18, color=C['text_main'])
add_bullets(slide, 0.6, 2.8, 10, 2, [
    "算力贵不贵？—— 电费占数据中心运营成本 30-60%",
    "响应快不快？—— 训练可等20ms，推理必须<10ms",
    "供应链稳不稳？—— 停电一次，损失数十万",
], font_size=16)
add_textbox(slide, 0.6, 5.5, 10, 0.8,
            "这三件事，决定了算电协同所有策略选择的底层逻辑。",
            font_size=16, bold=True, color=C['blue'])
add_footer(slide)

# ── Slide 3: 三类负载 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 8, 0.6, "三类算力负载 → 三种选址逻辑", font_size=26, bold=True, color=C['text_dark'])

categories = [
    ("AI训练", "时延容忍 >20ms\n成本极度敏感", "西部能源富集区\n电价 0.2-0.3元", C['blue']),
    ("在线推理", "时延敏感 <10ms\n需要靠近用户", "一线城市/区域中心\n分布式部署", C['accent2']),
    ("离线批量", "中度时延 50-100ms\n成本敏感", "电价低+网络好的\n地区均可", C['accent3']),
]
for i, (title, desc, loc, color) in enumerate(categories):
    x = 0.6 + i * 4.1
    # Card bg
    card = slide.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(3.8), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = C['bg']
    card.line.fill.background()
    # Title accent
    bar = slide.shapes.add_shape(1, Inches(x), Inches(1.6), Inches(3.8), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    # Title
    add_textbox(slide, x + 0.3, 1.9, 3.2, 0.5, title, font_size=20, bold=True, color=C['text_dark'])
    # Description
    add_textbox(slide, x + 0.3, 2.6, 3.2, 1.2, desc.replace('\n', '\n'), font_size=14, color=C['text_main'])
    # Location
    add_textbox(slide, x + 0.3, 4.2, 3.2, 1.2, loc.replace('\n', '\n'), font_size=14, color=color)
add_footer(slide)

# ── Slide 4: 核心洞察 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 10, 0.6, "核心洞察", font_size=28, bold=True, color=C['text_dark'])
add_textbox(slide, 0.6, 1.6, 11, 1.2,
            "同一家AI公司同时需要这三类算力",
            font_size=24, bold=True, color=C['blue'])
add_textbox(slide, 0.6, 3.0, 11, 1.5,
            '算电协同不是"选一个地方"的问题，\n是"建一张网"的问题。\n\n单一节点的成本最优 ≠ 全局最优',
            font_size=18, color=C['text_main'])
add_footer(slide)

# ── Slide 5: 六大决策变量 ──
make_section_slide(prs, 2, "六个决策变量", "从需求到选址")

slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 8, 0.6, "六个关键决策变量", font_size=26, bold=True)

vars_data = [
    ("时延约束", "不可妥协的硬限制", C['blue']),
    ("电费成本", "TCO 而非仅电价", C['accent2']),
    ("绿电碳合规", "从加分项到必选项", C['accent3']),
    ("政府政策", "正反两面都在影响", C['blue']),
    ("交付时间", "最被低估的决策变量", C['accent2']),
    ("供应链可靠性", "最后一个，最致命", C['accent3']),
]
for i, (name, desc, color) in enumerate(vars_data):
    row = i // 3
    col = i % 3
    x = 0.6 + col * 4.2
    y = 1.6 + row * 2.6
    # Card
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = C['bg']
    card.line.fill.background()
    # Accent dot
    dot = slide.shapes.add_shape(9, Inches(x + 0.2), Inches(y + 0.3), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_textbox(slide, x + 0.5, y + 0.2, 3, 0.5, name, font_size=16, bold=True, color=C['text_dark'])
    add_textbox(slide, x + 0.5, y + 0.8, 3, 0.8, desc, font_size=13, color=C['text_sec'])
add_footer(slide)

# ── Slide 6: 交付时间 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 8, 0.6, "最被低估的变量：交付时间", font_size=26, bold=True, color=C['text_dark'])
add_textbox(slide, 0.6, 1.3, 10, 0.8,
            "新建智算中心 2-4年 vs AI算力需求年增长 50-80%——等不起。",
            font_size=16, color=C['text_main'])

delivery = [
    ("存量机房改造", "3-6个月", "当前急需"),
    ("分布式算力仓", "1-3个月", "即时响应"),
    ("新建智算中心", "2-4年", "长期规划"),
]
for i, (mode, time, scene) in enumerate(delivery):
    x = 0.6 + i * 4.2
    card = slide.shapes.add_shape(1, Inches(x), Inches(2.3), Inches(3.8), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = C['bg']
    card.line.fill.background()
    add_textbox(slide, x + 0.3, 2.5, 3.2, 0.5, mode, font_size=18, bold=True, color=C['blue'])
    add_textbox(slide, x + 0.3, 3.2, 3.2, 0.5, time, font_size=22, bold=True, color=C['text_dark'])
    add_textbox(slide, x + 0.3, 3.9, 3.2, 0.5, f"适合: {scene}", font_size=14, color=C['text_sec'])

add_textbox(slide, 0.6, 5.5, 11, 0.8,
            '存量改造有天花板，但"新建慢"不等于"新建没价值"——两者是时间维度上的互补。',
            font_size=15, bold=True, color=C['text_dark'])
add_footer(slide)

# ── Slide 7: 三条路径 ──
make_section_slide(prs, 3, "三条落地路径", "从变量到路径")

# ── Slide 8: 路径一 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide, color=C['blue'])
add_textbox(slide, 0.6, 0.5, 10, 0.6, "路径一：西部集中式新建", font_size=26, bold=True, color=C['text_dark'])
add_textbox(slide, 0.6, 1.1, 5, 0.5, "时间定位：中长期（2-4年）", font_size=14, bold=True, color=C['blue'])
add_bullets(slide, 0.6, 1.8, 11, 4, [
    "服务：长期战略 + 成本敏感的客户（大模型训练、科学计算）",
    "优势：电价 0.2-0.3元/kWh、绿电 40%+、千P级可扩展",
    "壁垒：数十亿级资本、政府关系（拿地+能耗指标）、绿电直供能力",
    "客户逻辑：今天的紧迫用旧改解决，明后年的需求需要战略性基础设施承接",
    "何时唯一选项：当需要千P级以上集群，存量建筑条件无法满足时",
], font_size=15)
add_footer(slide)

# ── Slide 9: 路径二 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide, color=C['accent2'])
add_textbox(slide, 0.6, 0.5, 10, 0.6, "路径二：分布式边缘旧改 ⭐", font_size=26, bold=True, color=C['text_dark'])
add_textbox(slide, 0.6, 1.1, 5, 0.5, "时间定位：短期（1-6个月交付）", font_size=14, bold=True, color=C['accent2'])
add_bullets(slide, 0.6, 1.8, 11, 4, [
    "服务：时延敏感 + 现在就要的用户（推理、自动驾驶、企业私有化）",
    "优势：存量机房盘活，绕过新建审批，3-6个月交付",
    "壁垒：存量机房资源获取能力、快速交付能力、本地运维网络",
    "为何当前最活跃：AI需求增速远超新建周期，旧改填补供需缺口",
    "远期灵活：未来新建节点投产后，旧改节点转向时延敏感负载",
], font_size=15)
add_textbox(slide, 0.6, 5.8, 11, 0.6, "最适合 PE/VC 入场的路径——轻资产、可复制、有网络效应",
            font_size=16, bold=True, color=C['accent2'])
add_footer(slide)

# ── Slide 10: 路径三 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide, color=C['accent3'])
add_textbox(slide, 0.6, 0.5, 10, 0.6, "路径三：存量改造升级", font_size=26, bold=True, color=C['text_dark'])
add_textbox(slide, 0.6, 1.1, 5, 0.5, "时间定位：中短期（3-6个月）", font_size=14, bold=True, color=C['accent3'])
add_bullets(slide, 0.6, 1.8, 11, 4, [
    "服务：已有资产提效（运营商自有数据中心、企业自建机房）",
    "关键手段：气流组织优化（PUE 1.7→1.5，投资几乎为零）",
    "液冷两条腿：半液冷适合存量改造（PUE 1.2），芯片级适合新建（PUE 1.05）",
    "综合能源平台：风光储+绿电直连+AI调度，实现用电成本最低",
    "结论：投入产出比最高的路径，但天花板可视",
], font_size=15)
add_textbox(slide, 0.6, 5.8, 11, 0.6, "好生意，但不是大生意——适合作为生态能力补充",
            font_size=16, bold=True, color=C['accent3'])
add_footer(slide)

# ── Slide 11: 投资变现对比 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 8, 0.6, "投资人视角：三条路径变现逻辑对比", font_size=24, bold=True, color=C['text_dark'])

# Table approach - use text boxes as a simple table
headers = ["维度", "西部集中式", "分布式边缘", "存量改造"]
rows_data = [
    ["谁付费", "AI公司/云厂商", "AI应用/企业", "IDC运营商"],
    ["商业模式", "算力租赁/托管", "API零售/节点租赁", "EPC/改造分成"],
    ["资本强度", "极高（数十亿）", "低（百万级）", "低到中"],
    ["利润率", "中", "中高（延迟溢价）", "中（技术溢价）"],
    ["资产属性", "重资产", "轻资产", "项目制"],
    ["退出通道", "REITs", "被收购/上市", "服务合同"],
]

# Draw header row
for j, h in enumerate(headers):
    x = 0.6 + j * 3.1
    if j == 0:
        w = 1.5
    else:
        w = 3.1
    add_textbox(slide, x, 1.5, w, 0.4, h, font_size=12, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    # Header bg
    hdr = slide.shapes.add_shape(1, Inches(x), Inches(1.5), Inches(w), Inches(0.4))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = C['blue']
    hdr.line.fill.background()
    hdr_element = hdr
    # Re-add text on top
    add_textbox(slide, x, 1.5, w, 0.4, h, font_size=12, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

for i, row in enumerate(rows_data):
    for j, cell in enumerate(row):
        x = 0.6 + j * 3.1
        if j == 0:
            w = 1.5
        else:
            w = 3.1
        y = 2.0 + i * 0.45
        bg_color = C['bg'] if i % 2 == 0 else C['white']
        cell_bg = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.42))
        cell_bg.fill.solid()
        cell_bg.fill.fore_color.rgb = bg_color
        cell_bg.line.fill.background()
        fc = C['blue'] if j == 0 else C['text_main']
        add_textbox(slide, x + 0.1, y + 0.02, w - 0.2, 0.38, cell, font_size=11, color=fc)

add_footer(slide)

# ── Slide 12: 综合判断 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 8, 0.6, "综合判断：投资优先级排序", font_size=26, bold=True, color=C['text_dark'])

priorities = [
    ("P0", "分布式边缘", "立即关注", "推理需求爆发 + 存量盘活绕开审批", C['blue']),
    ("P1", "西部集中式", "布局，选对合作方", "需求确定但门槛高，适合产业资本牵头", C['accent2']),
    ("P2", "存量改造服务", "选择性参与", "收入确定但非高速赛道，作为生态补充", C['accent3']),
]
for i, (level, name, action, reason, color) in enumerate(priorities):
    y = 1.6 + i * 1.8
    # Level badge
    badge = slide.shapes.add_shape(1, Inches(0.6), Inches(y), Inches(0.8), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_textbox(slide, 0.6, y + 0.05, 0.8, 0.4, level, font_size=16, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    # Content
    add_textbox(slide, 1.7, y, 3, 0.5, name, font_size=18, bold=True, color=C['text_dark'])
    add_textbox(slide, 4.5, y, 2.5, 0.5, action, font_size=12, color=color)
    add_textbox(slide, 1.7, y + 0.6, 8, 0.8, reason, font_size=13, color=C['text_sec'])
add_footer(slide)

# ── Slide 13: 未来推演 ──
make_section_slide(prs, 4, "未来场景推演", "AI时代 ≠ 互联网时代，技术迭代是指数级的")

# ── Slide 14: 未来推演内容 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 10, 0.6, "指数级变化 vs 线性思维", font_size=26, bold=True, color=C['text_dark'])

add_bullets(slide, 0.6, 1.4, 11, 5, [
    '今天认为"超前"的规划，三年后可能不够用',
    '"过度建设"风险被高估，"建设不足"风险被低估',
    "技术迭代速度持续领先于物理世界的建设速度",
    '交付速度本身就是稀缺能力——6个月 > 2年的"最优方案"',
    '基础设施从"成本中心"变成"战略瓶颈"',
], font_size=17, color=C['text_main'])
add_textbox(slide, 0.6, 5.5, 11, 0.8,
            '不要用"泡沫"框架去理解AI算力基础设施投资。\n需要警惕的不是泡沫，而是用线性思维做指数级世界的决策。',
            font_size=16, bold=True, color=C['blue'])
add_footer(slide)

# ── Slide 15: 确定性 ──
slide = add_blank_slide(prs)
add_bg(slide)
add_accent_bar(slide)
add_textbox(slide, 0.6, 0.5, 10, 0.6, "快速变化中的五个确定性", font_size=26, bold=True, color=C['text_dark'])

certainties = [
    "确定性1：算力总需求持续增长，增速可能超过大多数预测",
    "确定性2：能源成为算力的最终瓶颈——电网扩容是物理限制",
    "确定性3：基础设施交付速度越来越重要，快比优更有竞争力",
    "确定性4：绿电和碳合规成为算力准入的门槛，不是锦上添花",
    "确定性5：跨越技术—工程—商业三层面的团队极度稀缺",
]
add_bullets(slide, 0.6, 1.4, 11, 5, certainties, font_size=16, color=C['text_main'])
add_textbox(slide, 0.6, 5.8, 11, 0.6, '旧改填补"等不起"的当下，新建承载"停不下来"的未来',
            font_size=16, bold=True, color=C['blue'])
add_footer(slide)

# ── Slide 16: 结束 ──
slide = add_blank_slide(prs)
add_bg(slide, C['text_dark'])
add_accent_bar(slide, top=3.2, left=5.5, width=1.5, height=0.06, color=C['blue'])
add_textbox(slide, 2, 3.5, 9, 1, "谢谢", font_size=48, bold=True, color=C['white'],
            align=PP_ALIGN.CENTER)
add_textbox(slide, 2, 4.5, 9, 0.6, "陈道雷 | 汇竑资本 | 2026年6月", font_size=16, color=C['text_sec'],
            align=PP_ALIGN.CENTER)

# ── Save ──
os.makedirs(OUTPUT_DIR, exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"✅ PPT 已生成: {OUTPUT_PATH}")
print(f"   共 {len(prs.slides)} 页幻灯片")
